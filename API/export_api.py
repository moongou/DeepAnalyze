"""
Export API for DeepAnalyze API Server
Handles report generation in MD/PDF/DOCX/PPTX formats
"""

import json
import os
import time
import uuid
from pathlib import Path
from typing import List, Optional, Dict, Any

from fastapi import APIRouter, HTTPException, Depends, Body
from fastapi.responses import FileResponse

from storage import storage
from auth_api import get_current_user, require_auth
from utils import get_thread_workspace, generate_report_from_messages

router = APIRouter(prefix="/v1/export", tags=["export"])


@router.post("/report")
def export_report(
    messages: List[Dict[str, Any]] = Body(...),
    thread_id: Optional[str] = Body(None),
    format: str = Body(default="md"),
    username: Optional[str] = Depends(get_current_user),
):
    require_auth(username)
    if format not in ("md", "pdf", "docx", "pptx"):
        raise HTTPException(status_code=400, detail=f"Unsupported format: {format}")

    workspace_dir = "workspace/_exports"
    if thread_id:
        workspace_dir = get_thread_workspace(thread_id)
    os.makedirs(workspace_dir, exist_ok=True)

    # Extract assistant reply from messages
    assistant_reply = ""
    for msg in messages:
        if isinstance(msg, dict) and msg.get("role") == "assistant":
            assistant_reply = msg.get("content", "")
            break

    if not assistant_reply:
        raise HTTPException(status_code=400, detail="No assistant reply found in messages")

    # Generate the report using utils
    generated_files = []
    report_content = generate_report_from_messages(
        messages, assistant_reply, workspace_dir, thread_id or "", generated_files
    )

    if format == "md":
        filename = f"report_{uuid.uuid4().hex[:8]}.md"
        filepath = os.path.join(workspace_dir, filename)
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(assistant_reply)
        return FileResponse(filepath, filename=filename, media_type="text/markdown")

    # For PDF/DOCX/PPTX, try to use the existing demo modules
    if format == "pdf":
        try:
            from demo.chat.pdf_utils import generate_pdf as gen_pdf, extract_markdown_sections
            filename = f"report_{uuid.uuid4().hex[:8]}.pdf"
            filepath = os.path.join(workspace_dir, filename)
            gen_pdf(assistant_reply, filepath)
            return FileResponse(filepath, filename=filename, media_type="application/pdf")
        except ImportError:
            raise HTTPException(status_code=501, detail="PDF generation module not available")

    if format == "docx":
        try:
            from demo.chat.docx_utils import generate_docx as gen_docx
            filename = f"report_{uuid.uuid4().hex[:8]}.docx"
            filepath = os.path.join(workspace_dir, filename)
            gen_docx(assistant_reply, filepath)
            return FileResponse(
                filepath, filename=filename,
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
        except ImportError:
            raise HTTPException(status_code=501, detail="DOCX generation module not available")

    if format == "pptx":
        try:
            from pptx import Presentation
            from pptx.util import Inches
            filename = f"report_{uuid.uuid4().hex[:8]}.pptx"
            filepath = os.path.join(workspace_dir, filename)
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "DeepAnalyze Report"
            slide.placeholders[1].text = assistant_reply[:500] + "..."
            prs.save(filepath)
            return FileResponse(
                filepath, filename=filename,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except ImportError:
            raise HTTPException(status_code=501, detail="PPTX generation module not available")
