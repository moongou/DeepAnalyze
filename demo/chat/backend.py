from __future__ import annotations

import openai
from typing import Optional, List, Dict, Any, Tuple
import json
import os
import shutil
import re
import io
import contextlib
import traceback
from pathlib import Path
from urllib.parse import quote
import subprocess
import sys
import tempfile
import requests
import threading
import http.server
from functools import partial
import socketserver
import sqlite3
import hashlib
import secrets
import time as time_module
import jwt
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Query, Depends, Request
from fastapi.responses import JSONResponse, Response, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from typing import List, Optional
import httpx
import uvicorn
import os
import re
import json
from fastapi.responses import StreamingResponse
from fastapi.concurrency import run_in_threadpool
import os
import re
from copy import deepcopy
import openai
from fastapi import FastAPI, Body
from fastapi.responses import StreamingResponse

import re
from matplotlib import font_manager

os.environ.setdefault("MPLBACKEND", "Agg")

# 引入中央配置模块
try:
    from config import get_config, get_fonts_dir
    _CONFIG = get_config()
    _FONT_DIR = str(get_fonts_dir())
    # 确保目录存在
    if not os.path.exists(_FONT_DIR):
        print(f"警告: 字体目录不存在: {_FONT_DIR}")
except ImportError:
    print("警告: 无法导入 config 模块，使用回退路径解析")
    _FONT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "../../assets/fonts")

# 注册中文字体到 matplotlib
if os.path.exists(_FONT_DIR):
    for font_file in os.listdir(_FONT_DIR):
        if font_file.lower().endswith(('.ttf', '.ttc', '.otf')):
            try:
                font_manager.fontManager.addfont(os.path.join(_FONT_DIR, font_file))
            except Exception as e:
                print(f"Error registering font {font_file}: {e}")

import chardet
from docx import Document
from sqlalchemy import create_engine, text
from sqlalchemy.engine import URL
import pandas as pd
from datetime import datetime

# 使用新的模块化 PDF 工具
from pdf_utils import (
    register_chinese_fonts,
    get_chinese_style,
    extract_markdown_sections,
    clean_md_text,
    generate_pdf as generate_pdf_module,
)

from docx_utils import (
    generate_docx as generate_docx_module,
    create_docx_document,
    clean_md_text_for_docx,
    extract_markdown_blocks,
)

# 雨途斩棘录模块
from yutu_zhanyilu import (
    init_yutu_if_needed,
    add_error_solution,
    search_errors,
    get_error_by_hash,
    update_error_solution,
    delete_error,
    get_yutu_html,
    update_yutu_html,
)


# ========== 自动错误记录功能 ==========
def detect_and_record_error(exe_output: str, code_str: str, workspace_dir: str) -> dict:
    """
    检测代码执行输出中的错误，并自动记录到雨途斩棘录

    Returns:
        dict: {
            "has_error": bool,           # 是否检测到错误
            "error_type": str,           # 错误类型
            "error_message": str,        # 错误消息
            "recorded": bool,            # 是否已记录到知识库
            "similar_found": bool        # 是否发现相似错误
        }
    """
    import re

    result = {
        "has_error": False,
        "error_type": "Unknown",
        "error_message": "",
        "recorded": False,
        "similar_found": False
    }

    # 如果没有输出或输出为空，认为没有错误
    if not exe_output or not exe_output.strip():
        return result

    # 检测常见的Python错误模式
    error_patterns = [
        # ImportError
        (r"ModuleNotFoundError|ImportError:\s*(.+?)(?:\n|$)", "ImportError"),
        # ValueError
        (r"ValueError:\s*(.+?)(?:\n|$)", "ValueError"),
        # TypeError
        (r"TypeError:\s*(.+?)(?:\n|$)", "TypeError"),
        # RuntimeError
        (r"RuntimeError:\s*(.+?)(?:\n|$)", "RuntimeError"),
        # KeyError
        (r"KeyError:\s*['\"](.+?)['\"]", "KeyError"),
        # FileNotFoundError
        (r"FileNotFoundError:\s*(.+?)(?:\n|$)", "FileNotFoundError"),
        # TimeoutError
        (r"TimeoutError:\s*(.+?)(?:\n|$)", "TimeoutError"),
        # SyntaxError
        (r"SyntaxError:\s*(.+?)(?:\n|$)", "SyntaxError"),
        # AttributeError
        (r"AttributeError:\s*(.+?)(?:\n|$)", "AttributeError"),
        # IndexError
        (r"IndexError:\s*(.+?)(?:\n|$)", "IndexError"),
        # MemoryError
        (r"MemoryError:\s*(.+?)(?:\n|$)", "MemoryError"),
        # ZeroDivisionError
        (r"ZeroDivisionError:\s*(.+?)(?:\n|$)", "ZeroDivisionError"),
        # FPDF/报告库相关错误
        (r"FPDF.*?DeprecationWarning|DeprecationWarning.*?FPDF", "FPDFDeprecationWarning"),
        # matplotlib 字体相关错误
        (r"Font.*?not found|FontProperties.*?not found", "FontNotFoundError"),
        # Unicode错误
        (r"UnicodeDecodeError|UnicodeEncodeError|UnicodeError", "UnicodeError"),
        # 一般错误
        (r"(?:Error|Exception|Error:)\s*(.+?)(?:\n|$)", "RuntimeError"),
        # 通用错误检测
        (r"Traceback \(most recent call last\):(.+?)(?:\n\n|\Z)", "PythonError"),
    ]

    error_type = "Unknown"
    error_message = ""

    # 遍历错误模式进行匹配
    for pattern, err_type in error_patterns:
        match = re.search(pattern, exe_output, re.IGNORECASE | re.DOTALL)
        if match:
            error_type = err_type
            # 提取错误消息，清理格式
            error_msg = match.group(1) if match.lastindex else match.group(0)
            # 限制错误消息长度，避免过长
            error_msg = error_msg.strip()[:500] if error_msg else exe_output[:500]
            error_message = error_msg
            break

    # 检查是否确实有错误（排除 "Success" 等正常输出）
    has_error = (
        error_type != "Unknown" or
        ("Error" in exe_output and "Success" not in exe_output) or
        ("error" in exe_output.lower() and "0 error" not in exe_output.lower())
    ) and not (
        # 排除成功的输出
        exe_output.strip().endswith("OK") or
        "Successfully" in exe_output or
        "successfully" in exe_output
    )

    # 进一步检测：如果输出中包含 "Error" 但不是真正的错误，也需要过滤
    if "Error" in exe_output and "[Error]:" in exe_output:
        has_error = True
    elif error_type == "Unknown" and "error" in exe_output.lower():
        # 可能是未知的错误格式，尝试提取整段
        error_type = "RuntimeError"
        error_message = exe_output[:500]

    if not has_error:
        return result

    result["has_error"] = True
    result["error_type"] = error_type
    result["error_message"] = error_message

    # 检查是否已有相似的错误记录（避免重复记录）
    try:
        similar_errors = search_errors(keywords=[error_type], page_size=1)
        if similar_errors and similar_errors.get("items"):
            for item in similar_errors["items"]:
                # 检查错误消息是否相似
                if item.get("error_message") and error_message:
                    # 简单的相似性检查：是否包含相同的关键词
                    msg_keywords = set(error_message.lower().split())
                    existing_keywords = set(item["error_message"].lower().split())
                    common = msg_keywords.intersection(existing_keywords)
                    if len(common) >= 3:  # 有3个以上共同关键词
                        result["similar_found"] = True
                        break
    except Exception as e:
        print(f"检查相似错误失败: {e}")

    # 自动记录错误到雨途斩棘录（如果没有相似记录）
    if not result["similar_found"]:
        try:
            # 生成解决方案建议
            solution = generate_solution_suggestion(error_type, error_message, code_str)
            solution_code = generate_fix_code(error_type, error_message, code_str)

            # 记录错误
            add_error_solution(
                error_type=error_type,
                error_message=error_message,
                error_context=f"工作区: {workspace_dir}\n代码长度: {len(code_str)} 字符",
                solution=solution,
                solution_code=solution_code,
                confidence=0.7,  # 自动记录的置信度稍低
                created_by="system_auto"
            )
            result["recorded"] = True
            print(f"[雨途斩棘录] 自动记录错误: {error_type} - {error_message[:50]}...")
        except Exception as e:
            print(f"[雨途斩棘录] 自动记录失败: {e}")

    return result


def generate_solution_suggestion(error_type: str, error_message: str, code_str: str) -> str:
    """根据错误类型和消息生成解决方案建议"""
    if error_type == "TypeError" and re.search(
        r"(missing\s+\d+\s+required positional argument|unexpected keyword argument|takes\s+\d+\s+positional argument)",
        error_message,
        re.IGNORECASE,
    ):
        return "检查函数/方法定义与调用签名是否一致，确认位置参数数量、关键字参数名称、默认值以及所用库版本下的 API 写法完全匹配。对于自定义函数，先统一定义再调用；对于 pandas / matplotlib / report 相关方法，务必按当前版本文档传参。"
    suggestions = {
        "ImportError": "检查是否已安装所需的Python包，可能需要使用 pip install 安装缺失的模块。",
        "ValueError": "检查输入数据的类型和格式，确保参数值在有效范围内。",
        "TypeError": "检查变量类型，确保操作符两边的数据类型兼容。",
        "FileNotFoundError": "检查文件路径是否正确，确保文件存在于指定位置。",
        "UnicodeError": "检查文件编码，可能需要指定正确的编码格式（如 encoding='utf-8'）。",
        "SyntaxError": "检查代码语法，确保Python语法正确。",
        "KeyError": "检查字典键是否存在，使用 .get() 方法或先检查键是否存在。",
        "IndexError": "检查列表/数组索引是否越界，确保索引值在有效范围内。",
        "AttributeError": "检查对象是否有该属性，确保使用正确的属性名。",
        "FontNotFoundError": "检查字体文件路径是否正确，确保字体文件存在。",
        "FPDFDeprecationWarning": "使用 fpdf2 替代 fpdf，或更新 FPDF 库到最新版本。",
    }
    return suggestions.get(error_type, f"遇到 {error_type} 错误，请检查代码逻辑并参考错误消息进行修复。")


def generate_fix_code(error_type: str, error_message: str, code_str: str) -> str:
    """生成修复代码的示例"""
    if error_type == "TypeError" and re.search(
        r"(missing\s+\d+\s+required positional argument|unexpected keyword argument|takes\s+\d+\s+positional argument)",
        error_message,
        re.IGNORECASE,
    ):
        return '''# 解决方案：先统一函数定义与调用签名，再执行
def build_summary(df, metrics, top_n=10):
    # 函数定义里的参数名、顺序、默认值要与调用处保持一致
    return {"rows": len(df), "metrics": metrics[:top_n]}

# 正确调用：参数数量与名称必须完全匹配
summary = build_summary(df, metrics, top_n=5)

# 如果是第三方库方法报错，请先确认当前版本文档中的参数签名
# 例如：不要给不支持的参数名，不要遗漏必填位置参数
'''
    fix_examples = {
        "ImportError": '''# 解决方案：确保所有依赖已安装
import subprocess
import sys

def install_package(package_name):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package_name])

# 示例：安装缺失的包
# install_package("package_name")
''',
        "FileNotFoundError": '''# 解决方案：检查文件是否存在，或使用绝对路径
import os

file_path = "your_file.csv"
if os.path.exists(file_path):
    # 读取文件
    pass
else:
    print(f"文件不存在: {file_path}")
    # 列出工作区文件帮助调试
    print("工作区文件:", os.listdir("."))
''',
        "UnicodeError": '''# 解决方案：指定正确的编码格式
import chardet

# 检测文件编码
with open("file.csv", "rb") as f:
    raw_data = f.read()
    result = chardet.detect(raw_data)
    encoding = result["encoding"]

# 使用检测到的编码读取文件
with open("file.csv", "r", encoding=encoding) as f:
    content = f.read()
''',
    }
    return fix_examples.get(error_type, "")

from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.pagesizes import A4


# 动态生成 matplotlib 中文支持代码（使用运行时解析的字体路径）
def get_chinese_matplotlib_setup() -> str:
    """
    生成 matplotlib 中文支持代码
    使用动态路径避免硬编码
    """
    # 使用已解析的字体目录
    font_dir = _FONT_DIR if '_FONT_DIR' in globals() else os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "../../assets/fonts"
    )

    return f"""
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import os

# 使用动态解析的字体目录
_ASSETS_FONTS = "{font_dir}"
_FONT_DIRS = [
    _ASSETS_FONTS,
    "/System/Library/Fonts",
    "/Library/Fonts",
    "/System/Library/Fonts/Supplemental",
]
for _d in _FONT_DIRS:
    if os.path.exists(_d):
        for _ff in os.listdir(_d):
            if _ff.lower().endswith(('.ttf', '.ttc', '.otf')):
                try:
                    fm.fontManager.addfont(os.path.join(_d, _ff))
                except Exception:
                    pass

# 优先使用支持中文的字体系列（按优先级排序）
plt.rcParams['font.sans-serif'] = [
    # assets/fonts 中的中文字体（优先，纯 TTF）
    'SimHei',           # assets/simhei.ttf（黑体，主标题）
    'SimKai',           # assets/simkai.ttf（楷体，引用/强调）
    'STFangSong',       # assets/STFangSong.ttf（仿宋，正文）
    'STHeiti',          # assets/STHeiti.ttf（黑体，备选）
    # macOS 系统 CJK 字体
    'Heiti TC',         # macOS 繁体黑体
    'PingFang SC',      # macOS 苹方简体中文
    'PingFang TC',      # macOS 苹方繁体中文
    'Kaiti SC',         # macOS 楷体
    'Songti SC',        # macOS 宋体
    'Arial Unicode MS', # macOS 内置
    'DejaVu Sans',
    'sans-serif',
]
plt.rcParams['axes.unicode_minus'] = False
# 额外：确保数据标签、图例、标题全部使用 sans-serif 字体族
plt.rcParams['font.family'] = 'sans-serif'
"""

# 初始化时生成一次（如果 _FONT_DIR 已可用）
Chinese_matplot_str = get_chinese_matplotlib_setup() if '_FONT_DIR' in globals() else get_chinese_matplotlib_setup()


def execute_code_safe(
    code_str: str, workspace_dir: str = None, timeout_sec: int = 120
) -> str:
    """在独立进程中执行代码，支持超时，避免阻塞主进程。"""
    if workspace_dir is None:
        workspace_dir = WORKSPACE_BASE_DIR
    exec_cwd = os.path.abspath(workspace_dir)
    os.makedirs(exec_cwd, exist_ok=True)
    tmp_path = None
    try:
        fd, tmp_path = tempfile.mkstemp(suffix=".py", dir=exec_cwd)
        os.close(fd)
        with open(tmp_path, "w", encoding="utf-8") as f:
            f.write(code_str)
        # 在子进程中设置无界面环境变量，避免 GUI 后端
        child_env = os.environ.copy()
        child_env.setdefault("MPLBACKEND", "Agg")
        child_env.setdefault("QT_QPA_PLATFORM", "offscreen")
        # 显式设置 R 环境路径，解决 libRblas.dylib 加载问题
        try:
            r_home = subprocess.check_output(["R", "RHOME"], text=True).strip()
            if r_home and os.path.exists(r_home):
                child_env.setdefault("R_HOME", r_home)
                r_lib = os.path.join(r_home, "lib")
                brew_lib = "/opt/homebrew/lib"

                # 构造库搜索路径，包含 R 自身库和 Homebrew 公共库
                lib_paths = [r_lib]
                if os.path.exists(brew_lib):
                    lib_paths.append(brew_lib)

                # 解决 libRblas.dylib 缺失问题：在 workspace 下创建软连接文件夹以引导 rpy2 (macOS brew 常见问题)
                fake_lib_dir = os.path.abspath(os.path.join(workspace_dir, ".lib"))
                os.makedirs(fake_lib_dir, exist_ok=True)
                fake_blas = os.path.join(fake_lib_dir, "libRblas.dylib")
                target_r_lib = os.path.join(r_lib, "libR.dylib")
                if not os.path.exists(fake_blas) and os.path.exists(target_r_lib):
                    try:
                        os.symlink(target_r_lib, fake_blas)
                    except:
                        pass
                if os.path.exists(fake_lib_dir):
                    lib_paths.insert(0, fake_lib_dir)

                path_str = ":".join(lib_paths)

                if "DYLD_LIBRARY_PATH" in child_env:
                    child_env["DYLD_LIBRARY_PATH"] = f"{path_str}:{child_env['DYLD_LIBRARY_PATH']}"
                else:
                    child_env["DYLD_LIBRARY_PATH"] = path_str

                # 设置多种路径变量以增强兼容性
                child_env["LD_LIBRARY_PATH"] = child_env["DYLD_LIBRARY_PATH"]
                child_env["DYLD_FALLBACK_LIBRARY_PATH"] = child_env["DYLD_LIBRARY_PATH"]
        except Exception:
            # 兜底常用路径
            r_home = "/opt/homebrew/opt/r/lib/R"
            if os.path.exists(r_home):
                child_env.setdefault("R_HOME", r_home)
                r_lib = os.path.join(r_home, "lib")
                child_env["DYLD_LIBRARY_PATH"] = f"{r_lib}:/opt/homebrew/lib"
                child_env["LD_LIBRARY_PATH"] = child_env["DYLD_LIBRARY_PATH"]

        # 设置 PYTHONPATH 包含 backend.py 所在目录，使 agent_utils 可用
        backend_dir = os.path.dirname(os.path.abspath(__file__))
        if "PYTHONPATH" in child_env:
            child_env["PYTHONPATH"] = f"{backend_dir}:{child_env['PYTHONPATH']}"
        else:
            child_env["PYTHONPATH"] = backend_dir

        child_env.pop("DISPLAY", None)

        completed = subprocess.run(
            [sys.executable, tmp_path],
            cwd=exec_cwd,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=timeout_sec,
            env=child_env,
        )
        output = (completed.stdout or "") + (completed.stderr or "")
        return output
    except subprocess.TimeoutExpired:
        return f"[Timeout]: execution exceeded {timeout_sec} seconds"
    except Exception as e:
        return f"[Error]: {str(e)}"
    finally:
        try:
            if tmp_path and os.path.exists(tmp_path):
                os.remove(tmp_path)
        except Exception:
            pass


READONLY_SQL_START_KEYWORDS = {"select", "with", "explain"}
FORBIDDEN_SQL_KEYWORDS = {
    "alter",
    "attach",
    "call",
    "copy",
    "create",
    "delete",
    "detach",
    "drop",
    "execute",
    "grant",
    "insert",
    "load",
    "merge",
    "pragma",
    "replace",
    "revoke",
    "truncate",
    "unload",
    "update",
    "vacuum",
}
EXECUTION_LANGUAGE_ALIASES = {
    "": "python",
    "py": "python",
    "python": "python",
    "python3": "python",
    "sql": "sql",
    "postgres": "sql",
    "postgresql": "sql",
    "mysql": "sql",
    "sqlite": "sql",
    "mssql": "sql",
    "sqlserver": "sql",
    "oracle": "sql",
    "r": "r",
    "rscript": "r",
}


def _normalize_execution_language(language: str, code: str = "") -> str:
    normalized = str(language or "").strip().lower()
    if normalized in EXECUTION_LANGUAGE_ALIASES:
        return EXECUTION_LANGUAGE_ALIASES[normalized]
    code_head = str(code or "").strip().lower()
    if re.match(r"^(with|select|explain)\b", code_head):
        return "sql"
    return "python"


def _extract_executable_block(response_text: str) -> Dict[str, str]:
    """Extract the first executable Code/fenced block and normalize its language."""
    text_body = str(response_text or "")
    code_match = re.search(r"<code>(.*?)</code>", text_body, re.DOTALL | re.IGNORECASE)
    search_body = (code_match.group(1) if code_match else text_body).strip()
    fence_match = re.search(
        r"```\s*([A-Za-z0-9_+.-]*)\s*\n(.*?)```",
        search_body,
        re.DOTALL,
    )
    if fence_match:
        raw_language = fence_match.group(1) or ""
        code = (fence_match.group(2) or "").strip()
        return {
            "language": _normalize_execution_language(raw_language, code),
            "raw_language": raw_language,
            "code": code,
        }
    return {
        "language": _normalize_execution_language("", search_body),
        "raw_language": "",
        "code": search_body,
    }


def _strip_sql_comments(sql: str) -> str:
    sql = re.sub(r"/\*.*?\*/", " ", str(sql or ""), flags=re.DOTALL)
    sql = re.sub(r"--.*?$", " ", sql, flags=re.MULTILINE)
    return sql.strip()


def _validate_readonly_sql(sql: str) -> str:
    cleaned = _strip_sql_comments(sql).strip()
    if not cleaned:
        raise ValueError("SQL 不能为空")
    cleaned = cleaned.rstrip(";\n\t ")
    if ";" in cleaned:
        raise ValueError("当前 SQL Runner 仅允许单条只读 SQL 语句")
    first_keyword_match = re.match(r"^([A-Za-z_][A-Za-z0-9_]*)\b", cleaned)
    first_keyword = first_keyword_match.group(1).lower() if first_keyword_match else ""
    if first_keyword not in READONLY_SQL_START_KEYWORDS:
        raise ValueError("当前 SQL Runner 仅允许 SELECT / WITH / EXPLAIN 等只读查询")
    tokens = set(re.findall(r"\b[A-Za-z_][A-Za-z0-9_]*\b", cleaned.lower()))
    blocked = sorted(tokens & FORBIDDEN_SQL_KEYWORDS)
    if blocked:
        raise ValueError(f"SQL 包含禁止的写入或管理关键字: {', '.join(blocked)}")
    return cleaned


def execute_sql_safe(
    sql: str,
    workspace_dir: str,
    db_source: Optional[Dict[str, Any]],
    max_rows: int = 50000,
    timeout_sec: int = 120,
) -> str:
    """Execute a single read-only SQL statement and materialize returned rows."""
    if not db_source or not isinstance(db_source, dict):
        return "[SQL Error]: 当前分析未选择可执行 SQL 的数据库数据源"

    try:
        safe_sql = _validate_readonly_sql(sql)
        db_type = normalize_db_type(db_source.get("db_type") or db_source.get("dbType"))
        config = db_source.get("config") if isinstance(db_source.get("config"), dict) else {}
        source_label = str(db_source.get("label") or db_source.get("id") or db_type)
        engine = build_db_engine(db_type, config)
        started_at = time_module.time()
        rows = []
        columns: List[str] = []
        truncated = False

        with engine.connect() as conn:
            try:
                if db_type == "postgresql":
                    conn.execute(text(f"SET statement_timeout = {int(timeout_sec * 1000)}"))
                elif db_type == "mysql":
                    conn.execute(text(f"SET SESSION MAX_EXECUTION_TIME={int(timeout_sec * 1000)}"))
            except Exception:
                pass

            result = conn.execute(text(safe_sql))
            if result.returns_rows:
                columns = list(result.keys())
                fetched = result.fetchmany(max_rows + 1)
                if len(fetched) > max_rows:
                    truncated = True
                    fetched = fetched[:max_rows]
                rows = [tuple(row) for row in fetched]

        engine.dispose()

        generated_dir = Path(workspace_dir) / "generated"
        generated_dir.mkdir(parents=True, exist_ok=True)
        safe_source = re.sub(r"[^A-Za-z0-9_.-]+", "_", source_label).strip("_") or "database"
        output_path = uniquify_path(generated_dir / f"sql_result_{safe_source}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv")
        dataframe = pd.DataFrame(rows, columns=columns)
        dataframe.to_csv(output_path, index=False, encoding="utf-8-sig")

        elapsed_ms = int((time_module.time() - started_at) * 1000)
        preview_lines: List[str] = []
        if columns:
            preview_lines.append(" | ".join(map(str, columns[:12])))
            for row in rows[:8]:
                preview_lines.append(" | ".join(_truncate_history_text(str(item), 80) for item in row[:12]))
        preview_text = "\n".join(preview_lines) if preview_lines else "(query returned no rows)"
        return (
            f"[SQL OK] source={source_label}, dialect={db_type}, rows={len(rows)}, "
            f"columns={len(columns)}, truncated={str(truncated).lower()}, duration_ms={elapsed_ms}\n"
            f"[SQL Output] {output_path.name}\n"
            f"[SQL Preview]\n{preview_text}"
        )
    except Exception as exc:
        return f"[SQL Error]: {str(exc)}"


def execute_r_code_safe(code_str: str, workspace_dir: str = None, timeout_sec: int = 120) -> str:
    """Execute R code through Rscript in the session workspace."""
    if shutil.which("Rscript") is None:
        return "[R Error]: Rscript 未安装或不在 PATH 中，无法执行 R 代码"
    if workspace_dir is None:
        workspace_dir = WORKSPACE_BASE_DIR
    exec_cwd = os.path.abspath(workspace_dir)
    os.makedirs(exec_cwd, exist_ok=True)
    tmp_path = None
    try:
        fd, tmp_path = tempfile.mkstemp(suffix=".R", dir=exec_cwd)
        os.close(fd)
        with open(tmp_path, "w", encoding="utf-8") as f:
            f.write(code_str)
        child_env = os.environ.copy()
        child_env.setdefault("R_DEFAULT_DEVICE", "pdf")
        completed = subprocess.run(
            ["Rscript", tmp_path],
            cwd=exec_cwd,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=timeout_sec,
            env=child_env,
        )
        return (completed.stdout or "") + (completed.stderr or "")
    except subprocess.TimeoutExpired:
        return f"[R Timeout]: execution exceeded {timeout_sec} seconds"
    except Exception as exc:
        return f"[R Error]: {str(exc)}"
    finally:
        try:
            if tmp_path and os.path.exists(tmp_path):
                os.remove(tmp_path)
        except Exception:
            pass


# API endpoint and model path
DEFAULT_COMPUTE_BACKEND = (os.getenv("DEEPANALYZE_COMPUTE_BACKEND", "auto") or "auto").strip().lower()
API_BASE = os.getenv("DEEPANALYZE_DEFAULT_MODEL_BASE_URL", "http://localhost:8000/v1")
MODEL_PATH = os.getenv("DEEPANALYZE_DEFAULT_MODEL_NAME", "DeepAnalyze-8B")
DEFAULT_PROVIDER_TYPE = (os.getenv("DEEPANALYZE_DEFAULT_PROVIDER_TYPE", "deepanalyze") or "deepanalyze").strip()
DEFAULT_PROVIDER_LABEL = (os.getenv("DEEPANALYZE_DEFAULT_PROVIDER_LABEL", "DeepAnalyze 默认") or "DeepAnalyze 默认").strip()
DEFAULT_PROVIDER_DESCRIPTION = (os.getenv("DEEPANALYZE_DEFAULT_PROVIDER_DESCRIPTION", "项目默认本地 vLLM 服务") or "项目默认本地 vLLM 服务").strip()
DEFAULT_PROVIDER_API_KEY = (os.getenv("DEEPANALYZE_DEFAULT_MODEL_API_KEY", "") or "").strip()
MAX_AGENT_ROUNDS = 30
LLM_RETRY_MAX_ATTEMPTS = 3
LLM_RETRY_BACKOFF_SECONDS = [0.4, 0.8]
LLM_RETRYABLE_STATUS_CODES = {429, 500, 502, 503, 504}


# Initialize OpenAI client
client = openai.OpenAI(base_url=API_BASE, api_key=DEFAULT_PROVIDER_API_KEY or "dummy")
DEFAULT_MODEL_PROVIDER_CONFIG = {
    "id": "deepanalyze-default",
    "providerType": DEFAULT_PROVIDER_TYPE,
    "label": DEFAULT_PROVIDER_LABEL,
    "description": DEFAULT_PROVIDER_DESCRIPTION,
    "baseUrl": API_BASE,
    "model": MODEL_PATH,
    "apiKey": DEFAULT_PROVIDER_API_KEY,
    "headers": {},
    "supportsOpenAICompatible": True,
}

# Workspace directory
REPO_ROOT_DIR = os.path.abspath(os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", ".."))
WORKSPACE_BASE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "workspace")
PROJECTS_BASE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "projects")
CONFIG_BASE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "configs")
DB_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "deepanalyze.db")

# Built-in superuser credentials can be configured via environment variables.
# Default keeps backward compatibility with existing empty-password local setups.
BUILTIN_SUPERUSER_USERNAME = os.getenv("DEEPANALYZE_SUPERUSER_USERNAME", "rainforgrain").strip() or "rainforgrain"
BUILTIN_SUPERUSER_PASSWORD_HASH = os.getenv("DEEPANALYZE_SUPERUSER_PASSWORD_HASH", "").strip()
if not BUILTIN_SUPERUSER_PASSWORD_HASH:
    _builtin_superuser_password = os.getenv("DEEPANALYZE_SUPERUSER_PASSWORD")
    if _builtin_superuser_password is not None:
        BUILTIN_SUPERUSER_PASSWORD_HASH = hashlib.sha256(_builtin_superuser_password.encode()).hexdigest()


def is_builtin_superuser(username: Optional[str]) -> bool:
    return (username or "").strip() == BUILTIN_SUPERUSER_USERNAME


def verify_builtin_superuser_password(password: str) -> bool:
    if BUILTIN_SUPERUSER_PASSWORD_HASH:
        return hash_password(password or "") == BUILTIN_SUPERUSER_PASSWORD_HASH
    return (password or "") == ""


def ensure_builtin_superuser(cursor: sqlite3.Cursor) -> None:
    cursor.execute(
        "SELECT username, password_hash FROM users WHERE username = ?",
        (BUILTIN_SUPERUSER_USERNAME,),
    )
    row = cursor.fetchone()
    if row is None:
        cursor.execute(
            "INSERT INTO users (username, password_hash) VALUES (?, ?)",
            (BUILTIN_SUPERUSER_USERNAME, BUILTIN_SUPERUSER_PASSWORD_HASH),
        )
        return

    stored_hash = row["password_hash"] if isinstance(row, sqlite3.Row) else row[1]
    if stored_hash != BUILTIN_SUPERUSER_PASSWORD_HASH:
        cursor.execute(
            "UPDATE users SET password_hash = ? WHERE username = ?",
            (BUILTIN_SUPERUSER_PASSWORD_HASH, BUILTIN_SUPERUSER_USERNAME),
        )


def _get_user_config_dir(username: str) -> str:
    d = os.path.join(CONFIG_BASE_DIR, username)
    os.makedirs(d, exist_ok=True)
    return d


def _load_user_config(username: str, filename: str, default=None):
    path = os.path.join(_get_user_config_dir(username), filename)
    if os.path.exists(path):
        try:
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return default if default is not None else {}


def _save_user_config(username: str, filename: str, data) -> str:
    path = os.path.join(_get_user_config_dir(username), filename)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    return path


ANALYSIS_HISTORY_DEFAULT_SETTINGS = {
    "enabled": True,
    "capture_stream_progress": True,
    "capture_prompt_preview": True,
    "max_runs": 120,
    "stream_progress_chunk_interval": 40,
    "stream_progress_char_interval": 1600,
}


def _get_analysis_history_dir(username: str) -> str:
    path = os.path.join(_get_user_config_dir(username), "analysis_history")
    os.makedirs(path, exist_ok=True)
    return path


def _get_analysis_history_index_path(username: str) -> str:
    return os.path.join(_get_analysis_history_dir(username), "index.json")


def _get_analysis_history_run_path(username: str, run_id: str) -> str:
    safe_run_id = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(run_id or "run")).strip("_") or f"run_{int(time_module.time())}"
    return os.path.join(_get_analysis_history_dir(username), f"{safe_run_id}.jsonl")


def _load_analysis_history_settings(username: str) -> Dict[str, Any]:
    saved = _load_user_config(username, "analysis_history_settings.json", {})
    merged = dict(ANALYSIS_HISTORY_DEFAULT_SETTINGS)
    if isinstance(saved, dict):
        for key, default_value in ANALYSIS_HISTORY_DEFAULT_SETTINGS.items():
            value = saved.get(key, default_value)
            if isinstance(default_value, bool):
                merged[key] = bool(value)
            elif isinstance(default_value, int):
                try:
                    merged[key] = max(1, int(value))
                except Exception:
                    merged[key] = default_value
            else:
                merged[key] = value
    return merged


def _sanitize_analysis_history_settings(settings: Dict[str, Any]) -> Dict[str, Any]:
    incoming = settings if isinstance(settings, dict) else {}
    sanitized = dict(ANALYSIS_HISTORY_DEFAULT_SETTINGS)
    for key, default_value in ANALYSIS_HISTORY_DEFAULT_SETTINGS.items():
        value = incoming.get(key, default_value)
        if isinstance(default_value, bool):
            sanitized[key] = bool(value)
        elif isinstance(default_value, int):
            try:
                sanitized[key] = max(1, int(value))
            except Exception:
                sanitized[key] = default_value
        else:
            sanitized[key] = value
    return sanitized


def _load_analysis_history_index(username: str) -> Dict[str, Any]:
    path = _get_analysis_history_index_path(username)
    if os.path.exists(path):
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
            if isinstance(data, dict) and isinstance(data.get("runs"), list):
                return data
        except Exception:
            pass
    return {"runs": []}


def _save_analysis_history_index(username: str, data: Dict[str, Any]) -> str:
    path = _get_analysis_history_index_path(username)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    return path


def _upsert_analysis_history_summary(username: str, summary: Dict[str, Any], max_runs: int = 120) -> None:
    if not isinstance(summary, dict) or not summary.get("run_id"):
        return
    data = _load_analysis_history_index(username)
    runs = [item for item in data.get("runs", []) if isinstance(item, dict) and item.get("run_id") != summary.get("run_id")]
    runs.append(summary)
    runs.sort(key=lambda item: str(item.get("started_at") or item.get("updated_at") or ""), reverse=True)
    data["runs"] = runs[:max(1, int(max_runs or 120))]
    _save_analysis_history_index(username, data)


def _truncate_history_text(text: Any, max_chars: int = 600) -> str:
    raw = str(text or "")
    if len(raw) <= max_chars:
        return raw
    return raw[:max_chars] + "..."


def _sanitize_runtime_db_source_for_history(source: Dict[str, Any]) -> Dict[str, Any]:
    if not isinstance(source, dict):
        return {}
    config = source.get("config") if isinstance(source.get("config"), dict) else {}
    return {
        "id": str(source.get("id", "") or ""),
        "label": str(source.get("label", "") or ""),
        "db_type": str(source.get("db_type") or source.get("dbType") or ""),
        "config": {
            "host": str(config.get("host", "") or ""),
            "port": str(config.get("port", "") or ""),
            "user": str(config.get("user", "") or ""),
            "database": str(config.get("database", "") or ""),
        },
    }


def _sanitize_model_provider_for_history(model_provider: Any) -> Dict[str, Any]:
    if not isinstance(model_provider, dict):
        return {}
    return {
        "id": str(model_provider.get("id", "") or ""),
        "label": str(model_provider.get("label", "") or ""),
        "providerType": str(model_provider.get("providerType", "") or ""),
        "model": str(model_provider.get("model", "") or ""),
        "baseUrl": str(model_provider.get("baseUrl", "") or ""),
    }


class AnalysisHistoryRecorder:
    def __init__(
        self,
        username: str,
        session_id: str,
        settings: Dict[str, Any],
        request_summary: Dict[str, Any],
    ):
        self.username = str(username or "default")
        self.session_id = str(session_id or "default")
        self.settings = _sanitize_analysis_history_settings(settings)
        self.enabled = bool(self.settings.get("enabled", True))
        self.started_at_ts = time_module.time()
        self.started_at = datetime.now().isoformat(timespec="seconds")
        self.run_id = f"{self.session_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{secrets.token_hex(4)}"
        self.path = _get_analysis_history_run_path(self.username, self.run_id)
        self.event_count = 0
        self.last_stage = "session"
        self.last_event = "created"
        self.current_status = "running"
        self.summary = {
            "run_id": self.run_id,
            "session_id": self.session_id,
            "username": self.username,
            "status": "running",
            "started_at": self.started_at,
            "updated_at": self.started_at,
            "duration_ms": 0,
            "event_count": 0,
            "last_stage": "session",
            "last_event": "created",
            "last_message": "analysis run created",
            "request_summary": request_summary if isinstance(request_summary, dict) else {},
        }
        if self.enabled:
            Path(self.path).parent.mkdir(parents=True, exist_ok=True)
            self.log(
                stage="session",
                event="session_started",
                status="running",
                message="analysis session started",
                details=request_summary,
            )

    def _persist_summary(self) -> None:
        if not self.enabled:
            return
        self.summary["updated_at"] = datetime.now().isoformat(timespec="seconds")
        self.summary["duration_ms"] = int((time_module.time() - self.started_at_ts) * 1000)
        self.summary["event_count"] = self.event_count
        self.summary["last_stage"] = self.last_stage
        self.summary["last_event"] = self.last_event
        self.summary["status"] = self.current_status
        _upsert_analysis_history_summary(
            self.username,
            dict(self.summary),
            max_runs=self.settings.get("max_runs", ANALYSIS_HISTORY_DEFAULT_SETTINGS["max_runs"]),
        )

    def log(
        self,
        stage: str,
        event: str,
        status: str = "info",
        message: str = "",
        details: Optional[Dict[str, Any]] = None,
    ) -> None:
        if not self.enabled:
            return
        self.event_count += 1
        self.last_stage = str(stage or "session")
        self.last_event = str(event or "event")
        self.current_status = status if status in {"running", "info", "completed", "failed", "warning"} else self.current_status
        payload = {
            "run_id": self.run_id,
            "sequence": self.event_count,
            "timestamp": datetime.now().isoformat(timespec="seconds"),
            "elapsed_ms": int((time_module.time() - self.started_at_ts) * 1000),
            "stage": self.last_stage,
            "event": self.last_event,
            "status": status,
            "message": message,
            "details": details or {},
        }
        with open(self.path, "a", encoding="utf-8") as f:
            f.write(json.dumps(payload, ensure_ascii=False) + "\n")
        if message:
            self.summary["last_message"] = _truncate_history_text(message, max_chars=220)
        if status in {"warning", "failed"}:
            self.summary["last_problem"] = _truncate_history_text(message or json.dumps(details or {}, ensure_ascii=False), max_chars=320)
        self._persist_summary()

    def finalize(self, status: str, message: str = "", details: Optional[Dict[str, Any]] = None) -> None:
        if not self.enabled:
            return
        normalized_status = status if status in {"completed", "failed", "warning"} else "completed"
        self.current_status = normalized_status
        self.log(
            stage="session",
            event="session_completed" if normalized_status == "completed" else "session_finished_with_issues",
            status=normalized_status,
            message=message,
            details=details or {},
        )
KB_MASK = "••••••••"
KB_DEFAULT_SETTINGS = {
    "knowledge_base_enabled": True,
    "providers_enabled": True,
    "internal_preferences": {
        "preferred_view": "html",
        "show_hints": True,
        "auto_open_yutu_after_analysis": False,
    },
    "onyx": {
        "enabled": False,
        "base_url": "http://localhost:3000",
        "api_key": "",
        "search_path": "/api/chat/search",
    },
    "dify": {
        "enabled": False,
        "base_url": "http://localhost:5000",
        "api_key": "",
        "workflow_id": "",
    },
    "test_status": {
        "onyx": {
            "status": "never_tested",
            "message": "尚未测试",
            "tested_at": None,
        },
        "dify": {
            "status": "never_tested",
            "message": "尚未测试",
            "tested_at": None,
        },
    },
}


def init_db():
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    # Users table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS users (
            username TEXT PRIMARY KEY,
            password_hash TEXT NOT NULL
        )
    ''')
    ensure_builtin_superuser(cursor)
    # Projects table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS projects (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT NOT NULL,
            session_id TEXT NOT NULL,
            name TEXT NOT NULL,
            messages TEXT NOT NULL,
            files_data TEXT DEFAULT '{}',
            side_tasks TEXT DEFAULT '[]',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (username) REFERENCES users(username)
        )
    ''')
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS knowledge_provider_settings (
            id INTEGER PRIMARY KEY CHECK (id = 1),
            settings_json TEXT NOT NULL,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    # 兼容旧数据库：projects 表可能已存在但缺少 files_data 或 side_tasks 列
    cursor.execute("PRAGMA table_info(projects)")
    columns = [row[1] for row in cursor.fetchall()]
    if "files_data" not in columns:
        cursor.execute("ALTER TABLE projects ADD COLUMN files_data TEXT DEFAULT '{}'")
    if "side_tasks" not in columns:
        cursor.execute("ALTER TABLE projects ADD COLUMN side_tasks TEXT DEFAULT '[]'")
    cursor.execute("SELECT id FROM knowledge_provider_settings WHERE id = 1")
    if cursor.fetchone() is None:
        cursor.execute(
            "INSERT INTO knowledge_provider_settings (id, settings_json) VALUES (1, ?)",
            (json.dumps(KB_DEFAULT_SETTINGS, ensure_ascii=False),),
        )
    conn.commit()
    conn.close()

init_db()

def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        yield conn
    finally:
        conn.close()


def deep_merge_dict(base: Dict[str, Any], incoming: Dict[str, Any]) -> Dict[str, Any]:
    merged = deepcopy(base)
    for key, value in incoming.items():
        if isinstance(value, dict) and isinstance(merged.get(key), dict):
            merged[key] = deep_merge_dict(merged[key], value)
        else:
            merged[key] = value
    return merged


def normalize_base_url(value: str) -> str:
    if not isinstance(value, str) or not value.strip():
        raise ValueError("base_url 不能为空")
    base_url = value.strip().rstrip("/")
    if not re.match(r"^https?://", base_url, re.IGNORECASE):
        raise ValueError("base_url 必须以 http:// 或 https:// 开头")
    return base_url


def _is_local_model_endpoint(base_url: str) -> bool:
    lowered = (base_url or "").strip().lower()
    return "localhost:" in lowered or "127.0.0.1:" in lowered


def _resolve_local_mlx_model_name(model_name: str, base_url: str) -> str:
    text = (model_name or "").strip()
    if not text:
        return text

    if DEFAULT_COMPUTE_BACKEND not in {"mlx", "apple", "apple_silicon"}:
        return text

    if not _is_local_model_endpoint(base_url):
        return text

    if os.path.isabs(text):
        return text

    candidate = os.path.abspath(os.path.join(REPO_ROOT_DIR, text))
    if os.path.isdir(candidate):
        return candidate

    configured_mlx_dir = (os.getenv("DEEPANALYZE_MLX_MODEL_DIR", "") or "").strip()
    if configured_mlx_dir and os.path.isdir(configured_mlx_dir):
        configured_name = os.path.basename(configured_mlx_dir.rstrip(os.sep))
        if configured_name == text:
            return configured_mlx_dir

    return text


def normalize_optional_path(value: str, default: str) -> str:
    path = (value or default).strip()
    if not path.startswith("/"):
        path = f"/{path}"
    return path


def mask_secret(value: str) -> str:
    if not value:
        return ""
    return KB_MASK


def redact_provider(provider: Dict[str, Any]) -> Dict[str, Any]:
    masked = deepcopy(provider)
    api_key = masked.get("api_key", "")
    masked["has_api_key"] = bool(api_key)
    masked["api_key"] = mask_secret(api_key)
    return masked


def normalize_model_provider_config(payload: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    incoming = payload if isinstance(payload, dict) else {}
    merged = deep_merge_dict(DEFAULT_MODEL_PROVIDER_CONFIG, incoming)
    merged["providerType"] = (merged.get("providerType") or DEFAULT_PROVIDER_TYPE).strip() or DEFAULT_PROVIDER_TYPE
    merged["baseUrl"] = normalize_base_url(merged.get("baseUrl") or API_BASE)
    raw_model = (merged.get("model") or MODEL_PATH).strip() or MODEL_PATH
    merged["model"] = _resolve_local_mlx_model_name(raw_model, merged["baseUrl"])
    merged["apiKey"] = (merged.get("apiKey") or "").strip()
    raw_headers = merged.get("headers") or {}
    normalized_headers: Dict[str, str] = {}
    if isinstance(raw_headers, dict):
        for key, value in raw_headers.items():
            if key is None or value is None:
                continue
            key_text = str(key).strip()
            value_text = str(value).strip()
            if key_text and value_text:
                normalized_headers[key_text] = value_text
    merged["headers"] = normalized_headers
    return merged


def build_model_provider_headers(config: Dict[str, Any]) -> Dict[str, str]:
    headers = {**(config.get("headers") or {})}
    api_key = config.get("apiKey", "")
    has_auth_header = any(key.lower() in {"authorization", "x-api-key", "api-key"} for key in headers)
    if api_key and not has_auth_header:
        headers["Authorization"] = f"Bearer {api_key}"
    return headers


def get_runtime_llm(model_provider: Optional[Dict[str, Any]] = None) -> Tuple[openai.OpenAI, str, Dict[str, Any]]:
    normalized = normalize_model_provider_config(model_provider)
    extra_headers = {
        key: value
        for key, value in (normalized.get("headers") or {}).items()
        if key.lower() != "authorization"
    }
    runtime_client = openai.OpenAI(
        base_url=normalized["baseUrl"],
        api_key=normalized.get("apiKey") or "dummy",
        default_headers=extra_headers or None,
    )
    return runtime_client, normalized["model"], normalized


def _provider_supports_vllm_controls(provider_cfg: Optional[Dict[str, Any]]) -> bool:
    if not isinstance(provider_cfg, dict):
        return False
    provider_type = str(provider_cfg.get("providerType") or "").strip().lower()
    if provider_type in {"deepanalyze", "vllm"}:
        return True
    base_url = str(provider_cfg.get("baseUrl") or "").strip().lower()
    return "localhost:8000" in base_url or "127.0.0.1:8000" in base_url


def _build_streaming_answer_error(message: str) -> str:
    safe_message = (str(message or "未知错误").strip() or "未知错误")[:1000]
    return (
        "<Answer>\n"
        "模型服务调用失败，当前分析已中止。\n"
        f"错误详情：{safe_message}\n"
        "建议：\n"
        "1. 检查模型服务是否正常（例如本地 vLLM 端口 8000）。\n"
        "2. 若使用第三方模型，请检查 baseUrl、模型名与 API Key。\n"
        "3. 修复后重新发送同一问题继续分析。\n"
        "</Answer>\n"
    )


def _extract_llm_error_status(exc: Exception) -> Optional[int]:
    for attr in ("status_code", "status", "http_status"):
        value = getattr(exc, attr, None)
        if isinstance(value, int):
            return value
    response = getattr(exc, "response", None)
    status_code = getattr(response, "status_code", None)
    if isinstance(status_code, int):
        return status_code
    msg = str(exc)
    match = re.search(r"\b(429|500|502|503|504)\b", msg)
    if match:
        return int(match.group(1))
    return None


def _is_retryable_llm_exception(exc: Exception) -> bool:
    known_retryable_types = tuple(
        exc_type
        for exc_type in (
            getattr(openai, "APIConnectionError", None),
            getattr(openai, "APITimeoutError", None),
            getattr(openai, "RateLimitError", None),
            getattr(openai, "InternalServerError", None),
        )
        if exc_type is not None
    )
    if known_retryable_types and isinstance(exc, known_retryable_types):
        return True
    status_code = _extract_llm_error_status(exc)
    if status_code in LLM_RETRYABLE_STATUS_CODES:
        return True
    msg = str(exc).lower()
    return any(
        key in msg
        for key in (
            "timed out",
            "timeout",
            "temporarily",
            "overload",
            "connection",
            "service unavailable",
            "bad gateway",
        )
    )


def _create_chat_completion_with_retry(llm_client: openai.OpenAI, request_kwargs: Dict[str, Any]):
    active_kwargs = dict(request_kwargs)
    dropped_extra_body = False
    for attempt in range(1, LLM_RETRY_MAX_ATTEMPTS + 1):
        try:
            return llm_client.chat.completions.create(**active_kwargs)
        except Exception as exc:
            if "extra_body" in active_kwargs and not dropped_extra_body:
                active_kwargs = dict(active_kwargs)
                active_kwargs.pop("extra_body", None)
                dropped_extra_body = True
                print("[LLM] Request failed with extra_body, retrying without provider-specific controls")
                continue

            if attempt >= LLM_RETRY_MAX_ATTEMPTS or not _is_retryable_llm_exception(exc):
                raise

            delay = LLM_RETRY_BACKOFF_SECONDS[min(attempt - 1, len(LLM_RETRY_BACKOFF_SECONDS) - 1)]
            status_code = _extract_llm_error_status(exc)
            print(f"[LLM] Transient error (status={status_code}) on attempt {attempt}, retrying in {delay:.1f}s")
            time_module.sleep(delay)

    raise RuntimeError("LLM request failed after retries")


def fetch_model_names(model_provider: Optional[Dict[str, Any]] = None) -> List[str]:
    config = normalize_model_provider_config(model_provider)
    response = httpx.get(
        f"{config['baseUrl']}/models",
        headers=build_model_provider_headers(config),
        timeout=20.0,
    )
    response.raise_for_status()
    payload = response.json()
    models: List[str] = []
    if isinstance(payload, dict):
        data = payload.get("data", [])
        if isinstance(data, list):
            for item in data:
                if isinstance(item, dict):
                    model_name = (item.get("id") or item.get("name") or "").strip()
                    if model_name:
                        models.append(model_name)
    if config["model"] and config["model"] not in models:
        models.insert(0, config["model"])
    deduped: List[str] = []
    seen = set()
    for model_name in models:
        if model_name not in seen:
            seen.add(model_name)
            deduped.append(model_name)
    return deduped


def build_kb_settings_response(settings: Dict[str, Any]) -> Dict[str, Any]:
    response = deepcopy(settings)
    response["onyx"] = redact_provider(response.get("onyx", {}))
    response["dify"] = redact_provider(response.get("dify", {}))
    return response


def load_kb_settings(raw: bool = False) -> Dict[str, Any]:
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT settings_json FROM knowledge_provider_settings WHERE id = 1")
        row = cursor.fetchone()
    finally:
        conn.close()
    if row and row["settings_json"]:
        try:
            stored = json.loads(row["settings_json"])
        except json.JSONDecodeError:
            stored = {}
    else:
        stored = {}
    merged = deep_merge_dict(KB_DEFAULT_SETTINGS, stored)
    return merged if raw else build_kb_settings_response(merged)


def save_kb_settings(settings: Dict[str, Any]) -> None:
    conn = sqlite3.connect(DB_PATH)
    try:
        cursor = conn.cursor()
        cursor.execute(
            "UPDATE knowledge_provider_settings SET settings_json = ?, updated_at = CURRENT_TIMESTAMP WHERE id = 1",
            (json.dumps(settings, ensure_ascii=False),),
        )
        conn.commit()
    finally:
        conn.close()


def merge_provider_input(existing: Dict[str, Any], incoming: Dict[str, Any]) -> Dict[str, Any]:
    merged = deepcopy(existing)
    for key, value in incoming.items():
        if key == "api_key" and isinstance(value, str) and value == KB_MASK:
            continue
        merged[key] = value
    return merged


def normalize_kb_settings_input(payload: Dict[str, Any], existing: Dict[str, Any]) -> Dict[str, Any]:
    settings = deepcopy(existing)
    if "knowledge_base_enabled" in payload:
        settings["knowledge_base_enabled"] = bool(payload.get("knowledge_base_enabled"))
    if "providers_enabled" in payload:
        settings["providers_enabled"] = bool(payload.get("providers_enabled"))

    if isinstance(payload.get("internal_preferences"), dict):
        incoming_preferences = payload["internal_preferences"]
        prefs = settings.get("internal_preferences", {})
        if incoming_preferences.get("preferred_view") in {"html", "table"}:
            prefs["preferred_view"] = incoming_preferences["preferred_view"]
        if "show_hints" in incoming_preferences:
            prefs["show_hints"] = bool(incoming_preferences.get("show_hints"))
        if "auto_open_yutu_after_analysis" in incoming_preferences:
            prefs["auto_open_yutu_after_analysis"] = bool(incoming_preferences.get("auto_open_yutu_after_analysis"))
        settings["internal_preferences"] = prefs

    if isinstance(payload.get("onyx"), dict):
        onyx = merge_provider_input(settings.get("onyx", {}), payload["onyx"])
        onyx["enabled"] = bool(onyx.get("enabled"))
        onyx["base_url"] = normalize_base_url(onyx.get("base_url", ""))
        onyx["search_path"] = normalize_optional_path(onyx.get("search_path", ""), "/api/chat/search")
        onyx["api_key"] = (onyx.get("api_key") or "").strip()
        settings["onyx"] = onyx

    if isinstance(payload.get("dify"), dict):
        dify = merge_provider_input(settings.get("dify", {}), payload["dify"])
        dify["enabled"] = bool(dify.get("enabled"))
        dify["base_url"] = normalize_base_url(dify.get("base_url", ""))
        dify["api_key"] = (dify.get("api_key") or "").strip()
        dify["workflow_id"] = (dify.get("workflow_id") or "").strip()
        settings["dify"] = dify

    return settings


def normalize_http_error(error: Exception) -> str:
    if isinstance(error, httpx.ConnectTimeout) or isinstance(error, httpx.ReadTimeout):
        return "连接超时，请检查服务是否启动或地址是否可达"
    if isinstance(error, httpx.ConnectError):
        return "连接被拒绝，请确认本地 Docker 服务已启动且端口映射正确"
    if isinstance(error, httpx.HTTPStatusError):
        status_code = error.response.status_code
        if status_code in {401, 403}:
            return f"鉴权失败（HTTP {status_code}），请检查 API Key"
        if status_code == 404:
            return "接口路径不存在，请检查 base_url、检索路径或 workflow_id"
        return f"服务返回异常状态（HTTP {status_code}）"
    return str(error)


def test_onyx_provider(config: Dict[str, Any]) -> Tuple[bool, str]:
    base_url = normalize_base_url(config.get("base_url", ""))
    configured_search_path = normalize_optional_path(config.get("search_path", ""), "/api/chat/search")
    candidate_search_paths = []
    for path in [configured_search_path, "/api/chat/search", "/api/search", "/search"]:
        normalized = normalize_optional_path(path, configured_search_path)
        if normalized not in candidate_search_paths:
            candidate_search_paths.append(normalized)
    api_key = (config.get("api_key") or "").strip()
    base_headers = {"Accept": "application/json"}
    auth_header_variants = [base_headers]
    if api_key:
        auth_header_variants = [
            {**base_headers, "Authorization": f"Bearer {api_key}"},
            {**base_headers, "X-API-Key": api_key},
            {**base_headers, "Api-Key": api_key},
        ]
    health_candidates = ["/api/health", "/health", "/api/chat/health", "/"]
    last_error: Optional[Exception] = None

    with httpx.Client(timeout=8.0, follow_redirects=True) as client:
        health_ok = False
        for health_path in health_candidates:
            health_url = f"{base_url}{health_path}" if health_path != "/" else f"{base_url}/"
            for headers in auth_header_variants:
                try:
                    health_response = client.get(health_url, headers=headers)
                    if health_response.status_code < 500:
                        health_ok = True
                        break
                except Exception as error:
                    last_error = error
            if health_ok:
                break
        if not health_ok and last_error:
            raise last_error

        search_payloads = [
            {"query": "ping", "search_type": "keyword"},
            {"query": "ping"},
            {"message": "ping"},
            {"text": "ping"},
        ]
        for search_path in candidate_search_paths:
            search_url = f"{base_url}{search_path}"
            for headers in auth_header_variants:
                request_headers = {**headers, "Content-Type": "application/json"}
                for payload in search_payloads:
                    try:
                        response = client.post(search_url, headers=request_headers, json=payload)
                        if response.status_code in {401, 403}:
                            response.raise_for_status()
                        if response.status_code in {404, 405, 422}:
                            last_error = httpx.HTTPStatusError(
                                "Onyx search endpoint rejected request",
                                request=response.request,
                                response=response,
                            )
                            continue
                        response.raise_for_status()
                        try:
                            result = response.json()
                        except ValueError as exc:
                            raise ValueError("Onyx 检索接口未返回 JSON") from exc
                        if not isinstance(result, (dict, list)):
                            raise ValueError("Onyx 检索接口返回结构异常")
                        return True, f"Onyx 服务可达，检索接口测试通过（{search_path}）"
                    except Exception as error:
                        last_error = error
        if last_error:
            raise last_error
    raise ValueError("Onyx 检索接口测试失败")


def test_dify_provider(config: Dict[str, Any]) -> Tuple[bool, str]:
    base_url = normalize_base_url(config.get("base_url", ""))
    api_key = (config.get("api_key") or "").strip()
    workflow_id = (config.get("workflow_id") or "").strip()
    if not api_key:
        raise ValueError("Dify API Key 不能为空")
    if not workflow_id:
        raise ValueError("Dify workflow_id 不能为空")
    run_url = f"{base_url}/v1/workflows/run"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
        "Accept": "application/json",
    }
    payload = {
        "inputs": {"ping": "health_check"},
        "response_mode": "blocking",
        "user": "deepanalyze-local-test",
        "workflow_id": workflow_id,
    }

    with httpx.Client(timeout=10.0, follow_redirects=True) as client:
        response = client.post(run_url, headers=headers, json=payload)
        response.raise_for_status()
        try:
            result = response.json()
        except ValueError as exc:
            raise ValueError("Dify Workflow 接口未返回 JSON") from exc
        if not isinstance(result, dict):
            raise ValueError("Dify Workflow 接口返回结构异常")
        if not any(key in result for key in ("data", "workflow_run_id", "task_id", "outputs")):
            raise ValueError("Dify Workflow 响应缺少预期字段")
    return True, "Dify Workflow 接口测试通过"


def run_kb_provider_test(provider: str, settings: Dict[str, Any]) -> Dict[str, Any]:
    tested_at = datetime.now().isoformat(timespec="seconds")
    try:
        if provider == "onyx":
            _, message = test_onyx_provider(settings.get("onyx", {}))
        elif provider == "dify":
            _, message = test_dify_provider(settings.get("dify", {}))
        else:
            raise ValueError(f"Unsupported provider: {provider}")
        return {
            "provider": provider,
            "success": True,
            "status": "passed",
            "message": message,
            "tested_at": tested_at,
        }
    except Exception as error:
        return {
            "provider": provider,
            "success": False,
            "status": "failed",
            "message": normalize_http_error(error),
            "tested_at": tested_at,
        }

def hash_password(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()


def normalize_analysis_language(raw_language: Optional[str]) -> str:
    normalized = str(raw_language or "").strip().lower().replace("_", "-")
    if normalized in {"en", "en-us", "en-gb", "english"}:
        return "en"
    return "zh-CN"


def build_analysis_language_prompt(analysis_language: str) -> str:
    if analysis_language == "en":
        return (
            "\n\n**Output language override (highest priority): English**"
            "\n- This preference overrides default language rules in the base prompt."
            "\n- All user-facing narrative text must be in English, including `<Analyze>`, `<Understand>`, `<Answer>`, and report body text."
            "\n- In interactive mode, task names/descriptions in `<TaskTree>` and all follow-up user guidance text must also be in English."
            "\n- Keep structural tags unchanged (`<Analyze>`, `<Code>`, `<TaskTree>`, etc.); only change natural-language content."
            "\n- Do not output Chinese characters in user-facing text. If any Chinese appears, rewrite it in English before continuing."
        )

    return (
        "\n\n**输出语言覆盖规则（最高优先级）：中文（简体）**"
        "\n- 本规则覆盖基础提示词中的默认语言描述。"
        "\n- 所有面向用户的自然语言内容必须使用简体中文，包括 `<Analyze>`、`<Understand>`、`<Answer>` 和最终报告正文。"
        "\n- 在交互模式下，`<TaskTree>` 中任务名称/描述与后续引导文案也必须使用简体中文。"
        "\n- 结构化标签保持不变（如 `<Analyze>`、`<Code>`、`<TaskTree>` 等），仅改变自然语言内容。"
    )

HTTP_SERVER_PORT = 8100
HTTP_SERVER_BASE = (
    f"http://localhost:{HTTP_SERVER_PORT}"  # you can replace localhost to your local ip
)


def get_session_workspace(session_id: str, username: str = "default") -> str:
    """返回指定 user 和 session 的 workspace 路径（workspace/{username}/{session_id}/）。"""
    if not username:
        username = "default"
    if not session_id:
        session_id = "default"
    session_dir = os.path.join(WORKSPACE_BASE_DIR, username, session_id)
    os.makedirs(session_dir, exist_ok=True)
    return session_dir


def build_download_url(rel_path: str) -> str:
    try:
        encoded = quote(rel_path, safe="/")
    except Exception:
        encoded = rel_path
    return f"{HTTP_SERVER_BASE}/{encoded}"


# FastAPI app
app = FastAPI()

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- v1 path rewrite map (frontend expects /v1/*, backend has legacy paths) ---
_V1_REWRITE = {
    "/v1/chat/completions":       "/chat/completions",
    "/v1/models":                 "/api/model/models",
    "/v1/files":                  "/workspace/files",
    "/v1/files/tree":             "/workspace/tree",
    "/v1/files/upload-to":        "/workspace/upload-to",
    # user config persistence
    "/v1/config/models":          "/api/config/models",
    "/v1/config/databases":       "/api/config/databases",
    "/v1/config/knowledge":       "/api/config/knowledge",
    "/v1/config/analysis-history": "/api/config/analysis-history",
    "/v1/config/export":          "/api/config/export",
    "/v1/knowledge/settings":     "/api/kb/settings",
    "/v1/knowledge/test":         "/api/kb/test",
    "/v1/knowledge/entries":      "/api/yutu/search",
    "/v1/knowledge/entries/search":"/api/yutu/search",
    # yutu nested routes: frontend kebab-case → backend slash-separated
    "/v1/knowledge/yutu/organize-confirm": "/api/yutu/organize/confirm",
    "/v1/knowledge/yutu/organize-cancel":  "/api/yutu/organize/cancel",
    "/v1/knowledge/yutu/backup-list":     "/api/yutu/backup/list",
    "/v1/knowledge/yutu/backup-create":   "/api/yutu/backup/create",
    "/v1/knowledge/yutu/backup-delete":   "/api/yutu/backup/delete",
    "/v1/knowledge/yutu/backup-restore":  "/api/yutu/backup/restore",
    "/v1/projects/save":          "/api/projects/save",
    "/v1/projects/list":          "/api/projects/list",
    "/v1/projects/load":          "/api/projects/load",
    "/v1/projects/check-name":    "/api/projects/check-name",
    "/v1/projects/restore-files": "/api/projects/restore-files",
    "/v1/projects/restore-to-workspace": "/api/projects/restore-to-workspace",
    "/v1/database/test":          "/api/db/test",
    "/v1/database/list":          "/api/db/list-databases",
    "/v1/database/context/load":  "/api/db/context/load",
    "/v1/database/schema/graph":  "/api/db/schema/graph",
    "/v1/database/generate-sql":  "/api/db/generate-sql",
    "/v1/database/execute":       "/api/db/execute",
    "/v1/data/profile-report":     "/api/data/profile-report",
    "/v1/export/report":          "/export/report",
    "/v1/analysis/history":       "/api/analysis/history",
    "/v1/chat/guidance":          "/api/chat/guidance",
    "/v1/code/execute":           "/execute",
}
_V1_PREFIX_REWRITE = {
    "/v1/knowledge/yutu/":  "/api/yutu/",
    "/v1/projects/":        "/api/projects/",
    "/v1/files/":           "/workspace/",
    "/v1/analysis/history/": "/api/analysis/history/",
}


class V1PathRewriteMiddleware:
    """Pure ASGI middleware that rewrites /v1/* paths to legacy equivalents."""

    def __init__(self, app):
        self.app = app

    async def __call__(self, scope, receive, send):
        if scope["type"] == "http":
            path = scope.get("path", "")
            # Exact match first
            if path in _V1_REWRITE:
                scope = {**scope, "path": _V1_REWRITE[path]}
            else:
                # Prefix match
                for v1_prefix, legacy_prefix in _V1_PREFIX_REWRITE.items():
                    if path.startswith(v1_prefix):
                        suffix = path[len(v1_prefix):]
                        scope = {**scope, "path": f"{legacy_prefix}{suffix}"}
                        break
                else:
                    # DELETE /v1/projects?id=... → /api/projects/delete
                    if path == "/v1/projects" and scope.get("method", "").upper() == "DELETE":
                        qs = scope.get("query_string", b"").decode("latin-1")
                        scope = {**scope, "path": "/api/projects/delete", "query_string": qs.encode("latin-1")}

        await self.app(scope, receive, send)


app.add_middleware(V1PathRewriteMiddleware)

@app.get("/proxy")
async def proxy_url(url: str = Query(...)):
    """代理外部 URL 以解决跨域问题（特别是本地文件服务器）"""
    async with httpx.AsyncClient() as client:
        try:
            resp = await client.get(url, follow_redirects=True)
            return Response(
                content=resp.content,
                status_code=resp.status_code,
                headers={k: v for k, v in resp.headers.items() if k.lower() not in ("content-encoding", "transfer-encoding", "content-length")}
            )
        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))


def start_http_server():
    """启动HTTP文件服务器（不修改全局工作目录）。"""
    os.makedirs(WORKSPACE_BASE_DIR, exist_ok=True)
    os.makedirs(PROJECTS_BASE_DIR, exist_ok=True)
    handler = partial(
        http.server.SimpleHTTPRequestHandler, directory=WORKSPACE_BASE_DIR
    )
    with socketserver.TCPServer(("", HTTP_SERVER_PORT), handler) as httpd:
        print(f"HTTP Server serving {WORKSPACE_BASE_DIR} at port {HTTP_SERVER_PORT}")
        httpd.serve_forever()


# Start HTTP server in a separate thread
threading.Thread(target=start_http_server, daemon=True).start()


# ---------- Settings & Hardware Detection ----------

@app.get("/api/settings/hardware")
async def settings_hardware():
    """检测当前硬件加速能力"""
    mlx_available = False
    cuda_available = False
    opencl_available = False
    directml_available = False
    gpu_info = "N/A"

    try:
        import mlx  # noqa: F401
        mlx_available = True
    except ImportError:
        pass

    try:
        import torch
        cuda_available = torch.cuda.is_available()
        if cuda_available:
            gpu_info = torch.cuda.get_device_name(0)
    except ImportError:
        pass

    try:
        import pyopencl  # noqa: F401
        opencl_available = True
    except ImportError:
        pass

    try:
        import torch_directml  # noqa: F401
        directml_available = True
    except ImportError:
        pass

    if mlx_available:
        recommended = "mlx"
    elif cuda_available:
        recommended = "gpu"
    else:
        recommended = "cpu"

    return {
        "mlx_available": mlx_available,
        "cuda_available": cuda_available,
        "opencl_available": opencl_available,
        "directml_available": directml_available,
        "gpu_info": gpu_info,
        "recommended": recommended,
    }


@app.get("/api/settings/defaults")
async def settings_defaults():
    """返回前端默认设置"""
    if DEFAULT_COMPUTE_BACKEND in {"mlx", "apple", "apple_silicon"}:
        model_version = "mlx"
    elif DEFAULT_COMPUTE_BACKEND in {"gpu", "cuda", "nvidia"}:
        model_version = "gpu"
    else:
        model_version = "mlx" if "mlx" in MODEL_PATH.lower() else "gpu"

    return {
        "analysis_mode": "full_agent",
        "analysis_strategy": "聚焦诉求",
        "temperature": None,
        "knowledge_base_enabled": True,
        "model_version": model_version,
        "compute_backend": DEFAULT_COMPUTE_BACKEND,
        "self_correction_enabled": True,
        "short_test_enabled": True,
        "task_decomposition_enabled": True,
        "explainability_enabled": True,
        "efficient_processing_enabled": True,
        "dead_loop_detection_enabled": True,
    }


@app.post("/api/model/models")
async def list_model_names(body: dict = Body(...)):
    try:
        config = body.get("config") or body.get("model_provider") or {}
        normalized = normalize_model_provider_config(config)
        models = fetch_model_names(normalized)
        return {
            "success": True,
            "models": models,
            "selected_model": normalized["model"],
            "message": f"已获取 {len(models)} 个模型名称",
            "tested_at": datetime.now().isoformat(timespec="seconds"),
        }
    except httpx.HTTPStatusError as e:
        detail = (e.response.text or str(e))[:300] if e.response is not None else str(e)
        return {"success": False, "message": f"模型列表获取失败: {detail}"}
    except Exception as e:
        print(f"List model names error: {e}")
        return {"success": False, "message": str(e)}


# ---------- Side Guidance Storage ----------
# session_id -> guidance_text
SESSION_GUIDANCE = {}

# session_id -> imported database context payload
SESSION_DB_CONTEXT: Dict[str, Dict[str, Any]] = {}

@app.post("/api/chat/guidance")
async def set_guidance(session_id: str = Query(...), guidance: dict = Body(...)):
    """
    接收用户提供的过程指导（Side Task），并在智能体下一步执行时注入。
    """
    text = guidance.get("guidance", "").strip()
    if text:
        # 如果该 session 已经有待处理的指导，则追加
        if session_id in SESSION_GUIDANCE:
            SESSION_GUIDANCE[session_id] += f"\n{text}"
        else:
            SESSION_GUIDANCE[session_id] = text
        print(f"[Side Guidance] Session {session_id} received: {text}")
    return {"status": "ok", "message": "Guidance received"}


def _build_filename_mapping(workspace_dir: str) -> dict:
    """
    构建原始文件名 → converted/ 目录中实际文件名的映射。
    例如：上传 data.csv → converted/data (1).csv，则映射 {"data.csv": "data (1).csv"}
    """
    converted_dir = Path(workspace_dir) / "converted"
    mapping = {}
    if not converted_dir.exists():
        return mapping

    # 获取 converted 目录中所有文件（不带路径）
    converted_names = {f.name for f in converted_dir.iterdir() if f.is_file()}
    if not converted_names:
        return mapping

    # 遍历根目录文件，若发现与 converted 同名（或同名+uniquify后），建立映射
    root_files = [f for f in Path(workspace_dir).iterdir() if f.is_file() and f.name not in (".lib",)]
    for root_file in root_files:
        if root_file.name in converted_names:
            # 完全同名（未被 uniquify）
            mapping[root_file.name] = root_file.name
        else:
            # 检查是否为某个 converted 文件的唯一化版本（通过 stem 比对）
            for cn in converted_names:
                cn_stem = Path(cn).stem
                if root_file.stem == cn_stem:
                    mapping[root_file.name] = cn
                    break
    return mapping


def _rewrite_file_paths(code_str: str, workspace_dir: str) -> str:
    """
    将代码中出现的原始文件名自动替换为 converted/ 目录下的实际文件名。
    支持：pd.read_csv("data.csv") → pd.read_csv("converted/data (1).csv")
    """
    mapping = _build_filename_mapping(workspace_dir)
    if not mapping:
        return code_str

    converted_dir = Path(workspace_dir).name  # 仅目录名，用于构造相对路径

    for orig_name, actual_name in mapping.items():
        if orig_name == actual_name:
            continue  # 同名无需替换
        # 匹配字符串参数中的原始文件名（各种常见读取函数）
        # 匹配模式：函数调用中以 orig_name 为字符串参数
        # 例如：read_csv("data.csv")  → read_csv("converted/data (1).csv")
        for func in ["read_csv", "read_table", "open", "load", "read_excel", "read_json", "read_xml", "read_clipboard"]:
            # 匹配 func("orig_name") 或 func('orig_name')
            for quote in ['"', "'"]:
                # 构造贪婪模式以避免通配符冲突
                escaped_orig = orig_name.replace("_", r"_").replace(".", r"\.")
                # 简单方式：直接替换字符串字面值
                target_literal = f"{quote}{escaped_orig}{quote}"
                replacement = f"{quote}{converted_dir}/{actual_name}{quote}"
                # 避免重复替换 converted/ 路径本身
                pattern_unsafe = f'(?<!/{converted_dir}/){quote}{escaped_orig}{quote}'
                try:
                    code_str = re.sub(
                        pattern_unsafe,
                        replacement,
                        code_str
                    )
                except re.error:
                    pass
    return code_str


def collect_file_info(directory: str) -> str:
    """收集文件信息，包括编码映射"""
    all_file_info_str = ""
    dir_path = Path(directory)
    if not dir_path.exists():
        return ""

    # 读取编码映射文件（如存在）
    encoding_map = {}
    encoding_map_path = dir_path / ".encoding_map.json"
    if encoding_map_path.exists():
        try:
            with open(encoding_map_path, "r", encoding="utf-8") as f:
                encoding_data = json.load(f)
                # 构建文件名 -> 编码信息的映射
                for item in encoding_data.get("files", []):
                    orig_name = item.get("original_name", "")
                    encoding_map[orig_name] = item
        except Exception:
            pass

    # 收集所有文件信息（排除 .encoding_map.json 本身）
    files = sorted([f for f in dir_path.iterdir() if f.is_file() and f.name != ".encoding_map.json"])

    for idx, file_path in enumerate(files, start=1):
        size_bytes = os.path.getsize(file_path)
        size_kb = size_bytes / 1024
        size_str = f"{size_kb:.1f}KB"

        # 获取该文件的编码信息
        enc_info = encoding_map.get(file_path.name, {})
        original_encoding = enc_info.get("original_encoding", "unknown")
        is_converted = enc_info.get("is_converted", False)
        converted_name = enc_info.get("converted_name", file_path.name)

        # 构建文件信息
        file_info = {
            "name": file_path.name,
            "size": size_str,
            "encoding": original_encoding,
            "is_utf8_converted": is_converted
        }

        # 如果是转换过的文件，标记正确路径
        if is_converted:
            file_info["use_path"] = f"converted/{converted_name}"

        file_info_str = json.dumps(file_info, indent=4, ensure_ascii=False)
        all_file_info_str += f"File {idx}:\n{file_info_str}\n\n"

    # 添加编码状态汇总
    if encoding_map:
        all_utf8 = all(
            item.get("original_encoding") in ("utf-8", "binary")
            for item in encoding_map.values()
        )
        converted_count = sum(1 for item in encoding_map.values() if item.get("is_converted"))

        if all_utf8:
            all_file_info_str += "【编码状态】所有文件已是 UTF-8 编码，可直接使用。\n\n"
        else:
            all_file_info_str += f"【编码状态】已转换 {converted_count} 个非 UTF-8 文件为 UTF-8 编码。\n"
            all_file_info_str += "【重要】读取数据文件时请使用 converted/ 目录下的 UTF-8 版本文件。\n\n"

    return all_file_info_str


def get_file_icon(extension):
    """获取文件图标"""
    ext = extension.lower()
    icons = {
        (".jpg", ".jpeg", ".png", ".gif", ".bmp"): "🖼️",
        (".pdf",): "📕",
        (".doc", ".docx"): "📘",
        (".txt",): "📄",
        (".md",): "📝",
        (".csv", ".xlsx"): "📊",
        (".json", ".sqlite"): "🗄️",
        (".mp4", ".avi", ".mov"): "🎥",
        (".mp3", ".wav"): "🎵",
        (".zip", ".rar", ".tar"): "🗜️",
    }

    for extensions, icon in icons.items():
        if ext in extensions:
            return icon
    return "📁"


def uniquify_path(target: Path) -> Path:
    """若目标已存在，生成 'name (1).ext'、'name (2).ext' 形式的新路径。"""
    if not target.exists():
        return target
    parent = target.parent
    stem = target.stem
    suffix = target.suffix
    import re as _re

    m = _re.match(r"^(.*) \((\d+)\)$", stem)
    base = stem
    start = 1
    if m:
        base = m.group(1)
        try:
            start = int(m.group(2)) + 1
        except Exception:
            start = 1
    i = start
    while True:
        candidate = parent / f"{base} ({i}){suffix}"
        if not candidate.exists():
            return candidate
        i += 1





# API Routes
@app.post("/api/auth/register")
async def register(username: str = Form(...), password: str = Form("")):
    print(f"Registering user: {username}")
    # Built-in superuser registration compatibility: allow short password only for superuser.
    if not is_builtin_superuser(username) and len(password) < 8:
        raise HTTPException(status_code=400, detail="Password must be at least 8 characters long")

    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute("INSERT INTO users (username, password_hash) VALUES (?, ?)",
                       (username, hash_password(password)))
        conn.commit()
        conn.close()
    except sqlite3.IntegrityError:
        raise HTTPException(status_code=400, detail="Username already exists")
    except Exception as e:
        print(f"Registration error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
    return {"message": "Registered successfully"}

@app.post("/api/auth/login")
async def login(username: str = Form(...), password: str = Form("")):
    print(f"Login attempt: {username}")
    if is_builtin_superuser(username):
        if not verify_builtin_superuser_password(password):
            raise HTTPException(status_code=401, detail="Invalid username or password")

        # Ensure superuser exists in the DB for foreign key constraints
        try:
            conn = sqlite3.connect(DB_PATH)
            conn.row_factory = sqlite3.Row
            cursor = conn.cursor()
            ensure_builtin_superuser(cursor)
            conn.commit()
            conn.close()
            return {"username": BUILTIN_SUPERUSER_USERNAME, "is_superuser": True}
        except Exception as e:
            print(f"Superuser login error: {e}")
            traceback.print_exc()
            raise HTTPException(status_code=500, detail=str(e))

    # Any user exists in DB - allow login (for this specific internal tool environment)
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute("SELECT username, password_hash FROM users WHERE username = ?", (username,))
        row = cursor.fetchone()
        conn.close()

        if not row:
            # For this internal environment, allow auto-registration on first login
            try:
                conn = sqlite3.connect(DB_PATH)
                cursor = conn.cursor()
                cursor.execute("INSERT INTO users (username, password_hash) VALUES (?, ?)",
                               (username, hash_password(password)))
                conn.commit()
                conn.close()
                return {"username": username, "is_superuser": False}
            except Exception as reg_err:
                print(f"Auto-registration error: {reg_err}")
                raise HTTPException(status_code=401, detail="Invalid username or password")

        # Check password if one exists in DB (not empty hash and not empty string)
        stored_hash = row["password_hash"]
        empty_hash = hash_password("")

        if stored_hash and stored_hash != "" and stored_hash != empty_hash:
            if hash_password(password) != stored_hash:
                raise HTTPException(status_code=401, detail="Invalid username or password")

        return {"username": username, "is_superuser": False}
    except HTTPException:
        raise
    except Exception as e:
        print(f"Login error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

# --- JWT Auth configuration ---
JWT_SECRET = os.getenv("DEEPANALYZE_JWT_SECRET", secrets.token_hex(32))
JWT_ALGORITHM = "HS256"
JWT_EXPIRATION_HOURS = int(os.getenv("DEEPANALYZE_JWT_EXPIRATION_HOURS", "168"))


def create_jwt_token(username: str) -> str:
    now = int(time_module.time())
    payload = {"sub": username, "iat": now, "exp": now + JWT_EXPIRATION_HOURS * 3600}
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)


def decode_jwt_token(token: str):
    try:
        return jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
    except Exception:
        return None


def get_current_user_from_token(request: Request) -> str | None:
    auth_header = request.headers.get("Authorization", "")
    if auth_header.startswith("Bearer "):
        payload = decode_jwt_token(auth_header[7:])
        if payload:
            return payload.get("sub")
    api_key = request.headers.get("X-API-Key", "")
    if api_key:
        payload = decode_jwt_token(api_key)
        if payload:
            return payload.get("sub")
    return None


@app.post("/v1/auth/register")
async def v1_register(request: Request):
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON body")
    username = body.get("username", "").strip()
    password = body.get("password", "")
    if not username or len(username) < 3:
        raise HTTPException(status_code=400, detail="Username must be at least 3 characters")
    if len(password) < 8:
        raise HTTPException(status_code=400, detail="Password must be at least 8 characters")

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute("SELECT username FROM users WHERE username = ?", (username,))
    if cursor.fetchone():
        conn.close()
        raise HTTPException(status_code=409, detail="Username already exists")

    cursor.execute("INSERT INTO users (username, password_hash) VALUES (?, ?)",
                   (username, hash_password(password)))
    conn.commit()
    conn.close()

    token = create_jwt_token(username)
    return {"access_token": token, "token_type": "bearer", "username": username}


@app.post("/v1/auth/login")
async def v1_login(request: Request):
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON body")
    username = body.get("username", "").strip()
    password = body.get("password", "")

    if is_builtin_superuser(username):
        if not verify_builtin_superuser_password(password):
            raise HTTPException(status_code=401, detail="Invalid username or password")

        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        ensure_builtin_superuser(cursor)
        conn.commit()
        conn.close()

        token = create_jwt_token(BUILTIN_SUPERUSER_USERNAME)
        return {
            "access_token": token,
            "token_type": "bearer",
            "username": BUILTIN_SUPERUSER_USERNAME,
            "is_superuser": True,
        }

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute("SELECT username, password_hash FROM users WHERE username = ?", (username,))
    row = cursor.fetchone()
    conn.close()

    if not row:
        raise HTTPException(status_code=401, detail="Invalid username or password")

    stored_hash = row["password_hash"]
    if stored_hash and stored_hash != hash_password(""):
        if hash_password(password) != stored_hash:
            raise HTTPException(status_code=401, detail="Invalid username or password")

    token = create_jwt_token(username)
    return {
        "access_token": token,
        "token_type": "bearer",
        "username": username,
        "is_superuser": is_builtin_superuser(username),
    }


@app.get("/v1/auth/users")
async def v1_list_users():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    ensure_builtin_superuser(cursor)
    conn.commit()
    cursor.execute("SELECT username FROM users ORDER BY username")
    users = [row["username"] for row in cursor.fetchall()]
    conn.close()
    return {"users": users}


@app.get("/v1/auth/me")
async def v1_get_me(request: Request):
    username = get_current_user_from_token(request)
    if not username:
        raise HTTPException(status_code=401, detail="Authentication required")
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute("SELECT username, password_hash FROM users WHERE username = ?", (username,))
    row = cursor.fetchone()
    conn.close()
    if not row:
        raise HTTPException(status_code=401, detail="User not found")
    return {
        "username": username,
        "created_at": 0,
        "is_superuser": is_builtin_superuser(username),
    }


@app.get("/workspace/files")
async def get_workspace_files(session_id: str = Query("default"), username: str = Query("default")):
    """获取工作区文件列表（支持 user & session 隔离）"""
    workspace_dir = get_session_workspace(session_id, username)
    generated_dir = Path(workspace_dir) / "generated"
    # 获取 generated 目录下的文件名集合
    generated_files = (
        set(f.name for f in generated_dir.iterdir() if f.is_file())
        if generated_dir.exists()
        else set()
    )

    files = []
    for file_path in Path(workspace_dir).iterdir():
        if file_path.is_file():
            if file_path.name in generated_files:
                continue
            stat = file_path.stat()
            # 重新构建包含 username/session_id 的路径
            rel_path = f"{username}/{session_id}/{file_path.name}"
            files.append(
                {
                    "name": file_path.name,
                    "size": stat.st_size,
                    "extension": file_path.suffix.lower(),
                    "icon": get_file_icon(file_path.suffix),
                    "download_url": build_download_url(rel_path),
                    "preview_url": (
                        build_download_url(rel_path)
                        if file_path.suffix.lower()
                        in [
                            ".jpg",
                            ".jpeg",
                            ".png",
                            ".gif",
                            ".bmp",
                            ".pdf",
                            ".txt",
                            ".doc",
                            ".docx",
                            ".csv",
                            ".xlsx",
                        ]
                        else None
                    ),
                }
            )
    return {"files": files}


# ---------- Workspace Tree & Single File Delete ----------
def _rel_path(path: Path, root: Path) -> str:
    try:
        rel = path.relative_to(root)
        return rel.as_posix()
    except Exception:
        return path.name


def build_tree(path: Path, root: Optional[Path] = None) -> dict:
    if root is None:
        root = path
    node: dict = {
        "name": path.name or "workspace",
        "path": _rel_path(path, root),
        "is_dir": path.is_dir(),
    }
    if path.is_dir():
        children = []

        # 自定义排序：generated 和 converted 文件夹放在最后，其他按目录优先、名称排序
        def sort_key(p):
            is_generated = p.name == "generated"
            is_converted = p.name == "converted"
            is_dir = p.is_dir()
            return (is_generated or is_converted, not is_dir, p.name.lower())

        for child in sorted(path.iterdir(), key=sort_key):
            if child.name.startswith("."):
                continue
            children.append(build_tree(child, root))
        node["children"] = children
    else:
        node["size"] = path.stat().st_size
        node["extension"] = path.suffix.lower()
        node["icon"] = get_file_icon(path.suffix)
        rel = _rel_path(path, root)
        node["download_url"] = build_download_url(rel)
        # 标记 converted 目录下的文件
        node["is_converted"] = "converted" + os.sep in rel or rel.startswith("converted")
    return node


@app.get("/workspace/tree")
async def workspace_tree(session_id: str = Query("default"), username: str = Query("default")):
    workspace_dir = get_session_workspace(session_id, username)
    root = Path(workspace_dir)
    tree_data = build_tree(root, root)

    # 在下载链接前加上 username/session_id 前缀
    def prefix_urls(node, un, sid):
        if "download_url" in node and node["download_url"]:
            # 重新构建包含 username/session_id 的路径
            rel = node.get("path", "")
            node["download_url"] = build_download_url(f"{un}/{sid}/{rel}")
        if "children" in node:
            for child in node["children"]:
                prefix_urls(child, un, sid)

    prefix_urls(tree_data, username, session_id)
    return tree_data


@app.delete("/workspace/file")
async def delete_workspace_file(
    path: str = Query(..., description="relative path under workspace"),
    session_id: str = Query("default"),
    username: str = Query("default"),
):
    workspace_dir = get_session_workspace(session_id, username)
    abs_workspace = Path(workspace_dir).resolve()
    target = (abs_workspace / path).resolve()
    if abs_workspace not in target.parents and target != abs_workspace:
        raise HTTPException(status_code=400, detail="Invalid path")
    if not target.exists():
        raise HTTPException(status_code=404, detail="Not found")
    if target.is_dir():
        raise HTTPException(status_code=400, detail="Folder deletion not allowed")
    try:
        target.unlink()
        return {"message": "deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/workspace/move")
async def move_path(
    src: str = Query(..., description="relative source path under workspace"),
    dst_dir: str = Query("", description="relative target directory under workspace"),
    session_id: str = Query("default"),
    username: str = Query("default"),
):
    """在同一 workspace 内移动（或重命名）文件/目录。
    - src: 源相对路径（必填）
    - dst_dir: 目标目录（相对路径，空表示移动到根目录）
    """
    workspace_dir = get_session_workspace(session_id, username)
    abs_workspace = Path(workspace_dir).resolve()

    abs_src = (abs_workspace / src).resolve()
    if abs_workspace not in abs_src.parents and abs_src != abs_workspace:
        raise HTTPException(status_code=400, detail="Invalid src path")
    if not abs_src.exists():
        raise HTTPException(status_code=404, detail="Source not found")

    abs_dst_dir = (abs_workspace / (dst_dir or "")).resolve()
    if abs_workspace not in abs_dst_dir.parents and abs_dst_dir != abs_workspace:
        raise HTTPException(status_code=400, detail="Invalid dst_dir path")
    abs_dst_dir.mkdir(parents=True, exist_ok=True)

    target = abs_dst_dir / abs_src.name
    target = uniquify_path(target)
    try:
        shutil.move(str(abs_src), str(target))
        rel_new = str(target.relative_to(abs_workspace))
        return {"message": "moved", "new_path": rel_new}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Move failed: {e}")


@app.delete("/workspace/dir")
async def delete_workspace_dir(
    path: str = Query(..., description="relative directory under workspace"),
    recursive: bool = Query(True, description="delete directory recursively"),
    session_id: str = Query("default"),
    username: str = Query("default"),
):
    """删除 workspace 下的目录。默认递归删除，禁止删除根目录。"""
    workspace_dir = get_session_workspace(session_id, username)
    abs_workspace = Path(workspace_dir).resolve()
    target = (abs_workspace / path).resolve()
    if abs_workspace not in target.parents and target != abs_workspace:
        raise HTTPException(status_code=400, detail="Invalid path")
    if target == abs_workspace:
        raise HTTPException(status_code=400, detail="Cannot delete workspace root")
    if not target.exists():
        raise HTTPException(status_code=404, detail="Not found")
    if not target.is_dir():
        raise HTTPException(status_code=400, detail="Not a directory")
    try:
        if recursive:
            shutil.rmtree(target)
        else:
            target.rmdir()
        return {"message": "deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/proxy")
async def proxy(url: str):
    """Simple CORS proxy for previewing external files.
    WARNING: For production, add domain allowlist and authentication.
    """
    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=15) as client_httpx:
            r = await client_httpx.get(url)
        return Response(
            content=r.content,
            media_type=r.headers.get("content-type", "application/octet-stream"),
            headers={"Access-Control-Allow-Origin": "*"},
            status_code=r.status_code,
        )
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Proxy fetch failed: {e}")


def convert_to_utf8(file_path: Path, workspace_dir: str) -> tuple[Optional[Path], str]:
    """检查文件编码并转换为 UTF-8，另存为 converted/ 子目录下的文件。
    返回 (目标路径, 原始编码) 元组。
    """
    if not file_path.exists() or not file_path.is_file():
        return None, "unknown"

    # 已经是 converted 目录下的文件，跳过
    if "converted" + os.sep in str(file_path) or str(file_path).startswith("converted"):
        return file_path, "utf-8"

    # 仅转换文本类文件
    if file_path.suffix.lower() not in [".csv", ".txt", ".md", ".json", ".xml"]:
        return file_path, "binary"

    try:
        with open(file_path, "rb") as f:
            raw_data = f.read()

        result = chardet.detect(raw_data)
        encoding = result['encoding']

        if not encoding:
            encoding = 'utf-8' # 兜底

        # converted 子目录
        converted_dir = Path(workspace_dir) / "converted"
        converted_dir.mkdir(parents=True, exist_ok=True)

        # 目标路径：converted/原始文件名（不加 _utf8 后缀）
        utf8_path = converted_dir / file_path.name
        # 唯一化，避免覆盖
        utf8_path = uniquify_path(utf8_path)

        if encoding.lower() == 'utf-8':
            # 已经是 utf-8，直接复制到 converted 目录
            if not utf8_path.exists():
                shutil.copy2(file_path, utf8_path)
            return utf8_path, "utf-8"

        # 转换非 UTF-8 编码
        content = raw_data.decode(encoding, errors='replace')
        with open(utf8_path, "w", encoding="utf-8") as f:
            f.write(content)
        return utf8_path, encoding
    except Exception as e:
        print(f"Error converting {file_path} to UTF-8: {e}")
        return file_path, "unknown"

@app.post("/workspace/upload")
async def upload_files(
    files: List[UploadFile] = File(...), session_id: str = Query("default"), username: str = Query("default")
):
    """上传文件到工作区（支持 user & session 隔离），并自动转换为 UTF-8"""
    workspace_dir = get_session_workspace(session_id, username)
    uploaded_files = []
    encoding_report = []  # 记录每个文件的编码状态

    for file in files:
        # 唯一化文件名，避免覆盖
        dst = uniquify_path(Path(workspace_dir) / file.filename)
        with open(dst, "wb") as buffer:
            content = await file.read()
            buffer.write(content)

        # 自动转换为 UTF-8（保存到 converted/ 子目录）
        utf8_dst, original_encoding = convert_to_utf8(dst, workspace_dir)

        # 记录编码信息
        encoding_report.append({
            "original_name": dst.name,
            "converted_name": utf8_dst.name if utf8_dst else dst.name,
            "original_encoding": original_encoding,
            "is_converted": (utf8_dst != dst) if utf8_dst else False,
            "path": str(utf8_dst.relative_to(Path(workspace_dir))) if utf8_dst else str(dst.relative_to(Path(workspace_dir)))
        })

        uploaded_files.append(
            {
                "name": dst.name,
                "size": len(content),
                "path": str(dst.relative_to(Path(workspace_dir))),
            }
        )
        if utf8_dst and utf8_dst != dst:
             uploaded_files.append(
                {
                    "name": utf8_dst.name,
                    "size": utf8_dst.stat().st_size,
                    "path": str(utf8_dst.relative_to(Path(workspace_dir))),
                }
            )

    # 生成编码状态汇总
    all_utf8 = all(item["original_encoding"] == "utf-8" or item["original_encoding"] == "binary" for item in encoding_report)
    converted_count = sum(1 for item in encoding_report if item["is_converted"])

    encoding_summary = {
        "all_files_utf8": all_utf8,
        "converted_count": converted_count,
        "total_files": len(encoding_report),
        "files": encoding_report
    }

    # 保存编码映射文件，供智能体后续分析时读取
    encoding_map_path = Path(workspace_dir) / ".encoding_map.json"
    with open(encoding_map_path, "w", encoding="utf-8") as f:
        json.dump(encoding_summary, f, indent=2, ensure_ascii=False)

    return {
        "message": f"Successfully uploaded {len(uploaded_files)} files",
        "files": uploaded_files,
        "encoding_info": encoding_summary
    }


@app.delete("/workspace/clear")
async def clear_workspace(session_id: str = Query("default"), username: str = Query("default")):
    """清空工作区（支持 user & session 隔离）"""
    workspace_dir = get_session_workspace(session_id, username)
    if os.path.exists(workspace_dir):
        shutil.rmtree(workspace_dir)
    os.makedirs(workspace_dir, exist_ok=True)
    return {"message": "Workspace cleared successfully"}


@app.post("/workspace/upload-to")
async def upload_to_dir(
    dir: str = Query("", description="relative directory under workspace"),
    files: List[UploadFile] = File(...),
    session_id: str = Query("default"),
    username: str = Query("default"),
):
    """上传文件到 workspace 下的指定子目录（仅限工作区内），并自动转换为 UTF-8。"""
    workspace_dir = get_session_workspace(session_id, username)
    abs_workspace = Path(workspace_dir).resolve()
    target_dir = (abs_workspace / dir).resolve()
    if abs_workspace not in target_dir.parents and target_dir != abs_workspace:
        raise HTTPException(status_code=400, detail="Invalid dir path")
    target_dir.mkdir(parents=True, exist_ok=True)

    saved = []
    encoding_report = []
    for f in files:
        dst = uniquify_path(target_dir / f.filename)
        try:
            with open(dst, "wb") as buffer:
                content = await f.read()
                buffer.write(content)

            # 自动转换为 UTF-8（保存到 converted/ 子目录）
            utf8_dst, original_encoding = convert_to_utf8(dst, workspace_dir)

            encoding_report.append({
                "original_name": dst.name,
                "converted_name": utf8_dst.name if utf8_dst else dst.name,
                "original_encoding": original_encoding,
                "is_converted": (utf8_dst != dst) if utf8_dst else False,
            })

            saved.append(
                {
                    "name": dst.name,
                    "size": len(content),
                    "path": str(dst.relative_to(abs_workspace)),
                }
            )
            if utf8_dst and utf8_dst != dst:
                 saved.append(
                    {
                        "name": utf8_dst.name,
                        "size": utf8_dst.stat().st_size,
                        "path": str(utf8_dst.relative_to(abs_workspace)),
                    }
                )
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Save failed: {e}")

    all_utf8 = all(item["original_encoding"] == "utf-8" or item["original_encoding"] == "binary" for item in encoding_report)
    converted_count = sum(1 for item in encoding_report if item["is_converted"])

    encoding_summary = {
        "all_files_utf8": all_utf8,
        "converted_count": converted_count,
        "total_files": len(encoding_report),
        "files": encoding_report
    }

    # 保存编码映射文件，供智能体后续分析时读取
    encoding_map_path = Path(workspace_dir) / ".encoding_map.json"
    with open(encoding_map_path, "w", encoding="utf-8") as f:
        json.dump(encoding_summary, f, indent=2, ensure_ascii=False)

    return {"message": f"uploaded {len(saved)} files", "files": saved, "encoding_info": encoding_summary}


@app.post("/execute")
async def execute_code_api(request: dict):
    """执行 Python 代码"""
    try:
        code = request.get("code", "")
        session_id = request.get("session_id", "default")
        username = request.get("username", "default")
        workspace_dir = get_session_workspace(session_id, username)

        if not code:
            raise HTTPException(status_code=400, detail="No code provided")

        # 使用子进程安全执行，避免 GUI/线程问题（在指定 session workspace 中）
        result = await run_in_threadpool(execute_code_safe, code, workspace_dir)

        return {
            "success": True,
            "result": result,
            "message": "Code executed successfully",
        }

    except Exception as e:
        return {
            "success": False,
            "result": f"Error: {str(e)}",
            "message": "Code execution failed",
        }


def fix_code_block(content):
    def fix_text(text):
        stack = []
        lines = text.splitlines(keepends=True)
        result = []
        for line in lines:
            stripped = line.strip()
            if stripped.startswith("```python"):
                if stack and stack[-1] == "```python":
                    result.append("```\n")
                    stack.pop()
                stack.append("```python")
                result.append(line)
            elif stripped == "```":
                if stack and stack[-1] == "```python":
                    stack.pop()
                result.append(line)
            else:
                result.append(line)
        while stack:
            result.append("```\n")
            stack.pop()
        return "".join(result)

    if isinstance(content, str):
        return fix_text(content)
    elif isinstance(content, tuple):
        text_part = content[0] if content[0] else ""
        return (fix_text(text_part), content[1])
    return content


def fix_tags_and_codeblock(s: str) -> str:
    """
    修复未闭合的tags，并确保</Code>后代码块闭合。
    """
    pattern = re.compile(
        r"<(Analyze|Understand|Code|Execute|Answer|TaskTree)>(.*?)(?:</\1>|(?=$))", re.DOTALL
    )

    # 找所有匹配
    matches = list(pattern.finditer(s))
    if not matches:
        return s  # 没有标签，直接返回

    # 检查最后一个匹配是否闭合
    last_match = matches[-1]
    tag_name = last_match.group(1)
    matched_text = last_match.group(0)

    if not matched_text.endswith(f"</{tag_name}>"):
        # 没有闭合，补上
        if tag_name == "Code":
            s = fix_code_block(s) + f"\n```\n</{tag_name}>"
        else:
            s += f"\n</{tag_name}>"

    return s


def get_system_prompt_with_fonts() -> str:
    """
    生成包含动态字体路径的 system prompt
    直接返回原始 system_prompt，使用字符串替换
    """
    # 获取字体目录
    try:
        fonts_dir = str(get_fonts_dir())
    except (NameError, ImportError):
        fonts_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "../../assets/fonts")

    # 原始 system prompt 模板（完整内容）
    system_prompt_template = """你是 DeepAnalyze，一位精通 Python 和 R 语言的顶尖数据科学家，同时也是专门从事中国海关风险管理和风险防控的数据分析专家。

**你的核心使命**：
忠于国家安全，服务海关履行职责，以风险管理和防控违法违规为目标，确保海关监管措施、管理规定、政策制度执行到位，维护市场公平竞争的秩序环境。

**你的核心任务**：
基于数据统计、比较、相关性和逻辑推理，深入分析进出口业务主体（包括经营企业、收发货人、货主单位、报关单位、代理单位、运输企业、跨境电商平台及其参与方等）的行为数据，挖掘并报告以下违法违规行为：
1. 进出口过程中的走私及违规行为；
2. 违反安全准入管理规定的行为；
3. 通过伪报、瞒报、虚报等方式逃避监管证件管理的行为；
4. 通过低报价格、伪报原产地、伪报HS编码归类逃避税税的行为。你的分析结果应明确指出可疑行为，并详细阐述推理原因。

**数据客观性准则（极重要）**：
你必须始终坚持以数据的客观性和用户指示要求为根本准则，严格遵循以下原则：
1. **忠于数据原则**：分析过程中只能使用数据中实际存在的内容加上合理的推测，而不能凭空引入各种无关的概念强加于数据之上，不得以幻想的方式赋予意义。如果数据中不包含某字段（如"通关""查验""放行"等），则分析过程中绝不能出现这些概念或基于这些概念的推理。
2. **数据探测优先（强制首步）**：分析开始后，必须**立即**执行数据探测代码，全面掌握数据形状（行数、列数）、字段名称、字段类型、字段定义、字段内容填制规范、字段内容代码字典（如有）等信息。具体步骤：
   - 执行 `df.shape`、`df.columns.tolist()`、`df.head(10)`、`df.dtypes`、`df.describe(include='all')`，获取完整的数据概貌。
   - 对每个字段查看唯一值样本（`df['col'].unique()[:20]`），理解字段内容的实际含义和编码方式。
   - **关键：将数据字段名称、字段定义、字段内容与用户分析目标关联起来**，初步推测在分析任务中可能用到的数据内容。
   - **用户表达与字段名称映射（极重要）**：用户的自然语言表达与数据字段名称之间经常存在微小差异，你必须智能匹配：例如用户说"运输工具"，但数据字段可能是"车牌号""车辆""运输方式"；用户说"国家"，但数据字段可能是"国别""原产国""启运国"；用户说"企业"，但数据字段可能是"经营单位""收发货人""申报单位"等。你必须基于已探测到的数据字段和字段内容，正确理解用户的需求，将用户语义准确映射到对应的数据字段，而不是因为字面不匹配就报告字段不存在。
   - 将字段映射关系（用户表达 → 实际字段名）明确记录在数据字典中，供后续所有分析任务引用。
3. **合理发散但不偏离数据**：保持分析角度的合理发散广度，但所有发散方向必须能在给定数据中找到支撑。不得凭空引入数据中不存在的业务概念。
4. **用户指示优先**：当用户提出具体分析需求时，首先在已建立的数据字典中查找语义匹配的字段（而不仅仅是字面匹配）。若确实不存在语义相关的字段，应向用户明确说明数据中缺少相关字段，并建议可行的替代分析方向。

**报告生成规范**：
- **阶段性分析**：在分析过程中，请通过聊天回复或代码执行结果展示阶段性发现和图表。每个阶段性分析都应按"文字解读在前 → '如下图：' → 图表在后"的顺序组织。
- **最终报告**：只有在**所有**分析维度（如：主体身份、价格风险、通关时效等）全部完成后，在任务的**最后阶段**，才调用 `generate_report_pdf` 或 `generate_report_docx` 将所有核心观点和可视化图形组织成一份完整的终期报告。
- **报告结构要求**：最终报告必须使用树形层次结构组织（一、（一）1. 2.（二）1. 2. 二、……），每个分析点按"文字解读 → 如下图 → 图表"的顺序排列。
- **避免重复**：禁止在分析过程中反复生成阶段性报告，确保用户在文件列表中只看到最终的、高质量的分析成果。
- **上下文积累（极重要）**：在整个分析过程中，你必须在内存中持续积累所有已完成的分析成果、关键发现、数据统计结果和图表文件路径。每完成一个分析维度后，必须在 <Answer> 标签中完整记录该维度的**全部详细内容**（包括具体数据数值、百分比、排名、图表引用路径等），而非仅写标题或占位符。
- **禁止空壳报告（极重要）**：生成最终报告时，报告内容**必须包含完整的分析结论、具体数据、推理过程和图表引用**。严禁生成仅包含报告结构框架而缺少实质内容的报告（如"[详细方法描述...]"、"[详细结果描述...]"这类占位文本）。如果某个章节的分析尚未完成，应明确标注"此部分尚未完成分析"而非使用占位符。
- **报告不得包含水印、版权声明或数据安全说明**：生成的报告中不需要添加任何水印、版权声明、免责声明或数据安全说明文字。报告不需要以压缩包方式打包输出。

**============================================
第一部分：可视化与美学规范（极重要）
============================================
你应以生动活泼且专业的方式展示分析结果。
1. **可视化风格**：
   - **Python**: 务必使用 `seaborn` 库提升图表美感。
     ```python
     import seaborn as sns
     import matplotlib.pyplot as plt
     sns.set_theme(style="whitegrid", palette="muted")
     # 使用专业颜色和清晰的标签
     ```
   - **R**: 务必使用 `ggplot2` 并配合 `theme_minimal()` 或 `theme_light()`。
   - **要求**：图表应清晰展示**趋势、关系、分布、规律**。禁止使用极其简陋的默认绘图样式。

2. **图表内容**：
   - 每个图表必须有清晰的中文标题、坐标轴标签和图例。
   - 对于趋势分析，使用折线图；对于分布分析，使用直方图或小提琴图；对于关系分析，使用散点图或热力图。

3. **空白图表检测（极重要）**：
   - 生成图表后，**必须**验证图表不是空白的。在 `plt.savefig()` 之前，检查数据是否有效。
   - 对于时间序列图表，必须先确认时间列已正确解析为 datetime 类型，使用 `pd.to_datetime(errors='coerce')` 进行转换，并检查转换后是否存在 NaT 值。
   - 验证方法：在保存图表前添加检查代码：
     ```python
     # 验证图表数据是否有效
     fig = plt.gcf()
     has_data = False
     for ax in fig.axes:
         for line in ax.get_lines():
             if len(line.get_xdata()) > 0:
                 has_data = True
                 break
         for container in ax.containers:
             if len(container) > 0:
                 has_data = True
                 break
         for collection in ax.collections:
             if len(collection.get_offsets()) > 0:
                 has_data = True
                 break
     if not has_data:
         print("警告：图表为空白，跳过保存")
         plt.close()
     else:
         plt.savefig(output_path, dpi=150, bbox_inches='tight')
         plt.close()
         print(f"图表已保存: {output_path}")
     ```
   - 如果图表为空白，不要将其包含在最终报告中。

**============================================
第二部分：工具库与增强组件（极重要）
============================================
为确保分析的高效与稳定，系统为你封装了专用工具库 `agent_utils`。在编写代码时，**务必优先导入并使用**以下功能：

1. **PDF/DOCX/PPTX 报告生成（一键式，支持中文）**：
   ```python
   from agent_utils import generate_report_pdf, generate_report_docx, generate_report_pptx

   # 直接从 Markdown 文本生成 PDF、DOCX 和 PPTX
   generate_report_pdf(report_md, "output.pdf", title="分析报告")
   generate_report_docx(report_md, "output.docx", title="分析报告")
   generate_report_pptx(report_md, "output.pptx", title="分析报告")
   ```

2. **FPDF 2.x 中文支持（解决 Undefined font 错误）**：
   如果你需要更细粒度的 FPDF 控制，请使用 `init_fpdf_chinese()`，它会自动注册所有必要的字体（SimHei, STFangSong 等）。
   ```python
   from agent_utils import init_fpdf_chinese
   pdf = init_fpdf_chinese() # 返回已预注册中文字体的 FPDF 实例
   pdf.add_page()
   pdf.set_font("SimHei", "B", 16)
   pdf.cell(0, 10, "中文标题", ln=True)
   ```

3. **时间处理（解决 Pandas 2.0+ Timestamp 减法错误）**：
   禁止直接使用 `pd.Timestamp - 30`。请使用 `pandas_date_sub` 或 `pd.Timedelta`。
   ```python
   from agent_utils import pandas_date_sub, pandas_date_add

   # 正确示例
   start_date = pandas_date_sub(end_date, 30) # 减去30天
   ```

**============================================
第三部分：PDF库弃用警告与正确用法（极重要）
============================================
使用 FPDF/fpdf2 库时，**必须**遵守以下规则以避免已知的弃用警告和错误：

1. **FPDF 2.x 弃用警告（"DeprecationWarning" / "FPDF internal: ..."）**：
   - 原因：FPDF 2.x 弃用了 `add_font(dejaVu condensed=True)` 语法。应使用 `style='B'` 参数。
   - **完整正确写法示例（生成中文PDF）**：
     ```python
     from fpdf import FPDF
     pdf = FPDF()
     # 添加中文字体（动态解析字体路径，推荐方式）
     fonts_dir = "{FONTS_DIR_PLACEHOLDER}"

     pdf.add_font('SimHei', '', os.path.join(fonts_dir, 'simhei.ttf'))
     pdf.add_font('SimHei', 'B', os.path.join(fonts_dir, 'simhei.ttf'))
     pdf.add_font('SimKai', '', os.path.join(fonts_dir, 'simkai.ttf'))
     pdf.add_font('STFangSong', '', os.path.join(fonts_dir, 'STFangSong.ttf'))

     pdf.add_page()
     # 设置字体，大小
     pdf.set_font('SimHei', '', 16)
     pdf.cell(0, 10, '中文标题', ln=True, align='C')
     pdf.set_font('STFangSong', '', 12)
     pdf.multi_cell(0, 10, '这是一段中文正文，可以正常显示。请注意使用 multi_cell 来自动换行。')
     pdf.output('output.pdf')
     ```
   - **禁止写法**：`add_font('SimHei', condensed=True)` → 会产生弃用警告
   - **禁止写法**：`add_font('SimHei', '', '/path/font.ttf', condensed=True)` → 会产生弃用警告
   - **重要**：必须使用 `pdf.set_font()` 设置已注册的字体，并使用 `pdf.multi_cell()` 渲染中文文本。

2. **FPDF 2.x 图片嵌入弃用（"NOTSUBSET" / "ftface warning"）**：
   - 原因：旧版 FPDF 的 `image()` 方法在处理某些 TTF 字体时会产生 NOTSUBSET 警告。
   - 解决方案：使用 `fpdf2`（而非旧版 `fpdf`），并确保字体在图片渲染前已正确注册。
   - 如使用 matplotlib 生成图片再嵌入 PDF，建议在 matplotlib 中直接使用支持 CJK 的字体渲染文字，避免依赖 FPDF 的字体子集化。

3. **reportlab 注意事项**：
   - reportlab 的 `TTFont` 不支持 TTC 格式字体（如 macOS 的 STHeiti、PingFang SC 等）。
   - 必须使用 assets/fonts/ 目录下的纯 TTF 字体文件。

4. **PDF 渲染中文多行文本**：
   - 使用 `pdf.multi_cell()` 而非手动换行。
   - 设置正确的 `ln` 参数控制换行后的光标位置。

**============================================
第三部分（续）：时间处理（极重要）
============================================
系统已安装丰富的时间处理包，包括但不限于：`pandas`, `numpy`, `datetime`, `dateutil`, `pytz`, `zoneinfo`, ` pendulum`, `delorean`, `timeuuid`。

**时间粒度与上下文敏感性原则**：
- 当用户要求"按月分析"时，**必须**先判断数据中时间字段的覆盖范围。
- **如果数据跨越多个年份**，则"按月分析"实际上是指"每年每月"（例如：2023年1月 vs 2024年1月是完全不同的时间段）。
- 绝不能将不同年份的同一月份混为一谈。
- 正确的做法：先探索数据的时间范围（使用 `df['date_column'].min()` / `.max()`），若跨年，按 "YYYY-MM" 格式分组，或按 "年份-月份" 组合维度分析。

**时间分析的输出结构**：
- 在分析报告中，必须明确标注每个时间段的年份。
- 例如：`2023年1月`、`2024年3月` 而非笼统的 `1月`、`3月`。

**============================================
第四部分：报告结构规范（极重要）
============================================
每次分析完成后，生成的报告**必须**遵循以下结构，内容必须具有层次感，富含推理和推导过程。**报告内所有内容必须使用简体中文**：

**报告整体结构（PDF/DOCX/PPTX/聊天输出均适用）**：
1. **第一部分：分析思路**（必须放在报告开头）
   - 阐述分析的背景和原因。
   - 描述推理逻辑和将要采取的步骤。

2. **第二部分：分析主体内容**（放在分析思路之后、分析小结之前）
   - **基于陈述原因和推理而得出结论**，而非直接给出结论。
   - 每一章节应包含：
     - **现状描述**（观察到的数据事实）
     - **推理推导**（为什么会出现这种情况，基于数据的逻辑推理）
     - **阶段性结论**（基于前两点得出的初步发现）
     - **可视化图表**（图表标题置于图表下方，引用编号如"图1："）
   - **图表与观点紧密结合（极重要）**：每个分析成果和观点所对应的数据图表必须紧跟在该观点的文字描述之后，而不是集中放在报告末尾。用户应能在阅读每个观点时立即看到支撑该观点的数据图表。
   - **报告层次组织规范**：
     - 每个分析维度作为独立章节
     - 章节内按"观点陈述 → 数据图表 → 推理说明"的顺序组织
     - 图表使用 Markdown 图片引用格式 `![图表描述](图表路径)` 嵌入
     - 确保每个关键发现都有对应的可视化支撑

3. **第三部分：分析小结**（必须放在报告结尾）
   - 总结核心发现。
   - 给出基于风险的针对性建议。
   - 附一个汇总表格，列出所有关键发现和对应的风险等级。

**重要提醒**：
- 严禁"直接得结论"，必须体现"数据 -> 推理 -> 结论"的层次感。
- 报告顺序：**分析思路** → **主体分析内容（推理+数据+图表交织）** → **分析小结**
- 严禁将图表集中堆放在报告某一处，必须将图表分散到各章节的对应观点位置。
- PDF/DOCX/PPTX排版要求：保持标题层级、使用正确的中文字体。

**============================================
第五部分：排版与美学规范
============================================
生成 PDF/DOCX 报告时，请严格遵守以下排版规范：

**字体使用规范**：
- **标题**：使用 SimHei（黑体），字号 16-22pt
- **副标题**：使用 SimHei（黑体），字号 14-16pt
- **正文**：使用 STFangSong（仿宋），字号 10-12pt
- **强调/引用**：使用 SimKai（楷体），字号 11-12pt
- **图表标题**：使用 SimHei（黑体），字号 10-12pt

**间距规范**：
- 段落间距：1.2-1.5 倍行距
- 标题与正文间距：12pt 以上
- 图表上下间距：10pt 以上

**页面布局**：
- 页边距：上下左右各 50pt
- 页码：居中，使用阿拉伯数字
- 章节分页：重要章节前应分页

**图表规范**：
- 图表编号置于标题前，如"图1："
- 图表标题置于图表下方
- 表格标题置于表格上方，编号如"表1："
- 确保中文显示正常，使用支持中文的字体

**表格数据输出规范（极重要）**：
- 在报告中输出表格数据时，**禁止使用 Markdown 表格语法**（即 `| 列1 | 列2 |` 格式），因为 PDF 渲染引擎无法正确对齐 Markdown 表格。
- **正确做法**：使用代码将表格数据导出为格式化的文本或图片。推荐方案：
  1. **使用 matplotlib 渲染表格为图片（首选）**：
     ```python
     import matplotlib.pyplot as plt
     import pandas as pd
     
     fig, ax = plt.subplots(figsize=(10, len(df)*0.4 + 1))
     ax.axis('off')
     table = ax.table(cellText=df.values, colLabels=df.columns, cellLoc='center', loc='center')
     table.auto_set_font_size(False)
     table.set_fontsize(10)
     table.scale(1.2, 1.5)
     plt.savefig('table_output.png', dpi=150, bbox_inches='tight', pad_inches=0.1)
     plt.close()
     ```
  2. **使用 `df.to_string()` 输出格式化文本**（备选）。
- 这确保表格在 PDF、DOCX、PPTX 中都能正确显示且列对齐。

**============================================
第六部分：任务执行与输出规范（极重要）
============================================

**执行流程**：
1. 【数据探测（必须首步）】→ 2. 【预测推理与短代码测试】→ 3. 【正式分析思路】→ 4. 【完整分析代码执行】→ 5. 【结果解读与报告生成】

**数据探测（强制首步，不可跳过）**：
- 在任何分析开始前，**第一个 <Code> 块必须**执行以下数据探测代码：
  ```python
  import pandas as pd
  import os
  
  # 1. 列出工作区文件
  for f in os.listdir('.'):
      print(f)
  
  # 2. 读取数据并探测结构
  df = pd.read_csv('文件名.csv')  # 或 read_excel
  print("形状:", df.shape)
  print("列名:", df.columns.tolist())
  print("数据类型:\\n", df.dtypes)
  print("前5行:\\n", df.head())
  print("缺失值:\\n", df.isnull().sum())
  print("数值列统计:\\n", df.describe())
  ```
- 探测结果必须被保存在变量中，后续所有分析代码必须引用探测到的实际列名。
- **绝不可在探测前就编写引用具体列名的分析代码**。

**代码健壮性规范（极重要）**：
- **所有代码块必须包含完整的 try/except 异常处理**，捕获后打印有意义的错误信息而非让程序崩溃。
- **列名引用必须使用变量**：使用 `cols = df.columns.tolist()` 然后通过列表操作引用列名，而不是硬编码字符串。
- **类型转换必须安全**：使用 `pd.to_numeric(errors='coerce')`、`pd.to_datetime(errors='coerce')` 等安全转换方法。
- **文件路径必须验证**：使用 `os.path.exists()` 验证文件路径后再操作。
- **避免常见错误模式**：
  - ❌ `df['不存在的列']` → 会抛出 KeyError
  - ✅ `if '列名' in df.columns: df['列名']` → 先检查再使用
  - ❌ `pd.Timestamp - 30` → 整数减法不支持
  - ✅ `pd.Timestamp - pd.Timedelta(days=30)` → 正确的时间运算
  - ❌ 在代码中硬编码未验证的列名
  - ✅ 先探测列名，再根据实际列名编写分析代码
  - ❌ 在临时文件中使用硬编码路径或未关闭的文件句柄
  - ✅ 使用 `with` 语句管理文件和使用 `os.path.join()` 构建路径
  - ❌ 忘记 `import` 所需的库
  - ✅ 每个 `<Code>` 块都必须包含完整的 import 语句，不要假设其他代码块已导入
  - ❌ 使用 `plt.show()` 在无显示器环境中
  - ✅ 使用 `plt.savefig()` 保存图表后调用 `plt.close()`

**预测推理与短代码测试（ mandatory ）**：
- 在生成完整的复杂分析代码前，**必须**先进行预测推理和短代码测试。
- 预测推理：思考接下来的代码执行可能会在哪些环节出错（如：路径、编码、数据类型、字段名等）。
- 短代码测试：编写极简代码（2-5行）来验证关键假设，如：
  - 文件路径/目录是否存在且正确。
  - 文件编码是否能被正确读取（尝试 UTF-8 或其他）。
  - 关键字段名是否如预期般存在于 DataFrame 中。
  - 函数调用是否能正确作用于对象。
- 根据测试结果修正后，再编写完整的分析代码。

**每个 <Analyze> 标签内应包含**：
```
# 预测推理
[对潜在执行错误的预判]

# 短代码测试与结果
[执行测试代码并说明结果]

# 正式分析思路
[基于测试结果修正后的正式分析路径]
```

**每个 <Answer> 标签内应包含**：
```
# 分析结论
[基于详尽推导过程得出的结论，体现层次感]

# 风险提示
[潜在风险点识别]

# 建议
[针对性的改进策略]
```

**代码执行规范**：
- 每个 <Code> 标签应完成一个完整的任务
- 执行后应输出关键信息（如图表路径、数据统计）
- 遇到错误应记录原因并尝试解决

**工作小结格式**：
   ```
   【工作小结】
   - 数据处理：描述数据处理过程和关键数据操作
   - 风险点识别：列出发现的主要风险点
   - 执行统计：共计完成X个子任务，生成Y个素材文件

   【运行问题记录】（如实填写，即使没有问题也需标注"无"）
   - 错误信息：[列出执行过程中遇到的各类错误及错误代码]
   - 死循环情况：[记录检测到的死循环次数、原因及跳出操作]
   - 环境组件问题：[组件缺失、版本冲突等问题的描述]
   - 分析重复原因：[如果出现分析结果重复，说明原因]
   - 其他异常：[其他导致运行不畅顺的因素]

   【改进建议】
   - 针对上述问题，提出具体的改进方向
   ```

---

**文件名与路径管理原则**：
- 所有由代码生成的文件，应在生成后立即打印其完整路径。
- 后续代码引用这些文件时，**必须使用之前打印的完整路径**，禁止重新猜测文件名。
- 使用 `os.listdir()` 或 `glob` 验证文件存在性。
- **已生成的图表和数据不得重复生成**，直接引用已有文件路径。

**============================================
第六部分（续）：环境就绪与字体规范（极重要）
============================================
- **环境就绪**：系统已为你安装了充足的 Python 和 R 语言工具包，包括但不限于 `fpdf2`, `python-docx`, `python-pptx`, `pandas`, `matplotlib`, `seaborn`, `chardet`, `reportlab` 等。
- **字体已预注册（极重要）**：系统在启动时已自动完成所有中文字体的注册工作（matplotlib、reportlab、fpdf2 均已预注册）。**在生成报告或绘图时，无需再次手动注册字体**。直接使用 `agent_utils` 中的 `generate_report_pdf`、`generate_report_docx`、`generate_report_pptx` 即可。matplotlib 绑定的字体也已在启动时通过 `font_manager.fontManager.addfont()` 完成注册。
- **R 语言中文/PDF 增强**：在 R 环境中，已为你安装了 `showtext`, `extrafont`, `Cairo`, `grDevices`, `ggplot2`, `lattice`, `knitr`, `rmarkdown`, `tinytex` 等核心包。在生成包含中文的 PDF 或图形时，请务必调用 `showtext_auto()`，并优先使用 `CairoPDF()` 或 `xelatex` 引擎进行渲染，确保中文字符完美显示。
- **UTF-8 编码优先**：系统已自动将上传的文本文件转换为 UTF-8 编码并保存到 `converted/` 子目录（文件名保持不变）。**无论用户输入的文件名是否带有编码转换标注，系统会自动将其映射到 `converted/` 目录下的正确文件进行分析**，请直接根据用户提到的文件名进行数据读取，系统会自动处理文件路径映射。
- **中文字符与编码处理**：在处理任何数据文件前，应确认使用 UTF-8 编码。对于任何包含中文的内容，必须确保在所有输出文件（Png, Jpg, Pdf, Txt, Csv, Docx 等）中正确显示中文。
- **可视化支持**：在 Python 绘图时，务必配置 `plt.rcParams['font.sans-serif']` 使用 `SimHei`, `PingFang SC` 或其他系统中文字体，防止出现乱码或方框。在 R 中使用 `showtext` 处理中文。
- **报告生成**：分析完成后，必须生成详细的最终报告。**最终报告必须同时包含 PDF、DOCX 和 PPTX 格式**，这是你的标准交付物。
  - **PDF 生成推荐方案（按优先级）**：
    1. **agent_utils 一键生成（首选）**：使用 `from agent_utils import generate_report_pdf, generate_report_docx`
    2. **reportlab + 中文字体**：使用 reportlab 库，注册 assets/fonts/ 下的纯 TTF 字体生成 PDF
    3. **fpdf2 + init_fpdf_chinese()**：使用 `from agent_utils import init_fpdf_chinese`
    4. **matplotlib PdfPages（图表为主）**：使用 matplotlib 的 PdfPages 生成包含图表的 PDF
  - **注意**：当前环境为 macOS，禁止使用 `comtypes` 或 `docx2pdf` 库。
- **字体与路径支持（极重要）**：请务必清楚当前环境中有以下字体可用，**不要随意猜测或尝试不存在的字体路径**：
  - **中文字体（assets/fonts/ 目录，纯 TTF，reportlab/matplotlib 全支持，已验证可用）**：
    - **动态解析字体路径（推荐）**：
      ```python
      import os
      from pathlib import Path

      # 方法1: 从项目根目录查找
      project_root = Path.cwd()
      while project_root != project_root.parent and not (project_root / "assets" / "fonts").exists():
          project_root = project_root.parent
      fonts_dir = str(project_root / "assets" / "fonts")

      # 方法2: 相对路径（当工作区已知时）
      fonts_dir = os.path.join(os.path.dirname(__file__), "..", "..", "assets", "fonts")

      # 使用字体
      font_path = os.path.join(fonts_dir, "simhei.ttf")
      ```
    - **可用字体文件**：
      - `SimHei` → `simhei.ttf`（黑体，主标题/重点内容，✅已验证）
      - `SimKai` → `simkai.ttf`（楷体，引用/强调/副标题，✅已验证）
      - `STFangSong` → `STFangSong.ttf`（仿宋，正文/报告，✅已验证）
      - `STHeiti` → `STHeiti.ttf`（黑体备选，✅已验证）
      - `LiSongPro` → `LiSongPro.ttf`（隶书，可选）
  - **macOS 系统 CJK 字体（TTC 格式，用于 R showtext / matplotlib）**：
    - `STHeiti Light` → `/System/Library/Fonts/STHeiti Light.ttc`（黑体）
    - `STHeiti Medium` → `/System/Library/Fonts/STHeiti Medium.ttc`（中黑）
    - `PingFang SC` → `/System/Library/Fonts/PingFang SC.ttc`（苹方）
    - `Songti SC` → `/System/Library/Fonts/Supplemental/Songti.ttc`（宋体）
  - **macOS 系统拉丁字体（纯 TTF，用于 reportlab）**：
    - `Arial Unicode` → `/System/Library/Fonts/Supplemental/Arial Unicode.ttf`
    - `Georgia` → `/System/Library/Fonts/Supplemental/Georgia.ttf`
    - `Times New Roman` → `/System/Library/Fonts/Supplemental/Times New Roman.ttf`
  - **字体选择原则**：
    - **PDF 报告生成**：使用 reportlab + assets/fonts/ 下的纯 TTF 字体（SimHei/SimKai/STFangSong）
    - Python matplotlib 绘图：使用 `SimHei`、`STHeiti`、`PingFang SC`、`Songti SC` 等中文字体
    - R showtext 渲染：使用 `SimHei`（assets）或 `STHeiti`/`Kaiti SC`（系统）
    - **正式报告场景**：使用 `SimHei` 作为标题，`STFangSong` 作为正文
    - **强调/引用场景**：使用 `SimKai`（楷体）
    - **绝对禁止**：尝试使用 `Arial.ttf`、`Helvetica.ttf` 等不存在的字体文件路径。禁止在 PDF 或图表中使用 `Arial`、`Helvetica`、`Times New Roman` 等不支持中文的字体名称。
    - **注意**：`simfang.ttf` 和 `fzbsk.ttf` 是损坏的字体文件（实际是 HTML 文件），请勿使用。`SimFang` 字体不可用，请使用 `STFangSong` 替代。
- **深度洞察**：能够穿透表面数据，通过多角度关联分析挖掘深层逻辑，明确指出可疑行为并详述推理原因。
- **自主思考**：能根据用户上传的数据，主动提出分析假设并验证。
- **工具专家**：熟练切换并结合 Python (Pandas, Scikit-learn, Seaborn) 和 R (Tidyverse, ggplot2, stats) 的优势进行建模与可视化。你可以通过 Python 的 `rpy2` 库直接调用 R 语言工具开展分析，在进行复杂可视化时，应充分发挥 R 语言 `ggplot2` 包的灵活性优势。**注意：使用 rpy2 (版本 3.x+) 时，请使用 `rpy2.robjects.pandas2ri.activate()` 或 `with rpy2.robjects.conversion.localconverter(rpy2.robjects.default_converter + rpy2.robjects.pandas2ri.converter):` 进行数据转换，不要使用已废弃的 `conversion.register` 属性。**
- **专业严谨**：始终保持专业、严谨的态度，提供具有前瞻性和决策价值的洞察。
- **语言规范**：与用户交流及生成报告时，**默认使用简体中文中文**。除非用户明确指定使用其他语言（如英文），否则不得使用英文或其他语言输出分析内容、报告文本。

**============================================
第七部分：雨途斩棘录 - 错误修正知识库（极重要）
============================================

雨途斩棘录是智能体的错误修正知识库。**系统已在会话开始时自动将历史已知错误及其解决方案注入到你的 context 中。**

**你必须在每次开始新任务前阅读这些历史记录，避免重蹈复辙！**

**雨途斩棘录使用原则**：
1. **启动必读**：无需再手动调用 API 查询，请直接查阅 context 中的历史记录。
2. **预测推理**：结合历史错误，预见当前任务可能出错的环节。
3. **自动记录**：系统会自动捕获并记录新错误的解决方案。

**超级用户管理功能**：
- 超级用户 `rainforgrain` 可以进行备份与恢复。
- 备份：将所有知识点导出为 JSON 文件。
- 恢复：可选择"追加"或"覆盖"方式。

**============================================
第八部分：并行试错与死循环检测（极重要）
============================================

**并行试错与死循环检测（极重要）**：

1. **并行优先原则**：当你需要尝试多种方案（例如：生成 PDF 时尝试不同字体、绘图时尝试不同引擎、导入时尝试不同路径），**必须将这些方案合并到一个代码块中**，通过 `if/elif/else` 或 `try/except` 结构同时验证多条路径，让代码在一次执行中就能确定哪个方案有效。
   - ✅ 正确示范：在一个 `<Code>` 中用 `try: reportlab... except: fpdf2... except: R...` 依次尝试，或用 `if os.path.exists("A"): ... elif os.path.exists("B"): ...`
   - ❌ 错误示范：将方案A放在一个 `<Code>`，等执行失败后在下一次 LLM 调用中再试方案B——这会导致循环和低效。

2. **死循环检测与强制跳出机制（极重要）**：
   - **死循环判定标准**（满足任一即判定）：
     - 相同的数据处理逻辑被执行 ≥3 次且结果相似/相同
     - 相同代码块被执行 ≥2 次且没有产生新的数据洞察
     - 反复尝试相同路径的分析方法 ≥3 次
     - 在同一子任务内停留超过预期的迭代次数
   - **强制跳出操作步骤**：
     1. 立即停止当前循环/重复操作
     2. 记录问题：`[死循环跳过] 任务X - 原因：{具体描述}`
     3. 转向任务清单中的下一项任务
     4. 继续执行，不得卡在原地
   - **禁止行为**：
     - ❌ 死循环中持续尝试同一方法超过2次
     - ❌ 在同一任务内无限迭代
     - ❌ 不记录问题就跳过

3. **终止逻辑**：每一轮完整的分析任务**必须**以 `<Answer>` 标签包裹的最终结论结束。
4. **禁止循环**：禁止在没有新进展的情况下重复生成相同的代码。如果上一次执行已成功或已确定某路径不可行，必须进入下一阶段或输出 `<Answer>`，不得重复相同代码。
5. **一次完成**：每个 `<Code>` 块应尽量完成一个完整的阶段性任务，避免拆分成多个小块依次执行。

**============================================
第九部分：九大原则（极重要）
============================================

**原则一：自我纠错与错误学习**
- 对代码执行中出现的错误，自动识别错误类型（数据缺失、类型错误、维度不匹配、编码问题等），并立即生成修复方案重试。
- 每次错误修复最多重试 3 次，若仍然失败则记录到雨途斩棘录并切换替代方案。
- **错误学习机制（极重要）**：每次遇到错误并找到解决办法后，必须在工作小结中明确记录该错误及其解决方案。系统会自动将其注入到雨途斩棘录中。后续遇到相同或相似错误时，必须先查阅 context 中已注入的雨途斩棘录记录，直接复用已验证的解决方案，不得一错再错。
- 自动修复策略：
  - 缺失字段：尝试模糊匹配或相似列名
  - 类型错误：自动推断并转换数据类型
  - 维度不匹配：自动调整数据形状或选择兼容的操作
  - 编码问题：自动检测并使用 chardet 修复

**原则二：先短代码测试可行性（高效测试）**
- 执行任何复杂分析前，先用 2-5 行代码验证关键假设（文件路径、字段名、数据类型、库可用性）。
- 使用代表性小样本（如前 5 行数据）进行测试。
- 确认可行后再执行完整代码。
- **短代码测试的核心目的**：用最小的代码成本验证最大的不确定性。例如：
  - 验证列名是否正确：`print(df.columns.tolist())`
  - 验证数据类型：`print(df['col'].dtype, df['col'].head(3))`
  - 验证文件路径：`print(os.path.exists('path'), os.listdir('.'))`
  - 验证库函数：`print(type(result), result.shape)`
- 短代码测试结果应作为后续完整分析的基础，不要丢弃测试结论。

**原则三：数据字典优先（极重要 - 分析基础）**
- **分析之初，必须首先建立数据字典（Data Dictionary）**，这是所有后续分析的基础。
- 数据字典建立流程：
  1. 读取所有上传的数据文件（CSV、Excel、JSON等），获取每个文件的字段列表。
  2. 对每个字段记录：字段名、数据类型（数值/类别/日期/文本）、非空率、唯一值数量、代表性样本值（至少5个典型值）。
  3. 分析各字段与用户分析目标之间的关系（核心指标、维度字段、辅助字段等）。
  4. **建立「用户语义 → 数据字段」映射表**：将用户分析目标中的关键词（如"运输工具""国家""企业""金额"等）与数据中实际的字段名（如"车牌号""国别/原产国""经营单位""总价"等）进行语义对应，记录在数据字典中。这一步可以大幅减少后续因用户表达与字段名不一致造成的反复探索和修正。
  5. 将数据字典以结构化方式存储在变量中（如 `data_dict`），供后续所有分析任务引用。
- **后续所有分析代码必须基于数据字典中的字段名进行**，严禁在后续任务中使用未经数据字典验证的字段名。
- 如果数据字典中字段名与代码中引用的名称不一致，必须先查阅数据字典纠正，不得猜测。
- 数据字典应在第一个 `<Code>` 块中建立，并在输出中向用户展示（包括字段映射表）。

**原则四：大任务拆分为小任务（结构化分析流程）**
- 将复杂分析目标分解为结构化任务树。
- 每个叶子节点为一个可独立执行的代码单元。
- 按依赖关系编排执行顺序，支持并行处理无依赖任务。
- 共性处理逻辑提取为公共函数，避免重复代码。
- **结构化分析流程（极重要）**：
  1. **数据探索与数据字典建立**（强制首步）：执行数据探测，建立数据字典，获取字段、类型、缺失值、统计信息。
  2. **分析规划阶段**：基于数据字典和用户指示，结合业务数据特点，设定有层次的分析角度。
  3. **任务确认阶段**：在交互式分析模式下，将分析计划以树形/列表方式发送给用户审查确认。**在这种模式下，只分析用户已经选定的任务列表项**，未被选中的任务不得擅自执行。
  4. **逐项执行阶段**：按任务列表有序执行，每完成一项立即标记完成，基于数据字典引用正确字段名。
  5. **报告生成与交付检查阶段**：汇总所有已完成任务的成果生成最终报告，并进行交付物完整性检查。

**原则五：两种分析模式**
- **交互式分析模式**：
  - 完成数据探索后，将分析计划以有层次的树形或列表方式展示给用户。
  - **必须等待用户选择要执行的分析角度后，才能开始实质性分析（生成代码、执行代码等）。**
  - 用户确认后，只分析用户选定的任务项，不得擅自扩展分析范围。
  - 在执行过程中如需调整计划，应先向用户说明并等待确认。
- **全程代理分析模式**：自主执行全部任务树，不展示中间选择界面。
- 根据当前设置的模式执行对应流程。

**原则六：输出规范与报告质量（极重要）**
- 所有输出使用简体中文。
- 报告支持 PDF、DOCX、PPTX 三种格式。
- 排版统一、美观，使用海关风格模板。
- 图表统一使用 seaborn 暗色/亮色主题。
- **报告内容质量要求（极重要）**：
  - 对于每一个分析角度/维度，报告中必须按以下顺序组织内容：
    1. **文字解读（在图表前面）**：先用一段完整的文字描述该分析维度的情况、趋势、特征和关键发现，引用具体的数据（数值、百分比、排名、增长率等）来支撑分析结论。文字解读应做到有理有据，让用户能够全面、详细地了解分析成果。
    2. **图表引导语**：完成文字解读后，接着写"如下图："，然后换行。
    3. **图表可视化**：在"如下图："之后显示图表。用恰当的图表（柱状图、折线图、饼图、热力图、散点图等）直观展示和印证上述文字解读的观点。图表必须有清晰的标题、坐标轴标签和图例。
  - 禁止出现"分析结果如上图所示"等空泛描述，必须写明具体发现。
  - 每个分析点的描述应做到"先解读、后图表、有数据"。
  - **报告层次结构要求（极重要）**：分析报告应以树形层次结构组织内容，使用以下编号格式：
    ```
    一、[主题]
    （一）[子主题]
    1. [分析点]
    2. [分析点]
    （二）[子主题]
    1. [分析点]
    2. [分析点]
    3. [分析点]
    二、[主题]
    （一）[子主题]
    （二）[子主题]
    三、[主题]
    ……
    ```
  - 每个层级都应有实质内容，不得只有标题而无分析。

**原则七：数据类型正确处理**
- 分析前自动检测数据列类型（数值、类别、日期、文本）。
- 执行类型转换与校验，对不一致或异常类型给出警告并自动修复。
- 日期列使用 pd.to_datetime()，数值列使用 pd.to_numeric(errors='coerce')，类别列使用 .astype('category')。

**原则八：高效处理与代码复用（极重要）**
- 优先复用已有的中间结果、缓存数据或已训练模型。
- **代码复用原则**：
  - 对于确定相同的代码逻辑，不应在同一分析会话中多次重复生成。
  - 一件分析工作应当充分利用已有的生成成果。例如：如果任务A已完成数据清洗得到 `df_clean`，后续任务B、C、D应直接使用 `df_clean`，而不是重新清洗。
  - 将多个任务中共同的处理逻辑（如数据读取、预处理、字体设置等）提取为公共函数，在第一个 `<Code>` 块中定义，后续直接调用。
  - **递进式构建**：分析成果应当递进式累积（如 1→12→123→1234），而非每次从头开始（如 1, 12, 123, 1234 各自独立）。前面的成果必须为后面的任务所用。
- 按依赖关系编排任务，支持并行执行。
- 输出的代码应避免冗余、精简高效。

**原则九：最终交付物检查（极重要）**
- **分析任务完成时，必须进行交付物完整性检查**：
  1. **报告文件检查**：检查用户在"报告类型"中勾选的所有文件格式（PDF、DOCX、PPTX）是否均已生成。如有缺失，必须补充生成。
  2. **分析目标检查**：回顾用户在聊天中提出的分析目标和需求，确认所有分析需求均已覆盖。
  3. **风调雨顺检查**：如果用户在分析过程中通过"风调雨顺"提交了额外需求或条件，检查这些需求是否已在分析中体现。
  4. **图表完整性检查**：确认所有图表均已正确生成并保存到 generated 目录，路径在报告中正确引用。
  5. **报告质量检查**：确认报告中每个分析点都包含"图表+文字描述+数据证据"三要素。
- 如果检查发现任何缺失，必须在输出 `<Answer>` 前补充完成。
- 在 `<Answer>` 中明确列出已交付的文件清单。

**============================================
第十部分：可解释性与智能增强
============================================

**可解释性要求**：
- 对机器学习模型的预测结果，必须输出特征重要性分析（使用 SHAP 或 feature_importances_）。
- 对规则类判断（如虚假申报识别），必须输出判断依据的完整链条。
- 所有结论必须附带数据支撑和推理过程。

**性能与资源管理**：
- 针对不同硬件自动调整批处理大小。
- 大文件分析时使用分块读取（chunksize）和内存监控。
- 数据量过小（<10条记录）或质量过差（>50%缺失值）时，给出明确提示而非强行分析。

**知识库持续学习**：
- 每次分析后，将新发现的规律、异常模式、成功经验总结记录。
- 这些经验将积累为下次分析的参考依据。

**异常情况处理**：
- 当模型无法完成某类分析时，主动建议替代方法。
- 当数据不足以支撑结论时，明确标注"证据不足"。
- 遇到死胡同时立即中止，更换分析思路并通知用户。

**============================================
第十一部分：关系型图数据库分析能力（Neo4j）
============================================

**核心能力**：当用户要求分析实体（如企业、个人、申报单位、运输工具等）之间的关系时，你具备基于图论和关系型数据库思维的穿透式分析能力。

**分析思路与方法**：
1. **实体识别与建模**：
   - 从数据中识别核心实体节点（如：企业、个人、申报单位、运输工具、货物、地址、银行账户等）
   - 识别实体间的关系边（如：申报关系、股权关系、代理关系、运输关系、交易关系等）
   - 使用 Python 的 `networkx` 库构建关系图谱

2. **图谱构建代码示范**：
   ```python
   import networkx as nx
   import matplotlib.pyplot as plt
   
   # 构建关系图
   G = nx.Graph()  # 或 nx.DiGraph() 用于有向关系
   
   # 添加节点（带属性）
   G.add_node("企业A", type="enterprise", risk_level="high")
   G.add_node("个人B", type="person", role="法定代表人")
   
   # 添加关系边（带权重和属性）
   G.add_edge("企业A", "个人B", relation="法定代表人", weight=1.0)
   
   # 分析指标
   degree_centrality = nx.degree_centrality(G)  # 度中心性
   betweenness = nx.betweenness_centrality(G)   # 介数中心性
   communities = nx.community.greedy_modularity_communities(G)  # 社区发现
   ```

3. **穿透式分析策略**：
   - **直接关系分析**：一度关联（直接交易、申报、代理关系）
   - **间接关系分析**：二度及以上关联（通过中间节点连接的隐性关系）
   - **路径分析**：计算任意两个实体间的最短路径和所有路径
   - **社区发现**：识别紧密关联的实体群组（可能的利益共同体）
   - **中心性分析**：找到关系网络中的核心节点（关键控制人/企业）
   - **异常模式识别**：检测星形结构（一个核心控制多个壳公司）、链式结构（层层代理规避监管）

4. **Neo4j 集成（当用户配置了 Neo4j 数据库时）**：
   ```python
   # 如果环境中安装了 neo4j 驱动
   try:
       from neo4j import GraphDatabase
       driver = GraphDatabase.driver(uri, auth=(user, password))
       with driver.session() as session:
           result = session.run("MATCH (a)-[r]->(b) RETURN a, r, b LIMIT 100")
   except ImportError:
       # 回退到 networkx 本地分析
       import networkx as nx
   ```

5. **关系可视化要求**：
   - 使用 `networkx` + `matplotlib` 绘制关系图谱
   - 节点大小按中心性（重要性）设定
   - 边的粗细按交易频次/金额设定
   - 不同类型的实体用不同颜色区分
   - 高风险节点用红色标注
   - 图表必须包含中文标签和图例

6. **实质性主体锁定**：
   - 通过关系距离（路径长度）发现隐性关联
   - 通过中心性指标锁定幕后控制人
   - 通过社区划分识别利益集团
   - 通过时序分析发现关系变化模式（如频繁更换代理的企业）

**注意**：使用 networkx 进行图分析时无需安装 Neo4j。对于大多数数据集，纯 Python 的 networkx 已经足够。仅当用户明确要求连接外部 Neo4j 数据库时才使用 neo4j 驱动。

6. 请始终以这种专业、敏锐且富有洞察力的风格与用户沟通。"""

    # 替换字体目录占位符
    return system_prompt_template.replace("{FONTS_DIR_PLACEHOLDER}", fonts_dir)


def get_compact_system_prompt_with_fonts() -> str:
    """更轻量的系统提示词，避免触发上下文长度上限。"""
    try:
        fonts_dir = str(get_fonts_dir())
    except (NameError, ImportError):
        fonts_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "../../assets/fonts")

    return (
        "你是 DeepAnalyze，负责中国海关风险分析与数据研判。"
        "你的结论必须数据驱动、可解释、可复现，不得臆造字段或结论。\n\n"
        "【输出协议】\n"
        "1) 按需使用 <Analyze>/<Understand>/<Code>/<Execute>/<Answer>/<TaskTree> 标签。\n"
        "2) 每轮应先说明思路，再执行可运行代码，再给结论。\n"
        "3) 代码必须包含必要 import、异常处理，并打印关键中间结果。\n\n"
        "【数据分析硬约束】\n"
        "1) 首次代码必须先做数据探测：数据规模、字段名、字段类型、样例值。\n"
        "2) 仅可使用真实字段；字段缺失时要明确说明并给替代分析方案。\n"
        "3) 图表需有中文标题/坐标轴/图例，保存后给出文件路径。\n"
        "4) 产出文件统一放在 generated 目录，并避免重复生成同类文件。\n\n"
        "【报告交付】\n"
        "1) 结果要体现 数据事实→推理→结论。\n"
        "2) 最终报告按用户勾选格式交付（PDF/DOCX/PPTX）。\n"
        "3) 禁止空壳报告；若某章节证据不足必须显式标注。\n\n"
        "【数据库知识使用】\n"
        "1) 优先使用数据库知识库检索上下文，不要重复粘贴整库结构。\n"
        "2) 若用户要求全库视角，先给摘要，再分层展开关键表字段。\n\n"
        "【工具提示】\n"
        "- 可优先使用 agent_utils.generate_report_pdf/docx/pptx。\n"
        f"- 中文字体目录: {fonts_dir}"
    )


def build_brief_yutu_context(max_items: int = 3, max_chars: int = 1000) -> str:
    """构建精简版雨途斩棘录上下文，避免无上限扩张。"""
    try:
        yutu_data = search_errors(page_size=max_items)
    except Exception as e:
        print(f"[雨途斩棘录] 启动注入失败: {e}")
        return ""

    items = yutu_data.get("items") or []
    if not items:
        return ""

    lines = ["\n\n【雨途斩棘录摘要（近期可复用修复经验）】"]
    for item in items:
        if not isinstance(item, dict):
            continue
        err_type = _shorten_text(item.get("error_type", "未知"), 40)
        err_msg = _shorten_text(item.get("error_message", ""), 80)
        solution = _shorten_text(item.get("solution", ""), 120)
        lines.append(f"- {err_type}: {err_msg} -> {solution}")

    text = "\n".join(lines)
    if len(text) > max_chars:
        text = text[: max_chars - 1] + "…"
    return text


def trim_text_for_budget(text: str, max_chars: int, tail_marker: str = "\n...(已按预算截断)") -> str:
    if len(text) <= max_chars:
        return text
    safe_max = max(0, max_chars - len(tail_marker))
    return text[:safe_max] + tail_marker


def assemble_system_prompt(modules: List[str], max_chars: int = 5600) -> str:
    prompt = "\n\n".join(part.strip() for part in modules if str(part or "").strip())
    return trim_text_for_budget(prompt, max_chars=max_chars)


def bot_stream(
    messages,
    workspace,
    session_id="default",
    username="default",
    strategy="聚焦诉求",
    temperature=0.4,
    analysis_mode="full_agent",
    report_types=None,
    model_provider=None,
    analysis_language="zh-CN",
    selected_database_sources=None,
    source_selection_explicit=False,
):
    analysis_language = normalize_analysis_language(analysis_language)

    # Strategy-specific prompts and default temperature
    strategy_temperatures = {
        "聚焦诉求": 0.2,
        "适度扩展": 0.4,
        "广泛延展": 0.6,
    }
    # Use provided temperature, or fall back to strategy default
    effective_temperature = temperature if temperature is not None else strategy_temperatures.get(strategy, 0.4)

    strategy_prompts = {
        "聚焦诉求": "\n**分析策略：聚焦诉求**。请严格遵守用户指令，仅针对用户直接提出的问题进行分析和回答，不要进行任何不必要的发散或多余的关联分析。保持回答简洁、高效、直击要点。",
        "适度扩展": "\n**分析策略：适度扩展**。在满足用户核心需求的基础上，请基于数据表现进行适量的关联性分析。你可以简要探讨与核心指标相关的其他因素，提供一些背景信息或浅层的风险提示，但请注意分寸，不要过度发散。",
        "广泛延展": "\n**分析策略：广泛延展**。请进行深度发散分析，充分挖掘数据间的潜在联系。你可以大胆发挥想象力，结合海关业务逻辑进行全方位的风险预判、趋势分析和关联挖掘。鼓励你提供多维度的洞察和前瞻性的建议。"
    }

    selected_strategy_prompt = strategy_prompts.get(strategy, strategy_prompts["聚焦诉求"])

    # Temperature-aware analysis depth hint
    if effective_temperature <= 0.15:
        temp_hint = "\n\n**温度/热度提示（当前：极低 {:.2f}）**：请以最高效率完成分析，严格聚焦用户明确提出的需求，不做任何延伸分析。每个分析点直奔结论，不做探索性展开。".format(effective_temperature)
    elif effective_temperature <= 0.3:
        temp_hint = "\n\n**温度/热度提示（当前：较低 {:.2f}）**：提高分析效率，聚焦于用户核心需求，减少不必要的延伸分析。仅对与用户目标直接相关的维度进行分析。".format(effective_temperature)
    elif effective_temperature <= 0.5:
        temp_hint = "\n\n**温度/热度提示（当前：适中 {:.2f}）**：在满足用户核心需求的基础上，可以适度延伸1-2个相关分析角度，但不要过度发散。".format(effective_temperature)
    else:
        temp_hint = "\n\n**温度/热度提示（当前：较高 {:.2f}）**：可以根据数据情况广泛延伸多个角度进行分析，充分挖掘数据价值，提供多维度洞察。鼓励探索性分析和创造性关联。".format(effective_temperature)

    # Report types awareness
    if not report_types:
        report_types = ["pdf"]
    report_types_prompt = "\n\n**用户要求的最终报告交付格式（极重要）**：{}。分析完成后，必须生成这些格式的最终报告文件。在 <Answer> 中确认所有要求的报告类型均已生成。".format("、".join([t.upper() for t in report_types]))

    # Analysis mode prompt injection
    mode_prompts = {
        "interactive": """\n\n**当前分析模式：交互式分析（极重要 - 必须严格遵守）**

在交互式分析模式下，你**必须**严格遵循以下工作流程，不得跳过任何步骤：

**第一步：数据探索与数据字典建立**
- 首先执行数据探测代码，获取所有数据文件的字段名、类型、缺失值、基本统计信息。
- 建立「用户语义 → 数据字段」映射表。
- 将数据探索结果整理为「数据字典」，向用户展示。

**第二步：分析规划与用户确认（强制步骤 - 不得跳过）**
- 基于数据字典和用户的分析目标，设计有层次的分析计划。
- **你必须使用 <TaskTree> 标签输出分析计划，这是强制要求。**
- 标签内部必须是且仅是一个合法的 JSON 对象，不得包含任何其他文字、说明、代码块标记或换行文字。
- **在 <TaskTree> 标签之外，也不要再以代码块或其他形式重复输出任务树 JSON。**
- 如果分析任务比较简单、无需分层次列举，可以将所有分析点平铺为一级列表（不含 children），让用户勾选以决定 TODO 清单。
- 如果分析任务复杂，则分层次组织为树形结构。
- 格式示例（注意：标签内只有纯 JSON）：
<TaskTree>{"tasks":[{"id":"1","name":"分析主题","description":"具体描述","children":[{"id":"1.1","name":"子分析","description":"子说明"}]},{"id":"2","name":"分析主题","description":"具体描述"}]}</TaskTree>
- **输出 <TaskTree> 标签后，立即停止所有输出。** 系统会自动截断后续内容。前端会自动弹出任务选择面板供用户勾选。
- **绝对禁止**在输出 <TaskTree> 后继续输出任何文字（包括"请选择…""以上是…"等）。
- 用户确认后，系统会自动将用户选择回传给你，格式为："用户选择了以下分析任务：[1] xxx, [1.1] xxx"。
- **只执行用户选定的任务，未选中的任务不得擅自执行。**
- **在收到用户选择之前，禁止生成任何分析代码或执行任何分析操作。**

**第三步：逐项执行并报告进度**
- 按用户选定的任务列表有序执行。
- 每完成一项任务，简要报告成果。
- 如果在执行过程中发现需要调整计划，应先向用户说明并等待确认。

**第四步：交付最终报告**
- 所有选定任务完成后，生成最终报告。""",
        "full_agent": "\n\n**当前分析模式：全程代理分析**\n请自主执行全部分析任务，不需要等待用户中间确认。按照任务依赖关系有序执行，确保覆盖所有必要的分析维度。\n\n**高效完成原则（极重要）**：\n- 分析任务应在完成所有必要维度后及时终结，不要无限制地扩展分析。\n- 完成数据探测、核心分析、风险识别后，立即进入报告生成阶段。\n- 生成报告时，必须将前面所有分析步骤中产生的具体数据、结论、图表路径全部整合到报告中，不得遗漏或使用占位符。\n- 报告中每个章节都必须包含实质性的分析内容（具体数值、比例、排名、趋势描述等），而非仅写章节标题。"
    }
    selected_mode_prompt = mode_prompts.get(analysis_mode, mode_prompts["full_agent"])

    if analysis_language == "en":
        strategy_prompts_en = {
            "聚焦诉求": "\n**Analysis strategy: Focused scope**. Follow the user's explicit request strictly, answer only directly asked questions, and avoid unnecessary expansion. Keep the response concise and high-signal.",
            "适度扩展": "\n**Analysis strategy: Moderate extension**. Satisfy the core request first, then provide limited related analysis and lightweight risk hints where relevant, without over-expanding.",
            "广泛延展": "\n**Analysis strategy: Broad exploration**. Perform deeper exploratory analysis, uncover latent relations, and provide multi-dimensional insights and forward-looking recommendations.",
        }
        selected_strategy_prompt = strategy_prompts_en.get(strategy, strategy_prompts_en["聚焦诉求"])

        if effective_temperature <= 0.15:
            temp_hint = "\n\n**Temperature hint (very low {:.2f})**: maximize efficiency, stay strictly on the user's explicit objective, and avoid exploratory detours.".format(effective_temperature)
        elif effective_temperature <= 0.3:
            temp_hint = "\n\n**Temperature hint (low {:.2f})**: prioritize efficient execution and focus on directly relevant dimensions only.".format(effective_temperature)
        elif effective_temperature <= 0.5:
            temp_hint = "\n\n**Temperature hint (medium {:.2f})**: after covering core requirements, you may extend to one or two closely related angles.".format(effective_temperature)
        else:
            temp_hint = "\n\n**Temperature hint (high {:.2f})**: broader exploratory analysis is allowed when it improves insight quality.".format(effective_temperature)

        report_types_prompt = "\n\n**Required final report formats (critical)**: {}. After analysis, generate all required report files and confirm completion in `<Answer>`.".format(", ".join([t.upper() for t in report_types]))

        mode_prompts_en = {
            "interactive": (
                "\n\n**Current mode: Interactive analysis (strictly required)**"
                "\nYou must follow this workflow and do not skip steps:"
                "\n1) Data exploration and data dictionary."
                "\n2) Planning and user confirmation. You MUST output the plan using `<TaskTree>`."
                "\n- Inside `<TaskTree>`, output only one valid JSON object."
                "\n- Do not repeat the task tree JSON outside `<TaskTree>`."
                "\n- After outputting `<TaskTree>`, stop immediately."
                "\n- Before user selection arrives, do not execute analysis code."
                "\n3) Execute only tasks selected by the user."
                "\n4) Deliver final report files in the required formats."
            ),
            "full_agent": (
                "\n\n**Current mode: Full-agent analysis**"
                "\nRun end-to-end analysis autonomously without intermediate user confirmation."
                "\nFinish after necessary dimensions are covered, then move to report generation."
                "\nReports must include concrete findings from earlier steps (numbers, rankings, trend evidence, chart references), not placeholders."
            ),
        }
        selected_mode_prompt = mode_prompts_en.get(analysis_mode, mode_prompts_en["full_agent"])

    requested_report_types = [
        item.strip().lower()
        for item in (report_types or ["pdf"])
        if isinstance(item, str) and item.strip()
    ]
    if not requested_report_types:
        requested_report_types = ["pdf"]
    if analysis_language == "en":
        report_prompt = (
            "\n\n**Report output requirements (mandatory)**:"
            f"\n- Selected formats: {', '.join(requested_report_types).upper()}."
            "\n- In report generation stage, output all selected formats (not a subset)."
            "\n- Keep this requirement throughout the run and produce all required artifacts."
        )
    else:
        report_prompt = (
            "\n\n**报告输出要求（必须遵守）**："
            f"\n- 用户选定的最终报告类型为：{', '.join(requested_report_types).upper()}。"
            "\n- 如果进入报告生成阶段，必须输出以上全部选定格式，不得只生成其中一部分。"
            "\n- 生成报告前请记住该要求，并在完成分析后产出对应成果文件。"
        )

    def _normalize_runtime_db_source(item: Any) -> Optional[Dict[str, Any]]:
        if not isinstance(item, dict):
            return None
        raw_type = item.get("dbType") or item.get("db_type")
        try:
            normalized_type = normalize_db_type(raw_type)
        except Exception:
            return None
        raw_config = item.get("config")
        if not isinstance(raw_config, dict):
            return None
        normalized_config = {
            "host": str(raw_config.get("host", "")).strip(),
            "port": str(raw_config.get("port", "")).strip(),
            "user": str(raw_config.get("user", "")).strip(),
            "password": str(raw_config.get("password", "")),
            "database": str(raw_config.get("database", "")).strip(),
        }
        if not normalized_config["database"]:
            return None
        return {
            "id": str(item.get("id", "")).strip(),
            "label": str(item.get("label", "")).strip(),
            "db_type": normalized_type,
            "config": normalized_config,
        }

    normalized_selected_sources: List[Dict[str, Any]] = []
    if isinstance(selected_database_sources, list):
        for source_item in selected_database_sources:
            normalized_item = _normalize_runtime_db_source(source_item)
            if normalized_item:
                normalized_selected_sources.append(normalized_item)

    configured_sources: List[Dict[str, Any]] = []
    try:
        saved_db_config = _load_user_config(username, "database_connections.json", {"connections": []})
        for conn_item in saved_db_config.get("connections", []):
            normalized_item = _normalize_runtime_db_source(conn_item)
            if normalized_item:
                configured_sources.append(normalized_item)
    except Exception:
        configured_sources = []

    active_sources = normalized_selected_sources
    if not active_sources and len(configured_sources) == 1:
        active_sources = [configured_sources[0]]

    database_source_prompt = ""
    database_data_context = ""
    if active_sources:
        source_lines = []
        for idx, source in enumerate(active_sources, start=1):
            cfg = source["config"]
            if source["db_type"] == "sqlite":
                source_lines.append(
                    f"{idx}. type={source['db_type']}, database={cfg['database']}"
                )
            else:
                source_lines.append(
                    f"{idx}. type={source['db_type']}, host={cfg['host'] or 'localhost'}, "
                    f"port={cfg['port'] or DB_DEFAULT_PORTS.get(source['db_type'], '')}, "
                    f"database={cfg['database']}, user={cfg['user']}, password={cfg['password']}"
                )

        if analysis_language == "en":
            database_source_prompt = (
                "\n\n**Connected database sources (use as primary analysis data source)**:"
                f"\n{chr(10).join(source_lines)}"
                "\n- Start from these connected databases directly. Do not wait for file uploads when they are available."
                "\n- Use these credentials exactly as provided by the user configuration; do not fabricate or hardcode credentials."
            )
            database_data_context = (
                "Connected database sources for this run:\n"
                f"{chr(10).join(source_lines)}\n"
                "Use these sources as the default analysis target."
            )
        else:
            database_source_prompt = (
                "\n\n**已连接数据库数据源（请优先作为本轮分析对象）**："
                f"\n{chr(10).join(source_lines)}"
                "\n- 当存在上述连接时，直接从数据库开始分析，不必等待用户上传文件。"
                "\n- 连接参数必须严格使用用户配置值，不得臆造或硬编码。"
            )
            database_data_context = (
                "本轮已选择数据库数据源：\n"
                f"{chr(10).join(source_lines)}\n"
                "请以这些数据库作为默认分析对象。"
            )
    elif configured_sources and len(configured_sources) > 1 and not source_selection_explicit:
        if analysis_language == "en":
            database_source_prompt = (
                "\n\nUser has configured multiple database sources, but no explicit selection was provided in this round. "
                "If no files are available, ask for source selection before deep analysis."
            )
        else:
            database_source_prompt = (
                "\n\n用户已配置多个数据库数据源，但本轮尚未明确选择。"
                "若当前无文件数据，请先要求用户选择至少一个数据库数据源后再深入分析。"
            )

    loaded_db_context = SESSION_DB_CONTEXT.get(session_id, {})
    loaded_db_source = str(loaded_db_context.get("source_label", "") or "").strip()
    loaded_db_table_count = int(loaded_db_context.get("table_count", 0) or 0)
    loaded_db_column_count = int(loaded_db_context.get("column_count", 0) or 0)
    loaded_db_snapshot_file = str(loaded_db_context.get("snapshot_file", "") or "").strip()

    preferred_source_labels: List[str] = [
        str(item.get("label", "") or "").strip()
        for item in active_sources
        if str(item.get("label", "") or "").strip()
    ]
    if loaded_db_source and loaded_db_source not in preferred_source_labels:
        preferred_source_labels.append(loaded_db_source)

    analysis_history_settings = _load_analysis_history_settings(username)
    history_recorder = AnalysisHistoryRecorder(
        username=username,
        session_id=session_id,
        settings=analysis_history_settings,
        request_summary={
            "strategy": strategy,
            "analysis_mode": analysis_mode,
            "analysis_language": analysis_language,
            "report_types": requested_report_types,
            "incoming_message_count": len(messages or []),
            "workspace_reference_count": len(workspace or []),
            "active_database_sources": [
                _sanitize_runtime_db_source_for_history(item)
                for item in active_sources
            ],
            "source_selection_explicit": bool(source_selection_explicit),
            "loaded_db_context": {
                "source_label": loaded_db_source,
                "table_count": loaded_db_table_count,
                "column_count": loaded_db_column_count,
                "snapshot_file": loaded_db_snapshot_file,
            },
            "model_provider": _sanitize_model_provider_for_history(model_provider),
        },
    )

    if loaded_db_source:
        if analysis_language == "en":
            database_source_prompt += (
                "\n\n**Session DB context loaded**: "
                f"{loaded_db_source} (tables={loaded_db_table_count}, columns={loaded_db_column_count})."
                f" Snapshot file: {loaded_db_snapshot_file or 'N/A'}."
            )
        else:
            database_source_prompt += (
                "\n\n**已加载会话数据库上下文**："
                f"{loaded_db_source}（表{loaded_db_table_count}张，字段{loaded_db_column_count}个）。"
                f" 快照文件：{loaded_db_snapshot_file or '无'}。"
            )

    latest_user_query = ""
    for msg in reversed(messages or []):
        if isinstance(msg, dict) and msg.get("role") == "user":
            latest_user_query = str(msg.get("content", "") or "")
            break

    database_context_started_ts = None
    if active_sources or loaded_db_source:
        database_context_started_ts = time_module.time()
        history_recorder.log(
            stage="database",
            event="database_context_started",
            status="running",
            message="database runtime context assembly started",
            details={
                "selected_source_count": len(active_sources),
                "loaded_context_source": loaded_db_source,
                "loaded_table_count": loaded_db_table_count,
                "loaded_column_count": loaded_db_column_count,
                "snapshot_file": loaded_db_snapshot_file,
            },
        )

    db_kb_context = build_database_knowledge_context(
        username=username,
        user_query=latest_user_query,
        preferred_source_labels=preferred_source_labels,
        analysis_language=analysis_language,
    )
    if db_kb_context:
        database_data_context = (
            f"{database_data_context}\n\n" if database_data_context else ""
        ) + db_kb_context
    if database_data_context:
        database_data_context = trim_text_for_budget(database_data_context, max_chars=2800)
    if database_context_started_ts is not None:
        history_recorder.log(
            stage="database",
            event="database_context_prepared",
            status="running",
            message="database runtime context prepared",
            details={
                "selected_source_count": len(active_sources),
                "loaded_context_source": loaded_db_source,
                "loaded_table_count": loaded_db_table_count,
                "loaded_column_count": loaded_db_column_count,
                "context_chars": len(database_data_context or ""),
                "has_knowledge_context": bool(db_kb_context),
                "duration_ms": int((time_module.time() - database_context_started_ts) * 1000),
            },
        )

    language_prompt = build_analysis_language_prompt(analysis_language)
    runtime_language_prompt = (
        "\n\n**Executable runtime capabilities**:"
        "\n- Python: use `<Code>```python ...```</Code>` for local data processing, modeling, visualization, and report generation."
        "\n- SQL: use `<Code>```sql ...```</Code>` for read-only database extraction. Only SELECT/WITH/EXPLAIN statements are allowed; results are materialized automatically under `workspace/generated/` for later Python/R analysis."
        "\n- R: use `<Code>```r ...```</Code>` for statistical tests, time series, and modeling when R packages are suitable."
        if analysis_language == "en"
        else "\n\n**可执行运行时能力**："
        "\n- Python：使用 `<Code>```python ...```</Code>` 做本地数据处理、建模、可视化和报告生成。"
        "\n- SQL：使用 `<Code>```sql ...```</Code>` 做只读数据库取数；仅允许 SELECT/WITH/EXPLAIN，查询结果会自动物化到 `workspace/generated/` 供后续 Python/R 分析读取。"
        "\n- R：使用 `<Code>```r ...```</Code>` 做统计检验、时间序列和适合 R 包的建模任务。"
    )
    sql_first_prompt = ""
    if active_sources or loaded_db_source:
        sql_first_prompt = (
            "\n\n**SQL-first extraction rule for connected databases**:"
            "\n- For non-trivial database analysis, do not start by scanning large tables from Python/R. Start with an extraction plan."
            "\n- First identify required entities, grains, metrics, filters, join paths, and expected output tables/files."
            "\n- Then execute one or more read-only SQL blocks to build smaller materialized extracts, cubes, samples, object lists, or candidate-risk sets."
            "\n- Use Python/R only after SQL has reduced the data to a manageable analytical dataset."
            "\n- Prefer multiple targeted SQL blocks over one huge query; each block should state its purpose in `<Analyze>` before execution."
            if analysis_language == "en"
            else "\n\n**已连接数据库的 SQL-first 取数规则**："
            "\n- 面对非平凡数据库分析任务时，不要一开始就用 Python/R 扫描大表；必须先制定取数计划。"
            "\n- 先识别分析所需主体、粒度、指标、过滤条件、Join 路径和预期输出数据集。"
            "\n- 然后用一段或多段只读 SQL 构建更小的物化取数结果、汇总立方体、样本、对象清单或风险候选集。"
            "\n- 只有在 SQL 已经把数据压缩为可承载的分析数据集之后，再使用 Python/R 做深度分析、统计建模和可视化。"
            "\n- 优先使用多段有明确目的的小 SQL，而不是一条巨大 SQL；每段 SQL 执行前应在 `<Analyze>` 中说明目的。"
        )
    signature_safety_prompt = (
        "\n\n**函数调用安全（必须遵守）**："
        "\n- 先核对函数签名与参数，再执行。"
        "\n- 对第三方库仅使用已支持参数，避免臆造 keyword。"
        "\n- 复用函数时确保定义与调用一致。"
    )
    methodology_integration_prompt = (
        "\n\n**方法执行要求（简版）**："
        "\n- 遵循 数据理解→假设→验证→结论→交付 的闭环。"
        "\n- `<Analyze>` 说明当前假设与验证计划；`<Code>` 执行可验证动作；`<Answer>` 给出可解释结论。"
        "\n- 证据不足时必须明确标注并提出补充数据建议。"
    )

    system_prompt = assemble_system_prompt(
        modules=[
            get_compact_system_prompt_with_fonts(),
            build_brief_yutu_context(max_items=3, max_chars=900),
            selected_strategy_prompt,
            selected_mode_prompt,
            temp_hint,
            report_types_prompt,
            signature_safety_prompt,
            methodology_integration_prompt,
            report_prompt,
            database_source_prompt,
            runtime_language_prompt,
            sql_first_prompt,
            language_prompt,
        ],
        max_chars=5600,
    )

    language_enforcer_prompt = (
        "CRITICAL OUTPUT RULE: Respond in English only for all user-facing narrative text. "
        "Do not output Chinese characters in Analyze/Understand/Answer/TaskTree/report text. "
        "Keep structural tags unchanged."
        if analysis_language == "en"
        else "关键输出规则：所有面向用户的自然语言文本必须使用简体中文（包括 Analyze/Understand/Answer/TaskTree/报告正文）。结构化标签保持不变。"
    )

    # Check if system prompt is already there, if not, insert it
    if not messages or messages[0]["role"] != "system":
        messages.insert(0, {"role": "system", "content": system_prompt})

    if len(messages) < 2 or messages[1].get("role") != "system" or messages[1].get("content") != language_enforcer_prompt:
        messages.insert(1, {"role": "system", "content": language_enforcer_prompt})

    original_cwd = os.getcwd()
    WORKSPACE_DIR = get_session_workspace(session_id, username)
    os.makedirs(WORKSPACE_DIR, exist_ok=True)
    # 创建 generated 子文件夹用于存放代码生成的文件
    GENERATED_DIR = os.path.join(WORKSPACE_DIR, "generated")
    os.makedirs(GENERATED_DIR, exist_ok=True)
    # 记录本轮开始前已有的报告文件，便于判断当前轮是否真正产出新报告
    preexisting_report_files: Dict[str, set[str]] = {
        "pdf": set(),
        "docx": set(),
        "pptx": set(),
    }
    try:
        generated_dir_path = Path(GENERATED_DIR)
        for existing in generated_dir_path.iterdir():
            if not existing.is_file():
                continue
            ext = existing.suffix.lower().lstrip(".")
            if ext in preexisting_report_files:
                preexisting_report_files[ext].add(existing.name)
    except Exception:
        pass
    # 创建 converted 子文件夹用于存放 UTF-8 转换后的文件
    CONVERTED_DIR = os.path.join(WORKSPACE_DIR, "converted")
    os.makedirs(CONVERTED_DIR, exist_ok=True)
    history_recorder.log(
        stage="prompt",
        event="prompt_prepared",
        status="running",
        message="system prompt assembled",
        details={
            "system_prompt_chars": len(system_prompt or ""),
            "language_enforcer_chars": len(language_enforcer_prompt or ""),
            "database_context_chars": len(database_data_context or ""),
            "database_source_prompt_chars": len(database_source_prompt or ""),
            "system_prompt_preview": _truncate_history_text(system_prompt, 1200)
            if analysis_history_settings.get("capture_prompt_preview", True)
            else "",
        },
    )
    # print(messages)
    if messages and messages[0]["role"] == "assistant":
        messages = messages[1:]
    elif len(messages) > 1 and messages[0]["role"] == "system" and messages[1]["role"] == "assistant":
        # Keep system prompt, but skip the first assistant message if it's misplaced
        messages = [messages[0]] + messages[2:]

    if messages and messages[-1]["role"] == "user":
        user_message = messages[-1]["content"]
        # 总是使用当前 session 的 WORKSPACE_DIR 获取文件信息
        file_info = collect_file_info(WORKSPACE_DIR)
        if file_info:
            file_info = trim_text_for_budget(file_info, max_chars=2200)
        data_sections: List[str] = []
        if file_info:
            data_sections.append(file_info)
        if database_data_context:
            data_sections.append(database_data_context)

        if data_sections:
            merged_data_sections = trim_text_for_budget("\n\n".join(data_sections), max_chars=3800)
            messages[-1][
                "content"
            ] = f"# Instruction\n{user_message}\n\n# Data\n{merged_data_sections}"
        else:
            messages[-1]["content"] = f"# Instruction\n{user_message}"
        history_recorder.log(
            stage="prompt",
            event="user_message_prepared",
            status="running",
            message="user instruction merged with runtime data context",
            details={
                "user_message_preview": _truncate_history_text(user_message, 600),
                "file_info_chars": len(file_info or ""),
                "merged_data_chars": len((messages[-1].get("content") or "")),
                "has_database_context": bool(database_data_context),
            },
        )
    # print("111",messages)
    initial_workspace = set(workspace)
    assistant_reply = ""
    finished = False
    exe_output = None
    llm_client, llm_model, runtime_provider = get_runtime_llm(model_provider)
    round_count = 0
    history_final_status = "completed"
    history_final_message = "analysis session finished"
    while not finished and round_count < MAX_AGENT_ROUNDS:
        round_count += 1
        history_recorder.log(
            stage="round",
            event="round_started",
            status="running",
            message=f"analysis round {round_count} started",
            details={
                "round": round_count,
                "message_count": len(messages),
                "workspace_dir": WORKSPACE_DIR,
            },
        )
        # 在每次 LLM 生成前，检查是否有用户提交的“过程指导/Side Task”
        guidance = SESSION_GUIDANCE.pop(session_id, None)
        if guidance:
            guidance_msg = {
                "role": "user",
                "content": f"【风调雨顺 - 过程指导/Side Task】: {guidance}\n\n请将此需求、方法或条件与你当前正在进行的任务结合起来考虑，并在下一步分析方案中予以体现。请确保不中断整体分析项目的连贯性。"
            }
            messages.append(guidance_msg)
            print(f"[Side Guidance] Injecting guidance into session {session_id}")
            history_recorder.log(
                stage="guidance",
                event="guidance_injected",
                status="running",
                message="side guidance injected into current round",
                details={
                    "round": round_count,
                    "guidance_preview": _truncate_history_text(guidance, 500),
                },
            )

        request_kwargs: Dict[str, Any] = {
            "model": llm_model,
            "messages": messages,
            "temperature": effective_temperature,
            "stream": True,
            "stop": [
                "</Code>",
                "</code>",
                "</TaskTree>",
                "</tasktree>",
                "</s>",
                "<|endoftext|>",
                "<|im_end|>",
            ],
        }
        if _provider_supports_vllm_controls(runtime_provider):
            request_kwargs["extra_body"] = {
                "add_generation_prompt": False,
                "max_new_tokens": 32768,
            }

        history_recorder.log(
            stage="llm",
            event="llm_request_started",
            status="running",
            message=f"llm request started for round {round_count}",
            details={
                "round": round_count,
                "provider": runtime_provider,
                "model": llm_model,
                "message_count": len(messages),
                "temperature": effective_temperature,
                "has_extra_body": "extra_body" in request_kwargs,
            },
        )

        try:
            response = _create_chat_completion_with_retry(llm_client, request_kwargs)
        except Exception as create_exc:
            history_final_status = "failed"
            history_final_message = f"llm request failed: {create_exc}"
            history_recorder.log(
                stage="llm",
                event="llm_request_failed",
                status="failed",
                message=history_final_message,
                details={"round": round_count},
            )
            error_block = _build_streaming_answer_error(str(create_exc))
            assistant_reply += error_block
            yield error_block
            break

        cur_res = ""
        last_finish_reason = None
        saw_final_answer_this_round = False
        stream_chunk_count = 0
        stream_char_count = 0
        next_chunk_progress = int(analysis_history_settings.get("stream_progress_chunk_interval", 40))
        next_char_progress = int(analysis_history_settings.get("stream_progress_char_interval", 1600))
        for chunk in response:
            if chunk.choices:
                choice = chunk.choices[0]
                last_finish_reason = choice.finish_reason
                if choice.delta.content is not None:
                    delta = choice.delta.content
                    cur_res += delta
                    assistant_reply += delta
                    stream_chunk_count += 1
                    stream_char_count += len(delta)
                    if analysis_history_settings.get("capture_stream_progress", True) and (
                        stream_chunk_count >= next_chunk_progress or stream_char_count >= next_char_progress
                    ):
                        history_recorder.log(
                            stage="llm",
                            event="llm_stream_progress",
                            status="running",
                            message=f"round {round_count} streamed {stream_char_count} chars",
                            details={
                                "round": round_count,
                                "chunk_count": stream_chunk_count,
                                "char_count": stream_char_count,
                            },
                        )
                        next_chunk_progress += int(analysis_history_settings.get("stream_progress_chunk_interval", 40))
                        next_char_progress += int(analysis_history_settings.get("stream_progress_char_interval", 1600))
                    yield delta

            if re.search(r"</answer>", cur_res, re.IGNORECASE):
                saw_final_answer_this_round = True
                break
            # ---- TaskTree 检测（流式中途发现完整闭合标签） ----
            if re.search(r"</tasktree>", cur_res, re.IGNORECASE):
                messages.append({"role": "assistant", "content": cur_res})
                finished = True
                break
        history_recorder.log(
            stage="llm",
            event="llm_response_completed",
            status="running",
            message=f"llm response completed for round {round_count}",
            details={
                "round": round_count,
                "finish_reason": last_finish_reason,
                "chunk_count": stream_chunk_count,
                "char_count": stream_char_count,
                "saw_final_answer": saw_final_answer_this_round,
                "has_tasktree": bool(re.search(r"<tasktree>", cur_res, re.IGNORECASE)),
                "has_code": bool(re.search(r"<code>|```", cur_res, re.IGNORECASE)),
            },
        )
        if finished:
            # 已因 </Answer> 或 </TaskTree> 结束，跳过后续代码执行逻辑
            continue

        # ---- TaskTree 检测（stop sequence 截断时，LLM 输出了 <TaskTree> 但未闭合） ----
        has_tasktree_open = re.search(r"<tasktree>", cur_res, re.IGNORECASE) is not None
        has_tasktree_close = re.search(r"</tasktree>", cur_res, re.IGNORECASE) is not None
        if has_tasktree_open and not has_tasktree_close:
            # stop sequence 截断了 </TaskTree>，手动补上闭合标签
            closing_tag = "</TaskTree>"
            cur_res += closing_tag
            assistant_reply += closing_tag
            yield closing_tag
            messages.append({"role": "assistant", "content": cur_res})
            history_recorder.log(
                stage="planner",
                event="tasktree_completed",
                status="running",
                message="task tree output was auto-closed after stop sequence",
                details={"round": round_count},
            )
            finished = True
            continue

        has_code_open = re.search(r"<code>", cur_res, re.IGNORECASE) is not None
        has_code_close = re.search(r"</code>", cur_res, re.IGNORECASE) is not None
        has_fenced_code = re.search(r"```\s*(?:python|py|sql|postgresql|mysql|sqlite|mssql|oracle|r|rscript)?\s*\n.*?```", cur_res, re.DOTALL | re.IGNORECASE) is not None

        if last_finish_reason == "stop" and has_code_open and not has_code_close:
            if not re.search(r"</code>\s*$", cur_res, re.IGNORECASE):
                missing_tag = "</Code>"
                cur_res += missing_tag
                assistant_reply += missing_tag
                yield missing_tag
            has_code_close = True

        should_execute_code = (has_code_open and has_code_close) or ((not has_code_open) and has_fenced_code)

        if should_execute_code:
            messages.append({"role": "assistant", "content": cur_res})
            executable_block = _extract_executable_block(cur_res)
            code_str = executable_block.get("code") or None
            code_language = executable_block.get("language") or "python"
            execution_stage = "sql" if code_language == "sql" else "r" if code_language == "r" else "code"
            execution_label = "SQL" if code_language == "sql" else "R" if code_language == "r" else "code"

            if code_str:
                execution_started_ts = time_module.time()
                if code_language == "python":
                    code_str = Chinese_matplot_str + "\n" + code_str
                    # 自动将用户提及的原始文件名映射为 converted/ 目录下的实际文件
                    code_str = _rewrite_file_paths(code_str, WORKSPACE_DIR)
                elif code_language == "r":
                    code_str = _rewrite_file_paths(code_str, WORKSPACE_DIR)
                history_recorder.log(
                    stage=execution_stage,
                    event=f"{execution_stage}_execution_started",
                    status="running",
                    message=f"{execution_label} execution started for round {round_count}",
                    details={
                        "round": round_count,
                        "language": code_language,
                        "raw_language": executable_block.get("raw_language") or "",
                        "code_chars": len(code_str),
                        "code_preview": _truncate_history_text(code_str, 700),
                    },
                )
                # 执行前快照（路径 -> (size, mtime)）
                try:
                    before_state = {
                        p.resolve(): (p.stat().st_size, p.stat().st_mtime_ns)
                        for p in Path(WORKSPACE_DIR).rglob("*")
                        if p.is_file()
                    }
                except Exception:
                    before_state = {}
                # 在固定工作区按语言执行
                if code_language == "sql":
                    sql_source = active_sources[0] if active_sources else None
                    exe_output = execute_sql_safe(code_str, WORKSPACE_DIR, sql_source)
                elif code_language == "r":
                    exe_output = execute_r_code_safe(code_str, WORKSPACE_DIR)
                else:
                    exe_output = execute_code_safe(code_str, WORKSPACE_DIR)
                error_info: Dict[str, Any] = {}

                # ========== 自动检测并记录错误到雨途斩棘录 ==========
                try:
                    error_info = detect_and_record_error(exe_output, code_str, WORKSPACE_DIR)
                    if error_info.get("has_error"):
                        print(f"[雨途斩棘录] 检测到错误: {error_info.get('error_type')}")
                        if error_info.get("recorded"):
                            print(f"[雨途斩棘录] 已自动记录到知识库")
                        elif error_info.get("similar_found"):
                            print(f"[雨途斩棘录] 发现相似错误记录，跳过重复记录")
                except Exception as e:
                    print(f"[雨途斩棘录] 错误检测失败: {e}")
                # ========== 错误检测结束 ==========

                # 执行后快照
                try:
                    after_state = {
                        p.resolve(): (p.stat().st_size, p.stat().st_mtime_ns)
                        for p in Path(WORKSPACE_DIR).rglob("*")
                        if p.is_file()
                    }
                except Exception:
                    after_state = {}
                # 计算新增与修改
                added_paths = [p for p in after_state.keys() if p not in before_state]
                modified_paths = [
                    p
                    for p in after_state.keys()
                    if p in before_state and after_state[p] != before_state[p]
                ]

                # 将新增和修改的文件移动到 generated 文件夹
                artifact_paths = []
                for p in added_paths:
                    try:
                        # 如果文件不在 generated 文件夹中，移动它
                        if not str(p).startswith(GENERATED_DIR):
                            dest_path = Path(GENERATED_DIR) / p.name
                            dest_path = uniquify_path(dest_path)
                            shutil.copy2(str(p), str(dest_path))
                            artifact_paths.append(dest_path.resolve())
                        else:
                            artifact_paths.append(p)
                    except Exception as e:
                        print(f"Error moving file {p}: {e}")
                        artifact_paths.append(p)

                # 为修改的文件生成副本并移动到 generated 文件夹
                for p in modified_paths:
                    try:
                        dest_name = f"{Path(p).stem}_modified{Path(p).suffix}"
                        dest_path = Path(GENERATED_DIR) / dest_name
                        dest_path = uniquify_path(dest_path)
                        shutil.copy2(p, dest_path)
                        artifact_paths.append(dest_path.resolve())
                    except Exception as e:
                        print(f"Error copying modified file {p}: {e}")

                # 旧：Execute 内部放控制台输出；新：追加 <File> 段落给前端渲染卡片
                exe_str = f"\n<Execute>\n```\n{exe_output}\n```\n</Execute>\n"
                file_block = ""
                if artifact_paths:
                    lines = ["<File>"]
                    for p in artifact_paths:
                        try:
                            rel = (
                                Path(p)
                                .relative_to(Path(WORKSPACE_DIR).resolve())
                                .as_posix()
                            )
                        except Exception:
                            rel = Path(p).name
                        # 在相对路径前加上 username/session_id 前缀
                        url = build_download_url(f"{username}/{session_id}/{rel}")
                        name = Path(p).name
                        lines.append(f"- [{name}]({url})")
                        if Path(p).suffix.lower() in [
                            ".png",
                            ".jpg",
                            ".jpeg",
                            ".gif",
                            ".webp",
                            ".svg",
                        ]:
                            lines.append(f"![{name}]({url})")
                    lines.append("</File>")
                    file_block = "\n" + "\n".join(lines) + "\n"
                full_execution_block = exe_str + file_block
                history_recorder.log(
                    stage=execution_stage,
                    event=f"{execution_stage}_execution_completed",
                    status="warning" if error_info.get("has_error") else "running",
                    message=(
                        f"{execution_label} execution finished with detected error: {error_info.get('error_type', 'unknown')}"
                        if error_info.get("has_error")
                        else f"{execution_label} execution finished for round {round_count}"
                    ),
                    details={
                        "round": round_count,
                        "language": code_language,
                        "duration_ms": int((time_module.time() - execution_started_ts) * 1000),
                        "output_chars": len(exe_output or ""),
                        "has_error": bool(error_info.get("has_error")),
                        "error_type": str(error_info.get("error_type", "") or ""),
                        "artifact_count": len(artifact_paths),
                        "artifacts": [Path(path).name for path in artifact_paths[:20]],
                        "added_file_count": len(added_paths),
                        "modified_file_count": len(modified_paths),
                    },
                )
                assistant_reply += full_execution_block
                yield full_execution_block
                messages.append({"role": "execute", "content": f"{exe_output}"})
                # 刷新工作区快照（路径集合）
                current_files = set(
                    [
                        os.path.join(WORKSPACE_DIR, f)
                        for f in os.listdir(WORKSPACE_DIR)
                        if os.path.isfile(os.path.join(WORKSPACE_DIR, f))
                    ]
                )
                new_files = list(current_files - initial_workspace)
                if new_files:
                    workspace.extend(new_files)
                    initial_workspace.update(new_files)

                if saw_final_answer_this_round:
                    finished = True

            continue

        if saw_final_answer_this_round:
            if cur_res.strip():
                messages.append({"role": "assistant", "content": cur_res})
                history_recorder.log(
                    stage="answer",
                    event="final_answer_recorded",
                    status="running",
                    message=f"final answer captured in round {round_count}",
                    details={
                        "round": round_count,
                        "answer_chars": len(cur_res),
                    },
                )
            finished = True
            continue

        # 无代码时也要保留本轮输出，允许模型进入下一轮（例如先 Analyze 再 Code）
        if cur_res.strip():
            messages.append({"role": "assistant", "content": cur_res})
            history_recorder.log(
                stage="round",
                event="assistant_round_saved",
                status="running",
                message=f"assistant content retained for next round {round_count}",
                details={
                    "round": round_count,
                    "content_chars": len(cur_res),
                },
            )

    if not finished and round_count >= MAX_AGENT_ROUNDS:
        history_final_status = "warning"
        history_final_message = f"analysis stopped after reaching max rounds ({MAX_AGENT_ROUNDS})"
        history_recorder.log(
            stage="session",
            event="max_round_limit_reached",
            status="warning",
            message=history_final_message,
            details={"round_count": round_count},
        )
        limit_block = (
            "<Answer>\n"
            f"已达到最大分析轮次（{MAX_AGENT_ROUNDS}轮），为避免无限循环已自动停止。\n"
            "请提供更具体的下一步指令，或缩小分析范围后继续。\n"
            "</Answer>\n"
        )
        assistant_reply += limit_block
        yield limit_block

    # 兜底：若本轮缺少所需报告产物，则服务端自动导出，确保前端文件输出区可见。
    try:
        requested_report_types = [
            str(item).strip().lower()
            for item in (report_types or ["pdf"])
            if str(item).strip().lower() in {"pdf", "docx", "pptx"}
        ]
        if not requested_report_types:
            requested_report_types = ["pdf"]

        has_final_answer = re.search(r"</answer>", assistant_reply, re.IGNORECASE) is not None
        if has_final_answer:
            current_report_files: Dict[str, set[str]] = {
                "pdf": set(),
                "docx": set(),
                "pptx": set(),
            }
            generated_dir_path = Path(GENERATED_DIR)
            if generated_dir_path.exists():
                for current_file in generated_dir_path.iterdir():
                    if not current_file.is_file():
                        continue
                    ext = current_file.suffix.lower().lstrip(".")
                    if ext in current_report_files:
                        current_report_files[ext].add(current_file.name)

            missing_types = [
                ext for ext in requested_report_types
                if not (current_report_files.get(ext, set()) - preexisting_report_files.get(ext, set()))
            ]

            if missing_types:
                history_recorder.log(
                    stage="report",
                    event="report_fallback_started",
                    status="running",
                    message="report fallback started to ensure requested artifacts exist",
                    details={
                        "missing_types": missing_types,
                        "requested_report_types": requested_report_types,
                    },
                )
                report_history = [
                    msg for msg in messages
                    if isinstance(msg, dict) and str(msg.get("role", "")).lower() in {"user", "assistant", "execute"}
                ]
                if not report_history or str(report_history[-1].get("role", "")).lower() != "assistant":
                    report_history.append({"role": "assistant", "content": assistant_reply})

                md_text = _extract_sections_from_messages(report_history)
                if not md_text:
                    md_text = (
                        "No structured report sections were extracted; fallback to final answer text."
                        if analysis_language == "en"
                        else "未提取到结构化报告段落，已回退使用最终回答文本。"
                    )
                    md_text += "\n\n" + (assistant_reply or "")

                ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                base_name = f"Auto_Report_{ts}"
                auto_generated_paths: List[Path] = []

                md_path = _save_md(md_text, base_name, GENERATED_DIR)
                if md_path:
                    auto_generated_paths.append(md_path)

                if "pdf" in missing_types:
                    try:
                        pdf_path = _save_pdf(md_text, base_name, GENERATED_DIR)
                        if pdf_path:
                            auto_generated_paths.append(pdf_path)
                    except Exception as pdf_error:
                        print(f"Auto report fallback pdf generation failed: {pdf_error}")
                if "docx" in missing_types:
                    try:
                        docx_path = _save_docx(md_text, base_name, GENERATED_DIR)
                        if docx_path:
                            auto_generated_paths.append(docx_path)
                    except Exception as docx_error:
                        print(f"Auto report fallback docx generation failed: {docx_error}")
                if "pptx" in missing_types:
                    try:
                        pptx_path = _save_pptx(md_text, base_name, GENERATED_DIR)
                        if pptx_path:
                            auto_generated_paths.append(pptx_path)
                    except Exception as pptx_error:
                        print(f"Auto report fallback pptx generation failed: {pptx_error}")

                if auto_generated_paths:
                    lines = ["<File>"]
                    for artifact in auto_generated_paths:
                        name = Path(artifact).name
                        url = build_download_url(f"{username}/{session_id}/generated/{name}")
                        lines.append(f"- [{name}]({url})")
                    lines.append("</File>")
                    auto_file_block = "\n" + "\n".join(lines) + "\n"
                    assistant_reply += auto_file_block
                    yield auto_file_block
                    history_recorder.log(
                        stage="report",
                        event="report_fallback_completed",
                        status="running",
                        message="report fallback generated missing artifacts",
                        details={
                            "generated_files": [Path(item).name for item in auto_generated_paths],
                            "missing_types": missing_types,
                        },
                    )
    except Exception as auto_report_error:
        history_final_status = "warning" if history_final_status == "completed" else history_final_status
        history_final_message = f"auto report fallback failed: {auto_report_error}"
        history_recorder.log(
            stage="report",
            event="report_fallback_failed",
            status="warning",
            message=history_final_message,
            details={},
        )
        print(f"Auto report fallback failed: {auto_report_error}")

    final_has_answer = re.search(r"</answer>", assistant_reply, re.IGNORECASE) is not None
    if history_final_status == "completed" and not final_has_answer:
        history_final_status = "warning"
        history_final_message = "analysis finished without a complete </Answer> block"
    history_recorder.finalize(
        status=history_final_status,
        message=history_final_message,
        details={
            "round_count": round_count,
            "assistant_reply_chars": len(assistant_reply),
            "has_final_answer": final_has_answer,
            "workspace_dir": WORKSPACE_DIR,
            "generated_dir": GENERATED_DIR,
        },
    )
    os.chdir(original_cwd)


@app.post("/chat/completions")
async def chat(body: dict = Body(...)):
    messages = body.get("messages", [])
    workspace = body.get("workspace", [])
    session_id = body.get("session_id", "default")
    username = body.get("username", "default")
    strategy = body.get("strategy", "聚焦诉求")
    temperature = body.get("temperature", None)  # Optional: user can override temperature
    analysis_mode = body.get("analysis_mode", "full_agent")
    analysis_language = normalize_analysis_language(body.get("analysis_language", "zh-CN"))
    report_types = body.get("report_types", ["pdf"])
    model_provider = body.get("model_provider")
    selected_database_sources = body.get("selected_database_sources", [])
    source_selection_explicit = bool(body.get("source_selection_explicit", False))
    model_name = MODEL_PATH
    if isinstance(model_provider, dict):
        model_name = (model_provider.get("model") or MODEL_PATH)

    # 动态构建 workspace 目录，确保能正确识别当前 session 的文件
    actual_workspace_dir = get_session_workspace(session_id, username)

    def generate():
        for delta_content in bot_stream(
            messages,
            workspace,
            session_id,
            username,
            strategy,
            temperature,
            analysis_mode,
            report_types,
            model_provider,
            analysis_language,
            selected_database_sources,
            source_selection_explicit,
        ):
            # print(delta_content)
            chunk = {
                "id": "chatcmpl-stream",
                "object": "chat.completion.chunk",  # 标识为流式块
                "created": 1677652288,
                "model": model_name,
                "choices": [
                    {
                        "index": 0,
                        # 3. 使用 delta 字段而非 message 字段
                        "delta": {
                            "content": delta_content  # 直接填入原始内容，不要调用 fix_tags
                        },
                        "finish_reason": None,  # 传输中为 None
                    }
                ],
            }

            yield json.dumps(chunk) + "\n"
            # 5. 循环结束后，发送一个结束标记 (Optional, 但推荐)
        end_chunk = {
            "id": "chatcmpl-stream",
            "object": "chat.completion.chunk",
            "created": 1677652288,
            "model": model_name,
            "choices": [{"index": 0, "delta": {}, "finish_reason": "stop"}]
        }
        yield json.dumps(end_chunk) + "\n"

    return StreamingResponse(generate(), media_type="text/plain")


# -------- Export Report (PDF + MD) --------
from datetime import datetime


def _extract_sections_from_messages(messages: list[dict]) -> str:
    """从历史消息中抽取报告内容。
    
    优先使用 <Answer> 标签中的内容作为报告主体。
    同时从 <Analyze> 标签中提取分析过程的实质内容作为补充。
    从 <Execute> 标签中提取执行结果中的关键数据。
    最终组装成一份完整的、有实质内容的分析报告。
    """
    if not isinstance(messages, list):
        return ""
    import re as _re

    tag_pattern = r"<(Analyze|Understand|Code|Execute|File|Answer|TaskTree)>([\s\S]*?)</\1>"

    # 收集所有助手消息内容
    all_content = ""
    for idx, m in enumerate(messages, start=1):
        role = (m or {}).get("role")
        if role != "assistant":
            continue
        all_content += str((m or {}).get("content") or "") + "\n"

    # 分类提取各类标签内容
    answer_sections = []
    analyze_sections = []
    execute_data = []

    for match in _re.finditer(tag_pattern, all_content, _re.DOTALL):
        tag, seg = match.groups()
        seg = seg.strip()
        if not seg:
            continue
        if tag == "Answer":
            answer_sections.append(seg)
        elif tag == "Analyze":
            # 去掉预测推理和短代码测试部分，只保留正式分析内容
            clean_seg = _re.sub(r"#\s*(预测推理|短代码测试与结果)\s*\n[\s\S]*?(?=\n#\s*正式分析|$)", "", seg, flags=_re.DOTALL).strip()
            if clean_seg and len(clean_seg) > 50:  # 只保留有实质内容的分析
                analyze_sections.append(clean_seg)
        elif tag == "Execute":
            # 提取执行结果中的有用数据（排除纯错误信息）
            if seg and "Traceback" not in seg and "Error" not in seg and len(seg) > 20:
                # 截断过长的执行输出
                if len(seg) > 2000:
                    seg = seg[:2000] + "\n...(输出已截断)"
                execute_data.append(seg)

    # 构建最终报告文本
    final_text = ""

    # 1. 主要使用 Answer 标签的内容
    if answer_sections:
        final_text = "\n\n".join(answer_sections).strip()

    # 2. 如果 Answer 内容过短或为空，用 Analyze 内容补充
    if not final_text or len(final_text) < 200:
        if analyze_sections:
            analyze_text = "\n\n".join(analyze_sections).strip()
            if final_text:
                final_text = analyze_text + "\n\n" + final_text
            else:
                final_text = analyze_text

    # 3. 检测占位符模式并尝试替换
    placeholder_pattern = r'\[详细.*?(?:描述|分析|评估|建议)\.{0,3}\]'
    if _re.search(placeholder_pattern, final_text):
        # 存在占位符，尝试从分析过程中提取实际内容填充
        if analyze_sections or execute_data:
            supplementary = "\n\n## 补充分析数据\n\n"
            if analyze_sections:
                for i, section in enumerate(analyze_sections[:5], 1):
                    supplementary += f"### 分析维度 {i}\n\n{section}\n\n"
            if execute_data:
                supplementary += "### 关键数据摘要\n\n"
                for data in execute_data[:3]:
                    supplementary += f"```\n{data}\n```\n\n"
            final_text += supplementary

    # 4. 如果仍然没有内容，从整体对话中提取
    if not final_text:
        # 最后兜底：取所有助手消息中有实质内容的部分
        for m in messages:
            role = (m or {}).get("role")
            if role == "assistant":
                content = str((m or {}).get("content") or "")
                # 移除标签
                clean = _re.sub(r'</?(?:Analyze|Understand|Code|Execute|File|Answer|TaskTree)>', '', content).strip()
                if clean and len(clean) > 100:
                    final_text += clean + "\n\n"

    return final_text.strip()


def _save_md(md_text: str, base_name: str, workspace_dir: str) -> Path:
    Path(workspace_dir).mkdir(parents=True, exist_ok=True)
    md_path = uniquify_path(Path(workspace_dir) / f"{base_name}.md")
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(md_text)
    return md_path


import pypandoc


def _get_available_fonts() -> dict:
    """
    返回所有可用字体信息的字典。
    返回结构：{
        "assets": [(font_path, font_name, font_description), ...],
        "system": [(font_path, font_name, font_description), ...],
        "reportlab_fallback": (font_path, font_name),  # 纯 TTF，无 TTC
    }
    """
    import reportlab.pdfbase.ttfonts as ttfonts

    fonts = {"assets": [], "system": [], "reportlab_fallback": None}
    fonts_dir = _FONT_DIR if '_FONT_DIR' in globals() else os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "../../assets/fonts"
    )

    # === assets 字体（纯 TTF，reportlab/fpdf2/matplotlib 均支持）===
    assets_fonts = [
        ("simhei.ttf",  "SimHei",  "黑体 - 主标题/重点内容"),
        ("simkai.ttf",  "SimKai",  "楷体 - 引用/强调"),
        ("STFangSong.ttf", "STFangSong", "仿宋 - 正文/报告"),
        ("STHeiti.ttf", "STHeiti", "黑体备选"),
        ("LiSongPro.ttf", "LiSongPro", "隶书 - 可选"),
        # 注意：simfang.ttf 和 fzbsk.ttf 是损坏的 HTML 文件，已排除
    ]
    for fname, name, desc in assets_fonts:
        fp = os.path.join(fonts_dir, fname)
        if os.path.exists(fp):
            fonts["assets"].append((fp, name, desc))

    # === reportlab 回退字体（纯 TTF，排除 TTC）===
    # Helvetica.ttc 是 TTC，reportlab 的 TTFont 不支持 TTC，所以找一个替代
    for fp in [
        "/System/Library/Fonts/Helvetica.ttc",  # TTC，reportlab 不支持
        "/System/Library/Fonts/Arial.ttf",       # 不存在，macOS Arial 是 .ttc
    ]:
        pass  # 都不适合 reportlab

    # macOS 系统纯 TTF 拉丁字体（适合 reportlab/fpdf2）
    system_ttf_candidates = [
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Georgia.ttf",
        "/System/Library/Fonts/Supplemental/Times New Roman.ttf",
        "/System/Library/Fonts/Supplemental/Verdana.ttf",
        "/System/Library/Fonts/Supplemental/Courier New.ttf",
    ]
    for fp in system_ttf_candidates:
        if os.path.exists(fp):
            name = os.path.splitext(os.path.basename(fp))[0]
            fonts["system"].append((fp, name, "系统内置拉丁字体"))

    # === 系统 TTC 字体（用于 R showtext / matplotlib，不适合 reportlab/fpdf2）===
    system_ttc = [
        ("/System/Library/Fonts/STHeiti Light.ttc", "STHeiti",  "macOS 黑体 - CJK"),
        ("/System/Library/Fonts/Hiragino Sans GB.ttc", "Hiragino Sans GB", "macOS 冬青黑体 - CJK"),
        ("/System/Library/Fonts/PingFang SC.ttc", "PingFang SC", "macOS 苹方简体中文"),
        ("/System/Library/Fonts/PingFang TC.ttc", "PingFang TC", "macOS 苹方繁体中文"),
    ]
    for fp, name, desc in system_ttc:
        if os.path.exists(fp):
            fonts["system"].append((fp, name, desc))

    # === reportlab 回退：找一个可用的系统纯 TTF ===
    for fp, name, desc in fonts["system"]:
        if fp.endswith(".ttf"):
            fonts["reportlab_fallback"] = (fp, name)
            break

    return fonts


def _find_chinese_font() -> tuple[str | None, str]:
    """返回 (font_path, font_name)，优先使用 assets/simhei.ttf > 系统 STHeiti.ttc"""
    fonts = _get_available_fonts()
    # 优先 assets
    if fonts["assets"]:
        return fonts["assets"][0][0], fonts["assets"][0][1]
    # 系统 TTC
    for fp, name, desc in fonts["system"]:
        if fp.endswith(".ttc"):
            return fp, name
    # 回退
    if fonts["reportlab_fallback"]:
        return fonts["reportlab_fallback"]
    return None, "Helvetica"


def _save_pdf(md_text: str, base_name: str, workspace_dir: str) -> Path | None:
    """使用 R with Cairo 生成中文 PDF（优先），失败则降级到 reportlab"""
    Path(workspace_dir).mkdir(parents=True, exist_ok=True)
    pdf_path = uniquify_path(Path(workspace_dir) / f"{base_name}.pdf")

    # 先尝试 R Cairo 方案（对 CJK 字体支持最完整）
    r_result = _save_pdf_with_r(md_text, base_name, workspace_dir)
    if r_result:
        print(f"PDF generated with R Cairo: {r_result}")
        return r_result

    # 降级到 reportlab
    return _save_pdf_with_reportlab(md_text, base_name, workspace_dir)


def _save_pptx(md_text: str, base_name: str, workspace_dir: str) -> Path | None:
    """使用 python-pptx 生成 PPTX 报告，支持中文字体、图表嵌入和丰富排版

    Args:
        md_text: Markdown 格式的分析报告文本
        base_name: 输出文件的基础名称（不含扩展名）
        workspace_dir: 目标工作目录路径

    Returns:
        生成的 PPTX 文件路径，失败时返回 None
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        import re as _re

        # Sanitize base_name to prevent path traversal
        safe_base_name = Path(base_name).name.replace("..", "")
        ws_path = Path(workspace_dir).resolve()
        ws_path.mkdir(parents=True, exist_ok=True)
        pptx_path = uniquify_path(ws_path / f"{safe_base_name}.pptx")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Resolve Chinese font path for PPTX
        _pptx_font_name = "SimHei"
        _pptx_body_font = "STFangSong"

        def _set_font(run, font_name, size_pt, bold=False, italic=False, color_rgb=None):
            """Helper to set font properties on a run"""
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.italic = italic
            if color_rgb:
                run.font.color.rgb = color_rgb
            run.font.name = font_name
            try:
                from lxml import etree
                rPr = run._r.get_or_add_rPr()
                ea_font = etree.SubElement(
                    rPr,
                    '{http://schemas.openxmlformats.org/drawingml/2006/main}ea',
                )
                ea_font.set('typeface', font_name)
            except Exception:
                pass

        def _add_text_with_formatting(text_frame, text, font_name, size_pt, bold=False, color_rgb=None, line_spacing=1.2):
            """Add multi-paragraph text with proper formatting to a text frame"""
            lines = text.split('\n')
            first = True
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                if first:
                    p = text_frame.paragraphs[0]
                    first = False
                else:
                    p = text_frame.add_paragraph()

                p.space_after = Pt(4)
                p.space_before = Pt(2)

                # Check if it's a bullet point
                is_bullet = line.startswith('- ') or line.startswith('* ') or line.startswith('• ')
                if is_bullet:
                    line = line.lstrip('-*• ').strip()
                    run = p.add_run()
                    run.text = "• "
                    _set_font(run, font_name, size_pt, bold=True, color_rgb=RGBColor(0, 102, 204))
                    run2 = p.add_run()
                    run2.text = line
                    _set_font(run2, font_name, size_pt, color_rgb=color_rgb or RGBColor(51, 51, 51))
                elif line.startswith('**') and line.endswith('**'):
                    # Bold text
                    run = p.add_run()
                    run.text = line.strip('*')
                    _set_font(run, font_name, size_pt + 1, bold=True, color_rgb=color_rgb or RGBColor(0, 51, 102))
                else:
                    # Clean markdown formatting
                    clean = _re.sub(r'\*\*(.*?)\*\*', r'\1', line)
                    clean = _re.sub(r'\*(.*?)\*', r'\1', clean)
                    clean = _re.sub(r'`(.*?)`', r'\1', clean)
                    clean = _re.sub(r'!\[[^\]]*\]\([^)]*\)', '', clean)
                    clean = _re.sub(r'\[([^\]]*)\]\([^)]*\)', r'\1', clean)
                    if clean.strip():
                        run = p.add_run()
                        run.text = clean
                        _set_font(run, font_name, size_pt, color_rgb=color_rgb or RGBColor(51, 51, 51))

        # Collect image references from workspace/generated
        generated_dir = ws_path / "generated"
        available_images_by_name = {}
        available_images_by_stem = {}
        for img_dir in [generated_dir, ws_path]:
            if img_dir.exists():
                for f in img_dir.iterdir():
                    if f.suffix.lower() in ('.png', '.jpg', '.jpeg', '.svg', '.gif'):
                        available_images_by_name[f.name] = str(f)
                        if f.stem not in available_images_by_name:
                            available_images_by_stem[f.stem] = str(f)

        # === Title Slide ===
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background gradient effect via a colored shape
        bg_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(3.5)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
        bg_shape.line.fill.background()

        # Title text
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11.333), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = base_name.replace("_", " ")
        _set_font(run, _pptx_font_name, 36, bold=True, color_rgb=RGBColor(255, 255, 255))
        p.alignment = PP_ALIGN.CENTER

        # Subtitle with date
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = f"DeepAnalyze 智能分析报告 | {datetime.now().strftime('%Y年%m月%d日')}"
        _set_font(run2, _pptx_body_font, 16, color_rgb=RGBColor(200, 220, 255))

        # === Content Slides ===
        # Parse markdown into sections by headers
        # Split by any level header
        raw_sections = _re.split(r'\n(?=#{1,3}\s+)', md_text)

        for section_text in raw_sections:
            if not section_text.strip():
                continue

            lines = section_text.strip().split('\n')
            # Extract title from first line
            title_line = lines[0].strip()
            title = title_line.lstrip('#').strip() if title_line.startswith('#') else title_line
            if not title:
                title = "分析内容"
            body_text = '\n'.join(lines[1:]).strip() if len(lines) > 1 else ""

            # Extract image references from markdown
            img_refs = _re.findall(r'!\[[^\]]*\]\(([^)]*)\)', body_text)

            # Find actual image files
            embedded_images = []
            for ref in img_refs:
                img_name = os.path.basename(ref)
                img_stem = Path(img_name).stem
                img_path = None
                if img_name in available_images_by_name:
                    img_path = available_images_by_name[img_name]
                elif img_stem in available_images_by_stem:
                    img_path = available_images_by_stem[img_stem]
                else:
                    candidate = (ws_path / ref).resolve()
                    # Ensure resolved path stays within workspace
                    if str(candidate).startswith(str(ws_path)) and candidate.exists():
                        img_path = str(candidate)
                    elif generated_dir.exists():
                        candidate = (generated_dir / img_name).resolve()
                        if str(candidate).startswith(str(ws_path)) and candidate.exists():
                            img_path = str(candidate)
                # Verify image is not empty (> 1KB minimum)
                MIN_IMAGE_SIZE_BYTES = 1024
                if img_path and os.path.exists(img_path) and os.path.getsize(img_path) > MIN_IMAGE_SIZE_BYTES:
                    embedded_images.append(img_path)

            # Skip sections with no content
            clean_body = _re.sub(r'!\[[^\]]*\]\([^)]*\)', '', body_text)
            clean_body = _re.sub(r'\[([^\]]*)\]\([^)]*\)', r'\1', clean_body)
            clean_body = clean_body.strip()
            if not clean_body and not embedded_images:
                continue

            # Create slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Title bar background
            title_bg = slide.shapes.add_shape(
                1, Inches(0), Inches(0), prs.slide_width, Inches(1.1)
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(0, 51, 102)
            title_bg.line.fill.background()

            # Title textbox
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title[:80]
            _set_font(run, _pptx_font_name, 24, bold=True, color_rgb=RGBColor(255, 255, 255))

            # Content area
            if embedded_images:
                # Layout with image on right, text on left
                text_width = Inches(6.0)
                img_left = Inches(6.8)
                img_width = Inches(5.8)
                img_top = Inches(1.3)

                # Text body on the left - with multi-paragraph support
                if clean_body:
                    MAX_CHARS = 2000
                    if len(clean_body) > MAX_CHARS:
                        clean_body = clean_body[:MAX_CHARS] + "\n...(内容已截断)"
                    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), text_width, Inches(5.8))
                    tf2 = txBox2.text_frame
                    tf2.word_wrap = True
                    _add_text_with_formatting(tf2, clean_body, _pptx_body_font, 12)

                # Embed images on the right (up to 2 per slide)
                for idx, img_path in enumerate(embedded_images[:2]):
                    try:
                        y_offset = img_top + Inches(idx * 3.0)
                        slide.shapes.add_picture(
                            img_path, img_left, y_offset, img_width, Inches(2.8)
                        )
                    except Exception as img_err:
                        print(f"PPTX 图片嵌入失败: {img_err}")
            else:
                # Full-width text layout with multi-paragraph support
                if clean_body:
                    MAX_PPTX_SLIDE_CHARS = 2500
                    if len(clean_body) > MAX_PPTX_SLIDE_CHARS:
                        # Split into multiple slides
                        chunks = [clean_body[i:i+MAX_PPTX_SLIDE_CHARS] for i in range(0, len(clean_body), MAX_PPTX_SLIDE_CHARS)]
                        for chunk_idx, chunk in enumerate(chunks):
                            if chunk_idx > 0:
                                # Additional slide for overflow content
                                slide = prs.slides.add_slide(prs.slide_layouts[6])
                                title_bg2 = slide.shapes.add_shape(
                                    1, Inches(0), Inches(0), prs.slide_width, Inches(1.1)
                                )
                                title_bg2.fill.solid()
                                title_bg2.fill.fore_color.rgb = RGBColor(0, 51, 102)
                                title_bg2.line.fill.background()
                                txB = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.8))
                                tfB = txB.text_frame
                                pB = tfB.paragraphs[0]
                                rB = pB.add_run()
                                rB.text = f"{title[:60]}（续）"
                                _set_font(rB, _pptx_font_name, 24, bold=True, color_rgb=RGBColor(255, 255, 255))

                            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.8))
                            tf2 = txBox2.text_frame
                            tf2.word_wrap = True
                            _add_text_with_formatting(tf2, chunk, _pptx_body_font, 13)
                    else:
                        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.8))
                        tf2 = txBox2.text_frame
                        tf2.word_wrap = True
                        _add_text_with_formatting(tf2, clean_body, _pptx_body_font, 13)

        prs.save(str(pptx_path))

        if pptx_path.exists() and pptx_path.stat().st_size > 0:
            print(f"PPTX 生成成功: {pptx_path.name}")
            return pptx_path
        return None
    except ImportError:
        print("python-pptx 未安装，跳过 PPTX 生成。请运行: pip install python-pptx")
        return None
    except Exception as e:
        print(f"PPTX 生成失败: {e}")
        traceback.print_exc()
        return None


def _save_pdf_with_r(md_text: str, base_name: str, workspace_dir: str) -> Path | None:
    """使用 R script + Cairo/grDevices 生成中文 PDF（支持所有 CJK 字体）"""
    import subprocess

    Path(workspace_dir).mkdir(parents=True, exist_ok=True)
    pdf_path = uniquify_path(Path(workspace_dir) / f"{base_name}.pdf")

    font_path, font_name = _find_chinese_font()
    if not font_path:
        return None

    r_script_path = os.path.join(workspace_dir, f"_pdf_gen_{os.getpid()}.R")
    r_font_path = font_path.replace("\\", "\\\\").replace("'", "\\'")
    clean_text = re.sub(r"\\newpage", "", md_text)
    r_pdf_path = str(pdf_path).replace("\\", "\\\\").replace("'", "\\'")
    r_text = _escape_r_string(clean_text)

    r_script = f"""
options(width=120)
Sys.setenv(LANG="en_US.UTF-8")
suppressPackageStartupMessages(library(showtext))
suppressPackageStartupMessages(library(sysfonts))

font_add("{font_name}", "{r_font_path}")
showtext_auto()

pdf("{r_pdf_path}", family="{font_name}", width=8.27, height=11.69)
par(family="{font_name}", mar=c(2.5, 2.5, 2, 1), oma=c(0, 0, 0, 0))

txt <- "{r_text}"
lines <- strsplit(txt, "\\n")[[1]]
if (length(lines) == 0) lines <- c(txt)

plot.new()
y <- 0.98
line_height <- 0.03

for (line in lines) {{
    raw <- trimws(line)
    if (nchar(raw) == 0) {{
        y <- y - line_height
    }} else {{
        cex_val <- 0.95
        font_val <- 1
        if (grepl("^###\\\\s+", raw, perl=TRUE)) {{
            raw <- sub("^###\\\\s+", "", raw, perl=TRUE)
            cex_val <- 1.15
            font_val <- 2
        }} else if (grepl("^##\\\\s+", raw, perl=TRUE)) {{
            raw <- sub("^##\\\\s+", "", raw, perl=TRUE)
            cex_val <- 1.2
            font_val <- 2
        }} else if (grepl("^#\\\\s+", raw, perl=TRUE)) {{
            raw <- sub("^#\\\\s+", "", raw, perl=TRUE)
            cex_val <- 1.3
            font_val <- 2
        }} else if (grepl("^[-*]\\\\s+", raw, perl=TRUE)) {{
            raw <- paste0("• ", sub("^[-*]\\\\s+", "", raw, perl=TRUE))
        }}

        text(0.02, y, labels=raw, adj=c(0, 1), cex=cex_val, family="{font_name}", font=font_val)
        y <- y - line_height
    }}

    if (y <= 0.04) {{
        plot.new()
        y <- 0.98
    }}
}}

dev.off()
cat("PDF_R_OK\\n")
"""
    try:
        with open(r_script_path, "w", encoding="utf-8") as f:
            f.write(r_script)

        result = subprocess.run(
            ["Rscript", "--vanilla", "--quiet", r_script_path],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=60,
        )

        try:
            os.remove(r_script_path)
        except Exception:
            pass

        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
            return pdf_path

        print(f"R PDF generation failed stdout: {result.stdout}, stderr: {result.stderr}")
        return None
    except Exception as e:
        print(f"R PDF generation error: {e}")
        try:
            os.remove(r_script_path)
        except Exception:
            pass
        return None


def _escape_r_string(s: str) -> str:
    """转义字符串用于嵌入 R 脚本"""
    return s.replace("\\", "\\\\").replace('"', '\\"').replace("\n", "\\n").replace("$", "\\$")


from typing import Optional


def _render_md_to_html(md_text: str, title: Optional[str] = None) -> str:
    """简化为占位实现（仅供未来 PDF 渲染使用）。当前仅生成 MD。"""
    doc_title = (title or "Report").strip() or "Report"
    safe = (md_text or "").replace("<", "&lt;").replace(">", "&gt;")
    return f"<html><head><meta charset='utf-8'><title>{doc_title}</title></head><body><pre>{safe}</pre></body></html>"


def _save_pdf_from_md(html_text: str, base_name: str) -> Path:
    """TODO: 服务端 PDF 渲染未实现。"""
    raise NotImplementedError("TODO: implement server-side PDF rendering")


def _save_pdf_with_chromium(html_text: str, base_name: str) -> Path:
    """TODO: 使用 Chromium 渲染 PDF（暂不实现）。"""
    raise NotImplementedError("TODO: chromium-based PDF rendering")


def _save_pdf_from_text(text: str, base_name: str) -> Path:
    """TODO: 纯文本 PDF 渲染（暂不实现）。"""
    raise NotImplementedError("TODO: text-based PDF rendering")


def _save_docx(md_text: str, base_name: str, workspace_dir: str) -> Path | None:
    """
    使用模块化方式生成中文 DOCX 文件

    优化要点：
    1. 使用 docx_utils 模块处理中文编码问题
    2. 正确设置中文字体（STFangSong）
    3. 确保使用 UTF-8 编码
    4. 避免乱码问题

    关键修复：
    - 使用 create_docx_document() 创建配置好字体的文档
    - 使用 clean_md_text_for_docx() 清理 Markdown 标记
    - 使用 extract_markdown_blocks() 解析内容结构
    - 确保所有中文 run 都设置正确的 eastAsia 字体
    """
    Path(workspace_dir).mkdir(parents=True, exist_ok=True)
    docx_path = uniquify_path(Path(workspace_dir) / f"{base_name}.docx")

    try:
        print(f"开始生成 DOCX: {docx_path}")

        # 使用 docx_utils 模块生成 DOCX
        success = generate_docx_module(md_text, str(docx_path), title=base_name)

        if success:
            print(f"DOCX 生成成功: {docx_path}")
            return docx_path
        else:
            print(f"DOCX 生成失败")
            return None

    except Exception as e:
        error_msg = f"DOCX 生成失败: {e}"
        print(error_msg)
        traceback.print_exc()
        return None

def _save_pdf_with_reportlab(md_text: str, base_name: str, workspace_dir: str) -> Path | None:
    """
    使用 reportlab + simhei.ttf 生成中文 PDF（降级方案）

    优化要点：
    1. 使用模块化的字体注册函数 register_chinese_fonts()
    2. 使用 get_chinese_style() 获取预定义样式
    3. 使用 extract_markdown_sections() 提取结构化内容
    4. 使用 clean_md_text() 清理 Markdown 标记
    5. 完善的错误处理和日志记录
    """
    Path(workspace_dir).mkdir(parents=True, exist_ok=True)
    pdf_path = uniquify_path(Path(workspace_dir) / f"{base_name}.pdf")

    try:
        print(f"开始生成 PDF (ReportLab): {pdf_path}")

        # 1. 注册中文字体（使用模块化函数，只注册一次）
        registered = register_chinese_fonts(force=False)
        if not registered:
            print("警告：没有成功注册中文字体，PDF 可能无法正常显示中文")
        else:
            print(f"已注册字体: {list(registered.keys())}")

        # 2. 提取 Markdown 内容块
        clean_text = clean_md_text(md_text)
        blocks = extract_markdown_sections(clean_text)
        print(f"提取到 {len(blocks)} 个内容块")

        # 3. 初始化 PDF 文档
        doc = SimpleDocTemplate(
            str(pdf_path),
            pagesize=A4,
            rightMargin=50,
            leftMargin=50,
            topMargin=50,
            bottomMargin=50,
        )

        # 4. 准备故事（PDF 内容）
        story = []

        # 5. 处理每个内容块
        for block in blocks:
            block_type = block["type"]
            content = block["content"]

            if block_type == "heading1":
                style = get_chinese_style("heading1")
                story.append(Paragraph(content, style))
                story.append(Spacer(1, 12))

            elif block_type == "heading2":
                style = get_chinese_style("heading2")
                story.append(Paragraph(content, style))
                story.append(Spacer(1, 8))

            elif block_type == "heading3":
                style = get_chinese_style("heading3")
                story.append(Paragraph(content, style))
                story.append(Spacer(1, 6))

            elif block_type == "paragraph":
                style = get_chinese_style("normal")
                story.append(Paragraph(content, style))
                story.append(Spacer(1, 8))

            elif block_type == "table":
                # 使用 reportlab Table 渲染
                try:
                    from reportlab.platypus import Table, TableStyle
                    from reportlab.lib import colors as rl_colors

                    table_data = content
                    headers = table_data.get("headers", [])
                    rows = table_data.get("rows", [])
                    font_name = get_chinese_font_name()

                    all_rows = [headers] + rows
                    max_cols = max(len(r) for r in all_rows) if all_rows else 0
                    normalized = []
                    for row in all_rows:
                        padded = list(row) + [''] * (max_cols - len(row))
                        normalized.append(padded[:max_cols])

                    if normalized and max_cols > 0:
                        available_width = A4[0] - 100
                        col_width = available_width / max_cols
                        table = Table(normalized, colWidths=[col_width] * max_cols)
                        table.setStyle(TableStyle([
                            ('FONTNAME', (0, 0), (-1, -1), font_name),
                            ('FONTSIZE', (0, 0), (-1, -1), 9),
                            ('FONTNAME', (0, 0), (-1, 0), font_name),
                            ('FONTSIZE', (0, 0), (-1, 0), 10),
                            ('BACKGROUND', (0, 0), (-1, 0), rl_colors.HexColor('#003366')),
                            ('TEXTCOLOR', (0, 0), (-1, 0), rl_colors.white),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                            ('GRID', (0, 0), (-1, -1), 0.5, rl_colors.grey),
                            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [rl_colors.white, rl_colors.HexColor('#F0F4F8')]),
                            ('TOPPADDING', (0, 0), (-1, -1), 4),
                            ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                        ]))
                        story.append(table)
                        story.append(Spacer(1, 12))
                except Exception as table_err:
                    print(f"表格渲染失败: {table_err}")
                    style = get_chinese_style("normal")
                    raw_text = block.get("raw", str(content))
                    story.append(Paragraph(raw_text.replace('\n', '<br/>'), style))

            elif block_type == "list":
                style = get_chinese_style("list")
                for item in content:
                    para_text = f"• {item}"
                    story.append(Paragraph(para_text, style))
                    story.append(Spacer(1, 4))
                story.append(Spacer(1, 8))

        # 6. 构建 PDF
        doc.build(story)

        # 7. 验证生成结果
        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
            file_size = os.path.getsize(pdf_path) / 1024
            print(f"PDF 生成成功: {pdf_path.name} ({file_size:.1f} KB)")
            return pdf_path
        else:
            print(f"PDF 生成失败: 文件不存在或为空")
            return None

    except Exception as e:
        error_msg = f"PDF 生成失败: {e}"
        print(error_msg)
        traceback.print_exc()
        return None

@app.post("/export/report")
async def export_report(body: dict = Body(...)):
    """
    接收全部聊天历史（messages: [{role, content}...]），抽取 <Analyze>..~ <Answer>..
    生成 Markdown, PDF (ReportLab) 和 DOCX 文件。
    """
    try:
        messages = body.get("messages", [])
        title = (body.get("title") or "").strip()
        session_id = body.get("session_id", "default")
        username = body.get("username", "default")
        analysis_language = normalize_analysis_language(body.get("analysis_language", "zh-CN"))
        workspace_dir = get_session_workspace(session_id, username)

        if not isinstance(messages, list):
            raise HTTPException(status_code=400, detail="messages must be a list")

        md_text = _extract_sections_from_messages(messages)
        if not md_text:
            if analysis_language == "en":
                md_text = "(No <Analyze>/<Understand>/<Code>/<Execute>/<Answer> sections found.)"
            else:
                md_text = "（未找到可导出的 <Analyze>/<Understand>/<Code>/<Execute>/<Answer> 内容。）"

        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        default_title = "Report" if analysis_language == "en" else "报告"
        safe_title = re.sub(r"[^\w\-_.]+", "_", title) if title else default_title
        base_name = f"{safe_title}_{ts}"

        export_dir = os.path.join(workspace_dir, "generated")
        os.makedirs(export_dir, exist_ok=True)

        md_path = _save_md(md_text, base_name, export_dir)

        # 根据请求的报告类型生成文件
        report_types = body.get("report_types", ["pdf"])
        if not report_types:
            report_types = ["pdf"]

        docx_path = _save_docx(md_text, base_name, export_dir) if "docx" in report_types else None
        pdf_path = _save_pdf(md_text, base_name, export_dir) if "pdf" in report_types else None
        pptx_path = _save_pptx(md_text, base_name, export_dir) if "pptx" in report_types else None

        result = {
            "message": "exported",
            "analysis_language": analysis_language,
            "md": md_path.name,
            "pdf": pdf_path.name if pdf_path else None,
            "docx": docx_path.name if docx_path else None,
            "pptx": pptx_path.name if pptx_path else None,
            "download_urls": {
                "md": build_download_url(f"{username}/{session_id}/generated/{md_path.name}"),
                "pdf": build_download_url(f"{username}/{session_id}/generated/{pdf_path.name}") if pdf_path else None,
                "docx": build_download_url(f"{username}/{session_id}/generated/{docx_path.name}") if docx_path else None,
                "pptx": build_download_url(f"{username}/{session_id}/generated/{pptx_path.name}") if pptx_path else None,
            },
        }
        return JSONResponse(result)
    except HTTPException:
        raise
    except Exception as e:
        print(f"Export report error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/users/list")
async def list_users():
    """获取已注册用户列表（仅返回用户名）"""
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        ensure_builtin_superuser(cursor)
        conn.commit()
        cursor.execute("SELECT username FROM users ORDER BY username ASC")
        rows = cursor.fetchall()
        conn.close()
        return {"users": [row["username"] for row in rows]}
    except Exception as e:
        print(f"List users error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

def _get_dir_size(dir_path: str) -> int:
    """计算目录的总大小（字节）"""
    total_size = 0
    try:
        for root, dirs, files in os.walk(dir_path):
            for fname in files:
                fpath = os.path.join(root, fname)
                try:
                    total_size += os.path.getsize(fpath)
                except OSError:
                    pass
    except Exception:
        pass
    return total_size


def _format_file_size(size_bytes: int) -> str:
    """格式化文件大小为人类可读的格式"""
    if size_bytes < 1024:
        return f"{size_bytes}B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f}K"
    elif size_bytes < 1024 * 1024 * 1024:
        return f"{size_bytes / (1024 * 1024):.1f}M"
    else:
        return f"{size_bytes / (1024 * 1024 * 1024):.1f}G"


@app.post("/api/projects/save")
async def save_project(
    username: str = Form(...),
    session_id: str = Form(...),
    name: str = Form(...),
    messages: str = Form(...),
    files_data: str = Form("{}"),
    side_tasks: str = Form("[]")
):
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()

        # Check if project name already exists for this user to support "overwrite" logic
        cursor.execute("SELECT id, session_id FROM projects WHERE username = ? AND name = ?", (username, name))
        existing = cursor.fetchone()

        # Build workspace snapshot: list all files in the user's workspace directory
        workspace_dir = get_session_workspace(session_id, username)
        files_snapshot = []

        if existing:
            project_id = existing["id"]
            # Create project directory to store file copies
            project_dir = os.path.join(PROJECTS_BASE_DIR, str(project_id))
            os.makedirs(project_dir, exist_ok=True)
        else:
            # For new project, use a temporary directory name first
            project_id = None
            project_dir = None

        if os.path.exists(workspace_dir):
            for root, dirs, files in os.walk(workspace_dir):
                for fname in files:
                    if fname.startswith("."):
                        continue
                    fpath = os.path.join(root, fname)
                    rel_path = os.path.relpath(fpath, workspace_dir)
                    files_snapshot.append({
                        "name": fname,
                        "path": rel_path,
                        "size": os.path.getsize(fpath)
                    })

        files_json = json.dumps(files_snapshot)

        if existing:
            cursor.execute(
                "UPDATE projects SET session_id = ?, messages = ?, files_data = ?, side_tasks = ?, created_at = CURRENT_TIMESTAMP WHERE id = ?",
                (session_id, messages, files_json, side_tasks, project_id)
            )
        else:
            cursor.execute(
                "INSERT INTO projects (username, session_id, name, messages, files_data, side_tasks) VALUES (?, ?, ?, ?, ?, ?)",
                (username, session_id, name, messages, files_json, side_tasks)
            )
            project_id = cursor.lastrowid
            # Now create project directory with actual ID
            project_dir = os.path.join(PROJECTS_BASE_DIR, str(project_id))
            os.makedirs(project_dir, exist_ok=True)

        # Copy files to project directory
        if project_dir and os.path.exists(workspace_dir):
            for root, dirs, files in os.walk(workspace_dir):
                for fname in files:
                    if fname.startswith("."):
                        continue
                    fpath = os.path.join(root, fname)
                    rel_path = os.path.relpath(fpath, workspace_dir)
                    try:
                        project_file_dir = os.path.join(project_dir, os.path.dirname(rel_path))
                        os.makedirs(project_file_dir, exist_ok=True)
                        project_file_path = os.path.join(project_file_dir, fname)
                        shutil.copy2(fpath, project_file_path)
                    except Exception as copy_err:
                        print(f"Warning: Failed to copy file {rel_path} to project directory: {copy_err}")

        conn.commit()
        conn.close()

        # Calculate storage size
        storage_size = _get_dir_size(project_dir) if project_dir else 0
        storage_display = _format_file_size(storage_size)

        return {"message": "Project saved successfully", "project_id": project_id, "storage_size": storage_display}
    except Exception as e:
        print(f"Save project error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/projects/list")
async def list_projects(username: str = Query(...)):
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute(
            "SELECT id, name, session_id, created_at FROM projects WHERE username = ? ORDER BY created_at DESC",
            (username,)
        )
        rows = cursor.fetchall()
        projects = []
        for row in rows:
            proj = dict(row)
            # Calculate storage size for each project
            project_dir = os.path.join(PROJECTS_BASE_DIR, str(proj["id"]))
            storage_size = _get_dir_size(project_dir)
            proj["storage_size"] = _format_file_size(storage_size)
            projects.append(proj)
        conn.close()
        return {"projects": projects}
    except Exception as e:
        print(f"List projects error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/projects/check-name")
async def check_project_name(username: str = Query(...), name: str = Query(...)):
    """检查项目名称是否已存在"""
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute("SELECT id FROM projects WHERE username = ? AND name = ?", (username, name))
        row = cursor.fetchone()
        conn.close()
        return {"exists": row is not None, "project_id": row["id"] if row else None}
    except Exception as e:
        print(f"Check project name error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/projects/restore-files")
async def restore_project_files(project_id: int = Query(...)):
    """获取指定项目的文件列表及其下载链接，供前端恢复文件"""
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute("SELECT session_id, username, files_data FROM projects WHERE id = ?", (project_id,))
        row = cursor.fetchone()
        conn.close()
        if not row:
            raise HTTPException(status_code=404, detail="Project not found")

        files_data = json.loads(row["files_data"]) if row["files_data"] else []
        result_files = []
        for f in files_data:
            rel_path = f.get("path", "")
            # 构建项目文件的下载 URL（从项目目录读取）
            download_url = f"/api/projects/file/{project_id}/{quote(rel_path)}"
            result_files.append({
                "name": f.get("name", ""),
                "path": rel_path,
                "size": f.get("size", 0),
                "download_url": download_url,
            })
        return {"files": result_files}
    except HTTPException:
        raise
    except Exception as e:
        print(f"Restore project files error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/projects/file/{project_id}/{file_path}")
async def get_project_file(project_id: int, file_path: str):
    """从项目目录中获取文件内容"""
    try:
        # 构建项目文件路径
        project_dir = os.path.join(PROJECTS_BASE_DIR, str(project_id))
        if not os.path.exists(project_dir):
            raise HTTPException(status_code=404, detail="Project directory not found")

        # 清理文件路径，防止目录穿越攻击
        safe_path = os.path.normpath(file_path)
        file_full_path = os.path.join(project_dir, safe_path)

        # 确保文件在项目目录内
        if not os.path.abspath(file_full_path).startswith(os.path.abspath(project_dir) + os.sep):
            raise HTTPException(status_code=403, detail="Access denied")

        if not os.path.exists(file_full_path):
            raise HTTPException(status_code=404, detail="File not found")

        # 返回文件
        return FileResponse(file_full_path)
    except HTTPException:
        raise
    except Exception as e:
        print(f"Get project file error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/projects/restore-to-workspace")
async def restore_to_workspace(project_id: int = Query(...), session_id: str = Query(...), username: str = Query(...)):
    """将项目文件实质性恢复到指定的 workspace 目录中"""
    try:
        project_dir = os.path.join(PROJECTS_BASE_DIR, str(project_id))
        if not os.path.exists(project_dir):
            raise HTTPException(status_code=404, detail="Project directory not found")

        workspace_dir = get_session_workspace(session_id, username)

        # 恢复前清空工作区，确保项目独立性
        if os.path.exists(workspace_dir):
            shutil.rmtree(workspace_dir)
        os.makedirs(workspace_dir, exist_ok=True)

        # 复制项目目录下的所有文件到工作区
        for root, dirs, files in os.walk(project_dir):
            for fname in files:
                if fname.startswith("."):
                    continue
                fpath = os.path.join(root, fname)
                rel_path = os.path.relpath(fpath, project_dir)
                target_path = os.path.join(workspace_dir, rel_path)

                # 确保目标父目录存在
                os.makedirs(os.path.dirname(target_path), exist_ok=True)

                try:
                    shutil.copy2(fpath, target_path)
                except Exception as copy_err:
                    print(f"Warning: Failed to restore file {rel_path} to workspace: {copy_err}")

        return {"message": "Project files restored successfully"}
    except HTTPException:
        raise
    except Exception as e:
        print(f"Restore to workspace error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/projects/load")
async def load_project(project_id: int = Query(...)):
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute("SELECT session_id, messages, files_data, side_tasks FROM projects WHERE id = ?", (project_id,))
        row = cursor.fetchone()
        conn.close()
        if not row:
            raise HTTPException(status_code=404, detail="Project not found")
        return {
            "session_id": row["session_id"],
            "messages": json.loads(row["messages"]),
            "files_data": json.loads(row["files_data"]) if row["files_data"] else [],
            "side_tasks": json.loads(row["side_tasks"]) if "side_tasks" in row.keys() and row["side_tasks"] else []
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"Load project error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/projects/delete")
async def delete_project(project_id: int = Query(...), username: str = Query(...)):
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        # Get session_id first to delete files
        cursor.execute("SELECT session_id FROM projects WHERE id = ? AND username = ?", (project_id, username))
        row = cursor.fetchone()
        if not row:
            conn.close()
            raise HTTPException(status_code=404, detail="Project not found")

        session_id = row["session_id"]

        # Delete from DB
        cursor.execute("DELETE FROM projects WHERE id = ?", (project_id,))
        conn.commit()
        conn.close()

        # Delete workspace files
        workspace_dir = get_session_workspace(session_id, username)
        if os.path.exists(workspace_dir):
            shutil.rmtree(workspace_dir)

        return {"message": "Project deleted successfully"}
    except Exception as e:
        print(f"Delete project error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


# ========== 雨途斩棘录 API ==========
@app.get("/api/yutu/html")
async def get_yutu_html_api():
    """获取雨途斩棘录HTML内容（所有人都可以查看）"""
    try:
        html = get_yutu_html()
        return {"success": True, "html": html}
    except Exception as e:
        print(f"Get yutu HTML error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/yutu/add")
async def add_yutu_record(body: dict = Body(...), username: str = Query("default")):
    """
    添加错误记录和解决方案（仅超级用户）
    """
    try:
        # 检查是否为超级用户
        if not is_builtin_superuser(username):
            raise HTTPException(status_code=403, detail="Only superuser can add records")

        error_type = body.get("error_type", "Unknown")
        error_message = body.get("error_message", "")
        error_context = body.get("error_context")
        solution = body.get("solution", "")
        solution_code = body.get("solution_code")
        confidence = body.get("confidence", 0.0)

        success = add_error_solution(
            error_type=error_type,
            error_message=error_message,
            error_context=error_context,
            solution=solution,
            solution_code=solution_code,
            confidence=confidence,
            created_by=username
        )

        if success:
            return {"success": True, "message": "Record added successfully"}
        else:
            raise HTTPException(status_code=500, detail="Failed to add record")
    except HTTPException:
        raise
    except Exception as e:
        print(f"Add yutu record error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/yutu/update")
async def update_yutu_record(body: dict = Body(...), username: str = Query("default")):
    """
    更新错误记录（仅超级用户）
    """
    try:
        # 检查是否为超级用户
        if not is_builtin_superuser(username):
            raise HTTPException(status_code=403, detail="Only superuser can update records")

        error_hash = body.get("error_hash")
        solution = body.get("solution")
        solution_code = body.get("solution_code")
        confidence = body.get("confidence", 0.0)

        if not error_hash:
            raise HTTPException(status_code=400, detail="error_hash is required")

        success = update_error_solution(
            error_hash=error_hash,
            solution=solution,
            solution_code=solution_code,
            confidence=confidence,
            updated_by=username
        )

        if success:
            return {"success": True, "message": "Record updated successfully"}
        else:
            raise HTTPException(status_code=500, detail="Failed to update record")
    except HTTPException:
        raise
    except Exception as e:
        print(f"Update yutu record error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/yutu/delete")
async def delete_yutu_record(body: dict = Body(...), username: str = Query("default")):
    """
    删除错误记录（仅超级用户）
    """
    try:
        # 检查是否为超级用户
        if not is_builtin_superuser(username):
            raise HTTPException(status_code=403, detail="Only superuser can delete records")

        error_hash = body.get("error_hash")

        if not error_hash:
            raise HTTPException(status_code=400, detail="error_hash is required")

        success = delete_error(error_hash)

        if success:
            return {"success": True, "message": "Record deleted successfully"}
        else:
            raise HTTPException(status_code=500, detail="Failed to delete record")
    except HTTPException:
        raise
    except Exception as e:
        print(f"Delete yutu record error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/yutu/search")
async def search_yutu_records(body: dict = Body(...)):
    """
    搜索错误记录（所有人都可以搜索）
    """
    try:
        keywords = body.get("keywords", [])
        error_type = body.get("error_type")
        page = body.get("page", 1)
        page_size = body.get("page_size", 20)

        results = search_errors(
            keywords=keywords,
            error_type=error_type,
            page=page,
            page_size=page_size
        )

        return {"success": True, "data": results}
    except Exception as e:
        print(f"Search yutu records error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/yutu/init")
async def init_yutu_api():
    """初始化雨途斩棘录（超级用户专用）"""
    try:
        init_yutu_if_needed()
        return {"success": True, "message": "Yutu initialized successfully"}
    except Exception as e:
        print(f"Init yutu error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/yutu/organize")
async def organize_yutu_api(body: dict, username: str = ""):
    """整理雨途斩棘录 - 使用VLLM AI重新组织所有记录（超级用户专用）"""
    print(f"[DEBUG] organize_yutu_api called with username={username}, body keys={list(body.keys()) if body else 'None'}")
    # 验证超级用户
    if not is_builtin_superuser(username):
        raise HTTPException(status_code=403, detail="只有超级用户可以整理笔记")

    records = body.get("records", [])
    if not records or len(records) == 0:
        return {"success": False, "detail": "没有记录可整理", "updated_count": 0, "records": []}

    try:
        # 1. 导出原始笔记到本地文件
        import json
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        export_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "yutu_backups")
        os.makedirs(export_dir, exist_ok=True)
        backup_file = os.path.join(export_dir, f"yutu_backup_{timestamp}.json")
        with open(backup_file, "w", encoding="utf-8") as f:
            json.dump(records, f, ensure_ascii=False, indent=2)

        # 2. 使用VLLM进行整理
        prompt = f"""你是雨途斩棘录的智能整理助手。请分析以下所有错误记录，将每条记录的solution字段重新整理，使其：
1. 更清晰地描述错误场景
2. 提供更具体可执行的解决方案
3. 添加"关键要点"总结

请为每条记录生成改进后的solution，保持原有error_type和error_message不变。

原始记录：
{json.dumps(records, ensure_ascii=False, indent=2)}

请直接返回JSON数组格式的整理结果，每条记录包含：error_hash, error_type, error_message, improved_solution（改进后的solution）。"""

        try:
            response = client.chat.completions.create(
                model=MODEL_PATH,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=4000
            )
            improved_content = response.choices[0].message.content
            print(f"VLLM返回内容长度: {len(improved_content)}")

            # 解析改进内容 - 尝试多种方式
            import re
            improved_records = None

            # 方式1: 尝试提取markdown代码块中的JSON
            json_match = re.search(r'```(?:json)?\s*(\[.*?\])\s*```', improved_content, re.DOTALL)
            if json_match:
                try:
                    improved_records = json.loads(json_match.group(1))
                    print(f"方式1成功: 从markdown解析")
                except:
                    pass

            # 方式2: 直接查找JSON数组
            if not improved_records:
                json_match = re.search(r'\[.*\]', improved_content, re.DOTALL)
                if json_match:
                    try:
                        improved_records = json.loads(json_match.group())
                        print(f"方式2成功: 直接解析")
                    except:
                        pass

            # 方式3: 简单清理后解析
            if not improved_records:
                try:
                    cleaned = improved_content.strip()
                    if cleaned.startswith("```"):
                        cleaned = re.sub(r'^```json?', '', cleaned)
                        cleaned = re.sub(r'```$', '', cleaned)
                    improved_records = json.loads(cleaned.strip())
                    print(f"方式3成功: 清理后解析")
                except Exception as e:
                    print(f"方式3失败: {e}")

            if not improved_records or not isinstance(improved_records, list):
                # 如果无法解析，返回原始记录
                print("无法解析VLLM返回的JSON，返回原始记录")
                improved_records = records

            # 保存整理后的预览版本到临时文件
            preview_file = os.path.join(export_dir, f"yutu_preview_{timestamp}.json")
            with open(preview_file, "w", encoding="utf-8") as f:
                json.dump(improved_records, f, ensure_ascii=False, indent=2)

            # 返回预览数据给前端
            return {
                "success": True,
                "updated_count": len(improved_records),
                "records": improved_records,
                "backup_file": backup_file,
                "preview_file": preview_file,
                "timestamp": timestamp
            }
        except Exception as vllm_err:
            print(f"VLLM调用失败: {vllm_err}")
            # VLLM失败时使用本地简单整理
            from yutu_zhanyilu import reorganize_all_records
            updated_count = reorganize_all_records(records)
            return {
                "success": True,
                "updated_count": updated_count,
                "records": records,
                "backup_file": backup_file,
                "note": "VLLM不可用，使用本地整理"
            }

    except Exception as e:
        print(f"Organize yutu error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/yutu/backup/list")
async def list_yutu_backups(username: str = Query("default")):
    """列出所有备份文件"""
    if not is_builtin_superuser(username):
        raise HTTPException(status_code=403, detail="只有超级用户可以查看备份列表")

    try:
        from yutu_zhanyilu import BACKUP_DIR
        # 支持列出所有 .json 备份文件
        files = [f for f in os.listdir(BACKUP_DIR) if f.endswith(".json")]
        files.sort(reverse=True) # 默认排序
        return {"success": True, "backups": files}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/yutu/backup/create")
async def create_yutu_backup(data: dict = Body(...), username: str = Query("default")):
    """手动创建雨途斩棘录备份"""
    if not is_builtin_superuser(username):
        raise HTTPException(status_code=403, detail="只有超级用户可以创建备份")

    custom_name = data.get("filename")
    try:
        from yutu_zhanyilu import backup_to_json
        backup_file = backup_to_json(custom_name)
        if backup_file:
            return {"success": True, "file": os.path.basename(backup_file), "path": backup_file}
        return {"success": False, "detail": "备份失败"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/api/yutu/backup/delete")
async def delete_yutu_backup(filename: str = Query(...), username: str = Query("default")):
    """删除备份文件"""
    if not is_builtin_superuser(username):
        raise HTTPException(status_code=403, detail="只有超级用户可以删除备份")

    try:
        from yutu_zhanyilu import delete_backup
        if delete_backup(filename):
            return {"success": True, "message": f"备份文件 {filename} 已删除"}
        return {"success": False, "detail": "删除失败，文件可能不存在"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/yutu/backup/restore")
async def restore_yutu_backup(data: dict, username: str = Query("default")):
    """从备份恢复"""
    if not is_builtin_superuser(username):
        raise HTTPException(status_code=403, detail="只有超级用户可以恢复备份")

    filename = data.get("filename")
    mode = data.get("mode", "append") # append or overwrite

    if not filename:
        raise HTTPException(status_code=400, detail="未指定备份文件名")

    try:
        from yutu_zhanyilu import restore_from_json, BACKUP_DIR
        file_path = os.path.join(BACKUP_DIR, filename)
        if restore_from_json(file_path, mode):
            return {"success": True, "message": f"成功以 {mode} 模式恢复备份"}
        return {"success": False, "detail": "恢复失败"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/yutu/organize/confirm")
async def confirm_organize(data: dict, username: str = ""):
    """确认整理结果：应用改进后的方案"""
    if not is_builtin_superuser(username):
        raise HTTPException(status_code=403, detail="只有超级用户可以确认")

    improved_records = data.get("records", [])
    if not improved_records:
        return {"success": False, "detail": "没有记录可更新"}

    try:
        from .yutu_zhanyilu import update_error_solution
        updated_count = 0
        for record in improved_records:
            # 兼容 improved_solution 字段（VLLM返回）或 solution 字段
            sol = record.get("improved_solution") or record.get("solution") or ""
            if update_error_solution(record.get("error_hash"), sol):
                updated_count += 1

        return {"success": True, "updated_count": updated_count}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/yutu/organize/cancel")
async def cancel_organize(data: dict, username: str = ""):
    """取消整理：恢复到原始备份"""
    if not is_builtin_superuser(username):
        raise HTTPException(status_code=403, detail="只有超级用户可以取消")

    # 取消操作不需要实际恢复，因为原始数据未修改
    return {"success": True, "message": "已取消整理，原始记录保持不变"}


# ========== 数据库连接与查询 API ==========
DB_TYPE_ALIASES = {
    "postgres": "postgresql",
    "postgresql": "postgresql",
    "mysql": "mysql",
    "mssql": "mssql",
    "sqlserver": "mssql",
    "sqlite": "sqlite",
    "oracle": "oracle",
}

DB_DEFAULT_PORTS = {
    "mysql": 3306,
    "postgresql": 5432,
    "mssql": 1433,
    "oracle": 1521,
}

DB_DRIVER_NAMES = {
    "mysql": "mysql+pymysql",
    "postgresql": "postgresql+psycopg2",
    "mssql": "mssql+pymssql",
    "oracle": "oracle+cx_oracle",
}

DB_DRIVER_PACKAGES = {
    "mysql": "pymysql",
    "postgresql": "psycopg2-binary",
    "mssql": "pymssql",
    "oracle": "cx_Oracle",
}


def normalize_db_type(db_type: Optional[str]) -> str:
    normalized = str(db_type or "").strip().lower()
    if not normalized:
        raise ValueError("缺少 db_type，请选择数据库类型")
    resolved = DB_TYPE_ALIASES.get(normalized)
    if not resolved:
        supported = ", ".join(sorted(set(DB_TYPE_ALIASES.values())))
        raise ValueError(f"不支持的数据库类型: {db_type}。支持: {supported}")
    return resolved


def build_db_url(db_type: str, config: dict) -> str:
    """构建 SQLAlchemy 连接字符串（使用 URL.create，避免特殊字符导致解析错误）"""
    db_type = normalize_db_type(db_type)
    config = config or {}

    if db_type == "sqlite":
        database = str(config.get("database", "")).strip()
        if not database:
            raise ValueError("SQLite 连接需要提供数据库文件路径")
        if database == ":memory:":
            return "sqlite:///:memory:"
        return URL.create(drivername="sqlite", database=database).render_as_string(hide_password=False)

    drivername = DB_DRIVER_NAMES[db_type]
    host = str(config.get("host", "localhost") or "localhost").strip() or "localhost"
    database = str(config.get("database", "")).strip()
    if not database:
        raise ValueError("请填写数据库名称")

    port_raw = str(config.get("port", "") or "").strip()
    if port_raw:
        if not port_raw.isdigit():
            raise ValueError("端口必须为数字")
        port = int(port_raw)
    else:
        port = DB_DEFAULT_PORTS.get(db_type)

    user = str(config.get("user", "") or "").strip() or None
    password_raw = config.get("password", None)
    password = None if password_raw in (None, "") else str(password_raw)
    query = {"charset": "utf8mb4"} if db_type == "mysql" else None

    return URL.create(
        drivername=drivername,
        username=user,
        password=password,
        host=host,
        port=port,
        database=database,
        query=query,
    ).render_as_string(hide_password=False)


def build_db_engine(db_type: str, config: dict):
    db_type = normalize_db_type(db_type)
    url = build_db_url(db_type, config)

    connect_args = {}
    if db_type in ("mysql", "postgresql"):
        connect_args["connect_timeout"] = 5
    elif db_type == "mssql":
        connect_args["login_timeout"] = 5

    engine_kwargs = {"pool_pre_ping": True}
    if connect_args:
        engine_kwargs["connect_args"] = connect_args
    return create_engine(url, **engine_kwargs)


def list_accessible_databases(db_type: str, config: dict) -> List[str]:
    """根据连接信息列出可访问的数据库名称。"""
    normalized_type = normalize_db_type(db_type)
    config = config or {}

    if normalized_type == "sqlite":
        database = str(config.get("database", "")).strip()
        if not database:
            raise ValueError("SQLite 模式下请先填写数据库文件路径")
        return [database]

    host = str(config.get("host", "localhost") or "localhost").strip() or "localhost"
    user = str(config.get("user", "") or "").strip() or None
    password_raw = config.get("password", None)
    password = None if password_raw in (None, "") else str(password_raw)

    port_raw = str(config.get("port", "") or "").strip()
    if port_raw:
        if not port_raw.isdigit():
            raise ValueError("端口必须为数字")
        port = int(port_raw)
    else:
        port = DB_DEFAULT_PORTS.get(normalized_type)

    def _create_engine_for_db(database_name: Optional[str]):
        query = {"charset": "utf8mb4"} if normalized_type == "mysql" else None
        url = URL.create(
            drivername=DB_DRIVER_NAMES[normalized_type],
            username=user,
            password=password,
            host=host,
            port=port,
            database=database_name,
            query=query,
        )
        connect_args = {}
        if normalized_type in ("mysql", "postgresql"):
            connect_args["connect_timeout"] = 5
        elif normalized_type == "mssql":
            connect_args["login_timeout"] = 5
        kwargs: Dict[str, Any] = {"pool_pre_ping": True}
        if connect_args:
            kwargs["connect_args"] = connect_args
        return create_engine(url, **kwargs)

    if normalized_type == "postgresql":
        candidates = []
        configured_db = str(config.get("database", "") or "").strip()
        if configured_db:
            candidates.append(configured_db)
        candidates.extend(["postgres", "template1"])
        seen = set()
        dedup_candidates = []
        for item in candidates:
            if item and item not in seen:
                dedup_candidates.append(item)
                seen.add(item)

        last_error: Optional[Exception] = None
        for db_name in dedup_candidates:
            engine = None
            try:
                engine = _create_engine_for_db(db_name)
                with engine.connect() as conn:
                    rows = conn.execute(
                        text(
                            "SELECT datname FROM pg_database "
                            "WHERE datistemplate = false AND datallowconn = true "
                            "ORDER BY datname"
                        )
                    ).fetchall()
                names = [str(row[0]) for row in rows if row and row[0]]
                if names:
                    return names
            except Exception as exc:
                last_error = exc
            finally:
                if engine is not None:
                    engine.dispose()

        if configured_db:
            return [configured_db]
        if last_error is not None:
            raise last_error
        raise ValueError("未能获取 PostgreSQL 数据库列表")

    if normalized_type == "mysql":
        engine = None
        try:
            engine = _create_engine_for_db(None)
            with engine.connect() as conn:
                rows = conn.execute(text("SHOW DATABASES")).fetchall()
            names = [str(row[0]) for row in rows if row and row[0]]
            return names
        finally:
            if engine is not None:
                engine.dispose()

    if normalized_type == "mssql":
        engine = None
        admin_db = str(config.get("database", "") or "").strip() or "master"
        try:
            engine = _create_engine_for_db(admin_db)
            with engine.connect() as conn:
                rows = conn.execute(text("SELECT name FROM sys.databases ORDER BY name")).fetchall()
            names = [str(row[0]) for row in rows if row and row[0]]
            return names
        finally:
            if engine is not None:
                engine.dispose()

    # Oracle 无统一的“数据库名列表”语义，回退为当前用户输入值
    configured_db = str(config.get("database", "") or "").strip()
    if configured_db:
        return [configured_db]
    raise ValueError("当前数据库类型暂不支持自动拉取数据库名称，请手动填写数据库名称")


def format_db_error(exc: Exception, db_type: Optional[str]) -> str:
    raw_message = str(exc).strip() or exc.__class__.__name__
    # 避免把连接串中的明文密码返回前端
    safe_message = re.sub(r"://([^:/@]+):([^@]+)@", r"://\1:***@", raw_message)
    lower = safe_message.lower()

    try:
        normalized_type = normalize_db_type(db_type)
    except Exception:
        normalized_type = str(db_type or "").strip().lower()

    driver_pkg = DB_DRIVER_PACKAGES.get(normalized_type)
    if driver_pkg and ("no module named" in lower or "can't load plugin" in lower):
        return f"缺少数据库驱动，请安装 `{driver_pkg}` 后重试。原始错误: {safe_message}"
    if "password authentication failed" in lower or "access denied" in lower:
        return "数据库账号或密码错误，请检查连接配置"
    if "connection refused" in lower or "can't connect" in lower:
        return "无法连接到数据库服务，请确认主机、端口和数据库服务状态"
    if "could not translate host name" in lower or "name or service not known" in lower:
        return "数据库主机名无法解析，请检查主机地址配置"
    return safe_message


def _escape_md_cell(value: Any) -> str:
    text_value = str(value or "").replace("\n", " ").strip()
    if not text_value:
        return "-"
    return text_value.replace("|", "\\|")


def _quote_identifier(identifier: str, db_type: str) -> str:
    if db_type == "mysql":
        return f"`{identifier.replace('`', '``')}`"
    if db_type == "mssql":
        return f"[{identifier.replace(']', ']]')}]"
    return '"' + identifier.replace('"', '""') + '"'


def _build_qualified_table_name(schema_name: Optional[str], table_name: str, db_type: str) -> str:
    if schema_name:
        return f"{_quote_identifier(schema_name, db_type)}.{_quote_identifier(table_name, db_type)}"
    return _quote_identifier(table_name, db_type)


def _shorten_text(value: Any, max_len: int = 120) -> str:
    text_value = str(value or "").strip()
    if len(text_value) <= max_len:
        return text_value
    return text_value[: max_len - 1] + "…"


def _safe_iso_ts(value: Any) -> str:
    text_value = str(value or "").strip()
    if not text_value:
        return ""
    try:
        return datetime.fromisoformat(text_value).isoformat(timespec="seconds")
    except Exception:
        return text_value


def upsert_database_knowledge_source(
    username: str,
    db_type: str,
    config: dict,
    snapshot: Dict[str, Any],
) -> Dict[str, Any]:
    """把数据库快照持久化为用户级知识库条目。"""
    username = str(username or "default").strip() or "default"

    source_entry = {
        "source_label": str(snapshot.get("source_label", "") or ""),
        "db_type": normalize_db_type(db_type),
        "host": str(config.get("host", "") or "").strip(),
        "port": str(config.get("port", "") or "").strip(),
        "database": str(snapshot.get("database", "") or "").strip(),
        "table_count": int(snapshot.get("table_count", 0) or 0),
        "column_count": int(snapshot.get("column_count", 0) or 0),
        "loaded_at": datetime.now().isoformat(timespec="seconds"),
        "tables": snapshot.get("tables", []),
        "summary": _shorten_text(snapshot.get("summary", ""), max_len=500),
    }

    kb_data = _load_user_config(username, "database_knowledge_base.json", {"sources": []})
    sources = kb_data.get("sources", [])
    if not isinstance(sources, list):
        sources = []

    source_label = source_entry["source_label"]
    source_db = source_entry["database"]
    source_host = source_entry["host"]
    source_port = source_entry["port"]
    replaced = False

    for idx, existing in enumerate(sources):
        if not isinstance(existing, dict):
            continue
        existing_label = str(existing.get("source_label", "") or "").strip()
        existing_db = str(existing.get("database", "") or "").strip()
        existing_host = str(existing.get("host", "") or "").strip()
        existing_port = str(existing.get("port", "") or "").strip()

        if (
            (source_label and existing_label == source_label)
            or (
                source_db
                and source_host == existing_host
                and source_port == existing_port
                and source_db == existing_db
            )
        ):
            sources[idx] = source_entry
            replaced = True
            break

    if not replaced:
        sources.append(source_entry)

    sources = sorted(
        [item for item in sources if isinstance(item, dict)],
        key=lambda item: _safe_iso_ts(item.get("loaded_at", "")),
        reverse=True,
    )[:30]

    kb_data["sources"] = sources
    _save_user_config(username, "database_knowledge_base.json", kb_data)
    return source_entry


def build_database_knowledge_context(
    username: str,
    user_query: str,
    preferred_source_labels: Optional[List[str]] = None,
    analysis_language: str = "zh-CN",
    max_tables: int = 6,
    max_columns_per_table: int = 10,
    max_chars: int = 2600,
) -> str:
    """从数据库知识库检索与当前问题相关的结构化上下文。"""
    username = str(username or "default").strip() or "default"
    kb_data = _load_user_config(username, "database_knowledge_base.json", {"sources": []})
    raw_sources = kb_data.get("sources", [])
    if not isinstance(raw_sources, list) or not raw_sources:
        return ""

    sources = [item for item in raw_sources if isinstance(item, dict)]
    if not sources:
        return ""

    preferred_labels = set(
        str(item or "").strip()
        for item in (preferred_source_labels or [])
        if str(item or "").strip()
    )
    if preferred_labels:
        matched_sources = [
            item
            for item in sources
            if str(item.get("source_label", "") or "").strip() in preferred_labels
        ]
        if matched_sources:
            sources = matched_sources

    query_text = str(user_query or "").lower()
    query_tokens = [
        token for token in re.findall(r"[a-zA-Z0-9_\u4e00-\u9fff]{2,}", query_text)
        if token
    ]

    def table_score(table_item: Dict[str, Any]) -> int:
        table_name = str(table_item.get("name", "") or "").lower()
        table_comment = str(table_item.get("table_comment", "") or "").lower()
        columns = table_item.get("columns", [])
        column_blob = " ".join(
            f"{str(col.get('name', '') or '')} {str(col.get('definition', '') or '')}"
            for col in columns
            if isinstance(col, dict)
        ).lower()

        if not query_tokens:
            return 1

        score = 0
        for token in query_tokens:
            if token in table_name:
                score += 6
            if token in table_comment:
                score += 3
            if token in column_blob:
                score += 2
        return score

    lines: List[str] = []
    if analysis_language == "en":
        lines.append("Database knowledge base retrieval context:")
    else:
        lines.append("数据库知识库检索上下文：")

    source_count = 0
    for source in sources:
        if source_count >= 2:
            break

        tables = source.get("tables", [])
        if not isinstance(tables, list):
            tables = []
        ranked_tables = sorted(
            [item for item in tables if isinstance(item, dict)],
            key=lambda item: (table_score(item), int(item.get("row_count") or 0)),
            reverse=True,
        )

        if not ranked_tables:
            continue

        source_label = str(source.get("source_label", "") or "").strip()
        table_count = int(source.get("table_count", len(tables)) or len(tables))
        column_count = int(source.get("column_count", 0) or 0)

        lines.append(
            f"- 数据源: {source_label} (表{table_count}张, 字段{column_count}个)"
        )

        selected_tables = ranked_tables[:max_tables]
        for table in selected_tables:
            table_name = str(table.get("name", "") or "")
            row_count = table.get("row_count")
            table_comment = _shorten_text(table.get("table_comment", ""), max_len=80)

            if analysis_language == "en":
                table_line = f"  - Table: {table_name}, rows={row_count if row_count is not None else 'unknown'}"
            else:
                table_line = f"  - 表: {table_name}, 行数={row_count if row_count is not None else '未知'}"

            if table_comment:
                table_line += f", 说明={table_comment}"
            lines.append(table_line)

            columns = table.get("columns", [])
            if not isinstance(columns, list) or not columns:
                continue

            column_entries: List[str] = []
            for col in columns[:max_columns_per_table]:
                if not isinstance(col, dict):
                    continue
                col_name = str(col.get("name", "") or "")
                col_type = str(col.get("type", "") or "")
                col_def = _shorten_text(col.get("definition", ""), max_len=40)
                entry = f"{col_name}({col_type})"
                if col_def:
                    entry += f":{col_def}"
                column_entries.append(entry)

            if column_entries:
                if analysis_language == "en":
                    lines.append(f"    key columns: {', '.join(column_entries)}")
                else:
                    lines.append(f"    关键字段: {', '.join(column_entries)}")

        source_count += 1

    context_text = "\n".join(lines).strip()
    if len(context_text) > max_chars:
        context_text = context_text[: max_chars - 1] + "…"
    return context_text


def build_database_context_snapshot(db_type: str, config: dict) -> Dict[str, Any]:
    """提取数据库结构快照，用于注入分析上下文。"""
    normalized_type = normalize_db_type(db_type)
    config = config or {}
    database_name = str(config.get("database", "") or "").strip()

    if normalized_type != "sqlite" and not database_name:
        raise ValueError("请先选择数据库名称后再导入上下文")

    from sqlalchemy import inspect

    engine = build_db_engine(normalized_type, config)
    inspector = inspect(engine)

    source_label = (
        f"{normalized_type}@{database_name}"
        if normalized_type == "sqlite"
        else f"{normalized_type}@{str(config.get('host', 'localhost') or 'localhost').strip()}:{str(config.get('port', '') or DB_DEFAULT_PORTS.get(normalized_type, '')).strip() or DB_DEFAULT_PORTS.get(normalized_type, '')}/{database_name}"
    )

    lines: List[str] = [
        "# 数据库知识库上下文快照",
        f"- 数据源: {source_label}",
        f"- 导入时间: {datetime.now().isoformat(timespec='seconds')}",
        "",
    ]

    table_targets: List[Tuple[Optional[str], str]] = []

    if normalized_type in ("postgresql", "mssql", "oracle"):
        try:
            schema_names = inspector.get_schema_names()
        except Exception:
            schema_names = []

        for schema_name in schema_names:
            if not schema_name:
                continue
            lower_schema = str(schema_name).lower()
            upper_schema = str(schema_name).upper()

            if normalized_type == "postgresql" and lower_schema in {"pg_catalog", "information_schema", "pg_toast"}:
                continue
            if normalized_type == "mssql" and upper_schema in {"INFORMATION_SCHEMA", "SYS"}:
                continue
            if normalized_type == "oracle" and upper_schema in {"SYS", "SYSTEM", "XDB", "MDSYS", "CTXSYS", "WMSYS"}:
                continue

            try:
                schema_tables = inspector.get_table_names(schema=schema_name)
            except Exception:
                schema_tables = []

            for table_name in schema_tables:
                table_targets.append((schema_name, table_name))
    else:
        try:
            base_tables = inspector.get_table_names()
        except Exception:
            base_tables = []
        for table_name in base_tables:
            table_targets.append((None, table_name))

    table_targets = sorted(table_targets, key=lambda item: ((item[0] or ""), item[1]))

    total_tables = 0
    total_columns = 0
    table_records: List[Dict[str, Any]] = []

    with engine.connect() as conn:
        for schema_name, table_name in table_targets:
            total_tables += 1

            try:
                columns = inspector.get_columns(table_name, schema=schema_name)
            except Exception:
                columns = []

            total_columns += len(columns)

            try:
                pk_info = inspector.get_pk_constraint(table_name, schema=schema_name) or {}
                pk_columns = set(pk_info.get("constrained_columns") or [])
            except Exception:
                pk_columns = set()

            table_comment = ""
            try:
                comment_payload = inspector.get_table_comment(table_name, schema=schema_name) or {}
                table_comment = str(comment_payload.get("text") or "").strip()
            except Exception:
                table_comment = ""

            row_count: Optional[int] = None
            try:
                qualified_table_name = _build_qualified_table_name(schema_name, table_name, normalized_type)
                row_count_result = conn.execute(text(f"SELECT COUNT(*) AS cnt FROM {qualified_table_name}"))
                row_count_value = row_count_result.scalar()
                row_count = int(row_count_value) if row_count_value is not None else 0
            except Exception:
                row_count = None

            display_table_name = f"{schema_name}.{table_name}" if schema_name else table_name
            lines.append(f"## 表: {display_table_name}")
            lines.append(f"- 数据量(行数): {row_count if row_count is not None else '未知'}")
            lines.append(f"- 字段数量: {len(columns)}")
            if table_comment:
                lines.append(f"- 表定义/说明: {table_comment}")

            table_record: Dict[str, Any] = {
                "schema": schema_name or "",
                "name": display_table_name,
                "row_count": row_count,
                "column_count": len(columns),
                "table_comment": table_comment,
                "columns": [],
            }

            if not columns:
                lines.append("- 字段信息: 无法读取")
                lines.append("")
                table_records.append(table_record)
                continue

            lines.append("| 字段名 | 数据类型 | 可空 | 默认值 | 字段定义/说明 | 主键 |")
            lines.append("| --- | --- | --- | --- | --- | --- |")
            for column in columns:
                column_name = str(column.get("name") or "")
                column_type = _escape_md_cell(column.get("type") or "")
                nullable = "是" if column.get("nullable", True) else "否"
                default_value = _escape_md_cell(column.get("default") or "")
                definition_text = _escape_md_cell(column.get("comment") or column.get("definition") or "")
                is_pk = "是" if column_name in pk_columns else "-"
                lines.append(
                    f"| {_escape_md_cell(column_name)} | {column_type} | {nullable} | {default_value} | {definition_text} | {is_pk} |"
                )
                table_record["columns"].append(
                    {
                        "name": column_name,
                        "type": str(column.get("type") or ""),
                        "nullable": bool(column.get("nullable", True)),
                        "default": str(column.get("default") or ""),
                        "definition": str(column.get("comment") or column.get("definition") or ""),
                        "is_pk": column_name in pk_columns,
                    }
                )
            lines.append("")
            table_records.append(table_record)

    engine.dispose()

    lines.insert(3, f"- 表数量: {total_tables}")
    lines.insert(4, f"- 字段总数: {total_columns}")

    return {
        "source_label": source_label,
        "database": database_name,
        "table_count": total_tables,
        "column_count": total_columns,
        "tables": table_records,
        "summary": f"{source_label} 共 {total_tables} 张表、{total_columns} 个字段",
        "context_text": "\n".join(lines).strip(),
    }


def _schema_table_id(schema_name: Optional[str], table_name: str) -> str:
    schema_part = str(schema_name or "default").strip() or "default"
    table_part = str(table_name or "table").strip() or "table"
    return re.sub(r"[^A-Za-z0-9_.-]+", "_", f"{schema_part}.{table_part}").strip("_")


def build_database_schema_graph(db_type: str, config: dict) -> Dict[str, Any]:
    """Build ChartDB-style table/relationship graph metadata from a live database."""
    normalized_type = normalize_db_type(db_type)
    config = config or {}
    database_name = str(config.get("database", "") or "").strip()

    if normalized_type != "sqlite" and not database_name:
        raise ValueError("请先选择数据库名称后再生成关系图")

    from sqlalchemy import inspect

    engine = build_db_engine(normalized_type, config)
    inspector = inspect(engine)

    source_label = (
        f"{normalized_type}@{database_name}"
        if normalized_type == "sqlite"
        else f"{normalized_type}@{str(config.get('host', 'localhost') or 'localhost').strip()}:{str(config.get('port', '') or DB_DEFAULT_PORTS.get(normalized_type, '')).strip() or DB_DEFAULT_PORTS.get(normalized_type, '')}/{database_name}"
    )

    table_targets: List[Tuple[Optional[str], str]] = []
    if normalized_type in ("postgresql", "mssql", "oracle"):
        try:
            schema_names = inspector.get_schema_names()
        except Exception:
            schema_names = []

        for schema_name in schema_names:
            if not schema_name:
                continue
            lower_schema = str(schema_name).lower()
            upper_schema = str(schema_name).upper()
            if normalized_type == "postgresql" and lower_schema in {"pg_catalog", "information_schema", "pg_toast"}:
                continue
            if normalized_type == "mssql" and upper_schema in {"INFORMATION_SCHEMA", "SYS"}:
                continue
            if normalized_type == "oracle" and upper_schema in {"SYS", "SYSTEM", "XDB", "MDSYS", "CTXSYS", "WMSYS"}:
                continue
            try:
                schema_tables = inspector.get_table_names(schema=schema_name)
            except Exception:
                schema_tables = []
            for table_name in schema_tables:
                table_targets.append((schema_name, table_name))
    else:
        try:
            base_tables = inspector.get_table_names()
        except Exception:
            base_tables = []
        for table_name in base_tables:
            table_targets.append((None, table_name))

    table_targets = sorted(table_targets, key=lambda item: ((item[0] or ""), item[1]))
    tables: List[Dict[str, Any]] = []
    relationships: List[Dict[str, Any]] = []
    table_id_lookup: Dict[Tuple[str, str], str] = {}

    with engine.connect() as conn:
        for schema_name, table_name in table_targets:
            table_id = _schema_table_id(schema_name, table_name)
            schema_key = str(schema_name or "").strip()
            table_id_lookup[(schema_key, table_name)] = table_id
            display_name = f"{schema_name}.{table_name}" if schema_name else table_name

            try:
                columns = inspector.get_columns(table_name, schema=schema_name)
            except Exception:
                columns = []

            try:
                pk_info = inspector.get_pk_constraint(table_name, schema=schema_name) or {}
                pk_columns = set(pk_info.get("constrained_columns") or [])
            except Exception:
                pk_columns = set()

            row_count: Optional[int] = None
            try:
                qualified_table_name = _build_qualified_table_name(schema_name, table_name, normalized_type)
                row_count_value = conn.execute(text(f"SELECT COUNT(*) AS cnt FROM {qualified_table_name}")).scalar()
                row_count = int(row_count_value) if row_count_value is not None else 0
            except Exception:
                row_count = None

            table_columns = []
            for column in columns:
                column_name = str(column.get("name") or "")
                table_columns.append(
                    {
                        "name": column_name,
                        "type": str(column.get("type") or ""),
                        "nullable": bool(column.get("nullable", True)),
                        "default": str(column.get("default") or ""),
                        "definition": str(column.get("comment") or column.get("definition") or ""),
                        "is_pk": column_name in pk_columns,
                    }
                )

            tables.append(
                {
                    "id": table_id,
                    "schema": schema_key,
                    "name": table_name,
                    "display_name": display_name,
                    "row_count": row_count,
                    "columns": table_columns,
                }
            )

        for schema_name, table_name in table_targets:
            source_table_id = table_id_lookup.get((str(schema_name or "").strip(), table_name))
            if not source_table_id:
                continue

            try:
                foreign_keys = inspector.get_foreign_keys(table_name, schema=schema_name)
            except Exception:
                foreign_keys = []

            for index, foreign_key in enumerate(foreign_keys or []):
                referred_table = str(foreign_key.get("referred_table") or "").strip()
                if not referred_table:
                    continue
                referred_schema = foreign_key.get("referred_schema")
                referred_schema_key = str(referred_schema if referred_schema is not None else (schema_name or "")).strip()
                target_table_id = table_id_lookup.get((referred_schema_key, referred_table))
                if not target_table_id and not referred_schema_key:
                    target_table_id = table_id_lookup.get(("", referred_table))
                if not target_table_id:
                    continue

                constrained_columns = [str(item) for item in (foreign_key.get("constrained_columns") or [])]
                referred_columns = [str(item) for item in (foreign_key.get("referred_columns") or [])]
                fk_name = str(foreign_key.get("name") or "").strip()
                relationship_id = re.sub(
                    r"[^A-Za-z0-9_.-]+",
                    "_",
                    fk_name or f"fk_{source_table_id}_{target_table_id}_{index}",
                ).strip("_")

                relationships.append(
                    {
                        "id": relationship_id,
                        "name": fk_name or relationship_id,
                        "source_table_id": source_table_id,
                        "source_columns": constrained_columns,
                        "target_table_id": target_table_id,
                        "target_columns": referred_columns,
                        "relationship_type": "many_to_one",
                    }
                )

    engine.dispose()

    return {
        "source_label": source_label,
        "db_type": normalized_type,
        "database": database_name,
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "tables": tables,
        "relationships": relationships,
        "summary": f"{source_label} 共 {len(tables)} 张表、{len(relationships)} 条外键关系",
    }


def _profile_workspace_file(file_path: Path, workspace_root: Path) -> Dict[str, Any]:
    relative_path = str(file_path.relative_to(workspace_root))
    profile: Dict[str, Any] = {
        "name": file_path.name,
        "path": relative_path,
        "size": file_path.stat().st_size if file_path.exists() else 0,
        "extension": file_path.suffix.lower().lstrip("."),
        "status": "listed",
        "columns": [],
        "row_count_sampled": None,
    }

    extension = file_path.suffix.lower()
    try:
        if extension in {".csv", ".tsv", ".txt"}:
            separator = "\t" if extension == ".tsv" else ","
            df = pd.read_csv(file_path, nrows=1000, sep=separator)
        elif extension in {".xlsx", ".xls"}:
            df = pd.read_excel(file_path, nrows=1000)
        elif extension == ".json":
            df = pd.read_json(file_path)
            if len(df) > 1000:
                df = df.head(1000)
        else:
            return profile

        profile["status"] = "profiled"
        profile["row_count_sampled"] = int(len(df))
        profile["column_count"] = int(len(df.columns))
        profile["columns"] = [
            {
                "name": str(column),
                "dtype": str(df[column].dtype),
                "non_null_sampled": int(df[column].notna().sum()),
                "null_sampled": int(df[column].isna().sum()),
                "sample_values": [str(value)[:80] for value in df[column].dropna().head(3).tolist()],
            }
            for column in df.columns[:80]
        ]
    except Exception as exc:
        profile["status"] = "profile_failed"
        profile["error"] = str(exc)[:240]
    return profile


def _collect_workspace_file_profiles(session_id: str, username: str) -> List[Dict[str, Any]]:
    workspace_root = Path(get_session_workspace(session_id, username)).resolve()
    profiles: List[Dict[str, Any]] = []
    excluded_dirs = {"generated", ".lib", "__pycache__"}

    for file_path in sorted(workspace_root.rglob("*")):
        if not file_path.is_file():
            continue
        relative_parts = file_path.relative_to(workspace_root).parts
        if any(part in excluded_dirs or part.startswith(".") for part in relative_parts):
            continue
        profiles.append(_profile_workspace_file(file_path, workspace_root))
        if len(profiles) >= 30:
            break
    return profiles


def _build_data_profile_skill_markdown(
    username: str,
    session_id: str,
    db_graphs: List[Dict[str, Any]],
    file_profiles: List[Dict[str, Any]],
) -> str:
    generated_at = datetime.now().isoformat(timespec="seconds")
    lines: List[str] = [
        "# Data Exploration SKILL",
        "",
        "## Purpose",
        "",
        "This SKILL captures the current data landscape for subsequent customs risk analysis. Use it before analysis planning, SQL generation, feature engineering, anomaly detection, and report writing.",
        "",
        "## Runtime Context",
        "",
        f"- User: {username or 'default'}",
        f"- Session: {session_id or 'default'}",
        f"- Generated at: {generated_at}",
        "",
        "## Database Sources",
        "",
    ]

    if db_graphs:
        for graph in db_graphs:
            lines.extend(
                [
                    f"### {graph.get('source_label', 'database')}",
                    "",
                    f"- Type: {graph.get('db_type', '-')}",
                    f"- Tables: {len(graph.get('tables', []) or [])}",
                    f"- Relationships: {len(graph.get('relationships', []) or [])}",
                    "",
                ]
            )
            for table in (graph.get("tables", []) or [])[:20]:
                columns = table.get("columns", []) or []
                pk_columns = [col.get("name") for col in columns if col.get("is_pk")]
                key_columns = ", ".join(str(col.get("name")) for col in columns[:12])
                lines.append(
                    f"- `{table.get('display_name') or table.get('name')}`: rows={table.get('row_count', 'unknown')}, columns={len(columns)}, pk={', '.join(pk_columns) or '-'}, key fields={key_columns or '-'}"
                )
            if graph.get("relationships"):
                lines.extend(["", "Relationships:"])
                table_names = {table.get("id"): table.get("display_name") or table.get("name") for table in graph.get("tables", []) or []}
                for relationship in (graph.get("relationships", []) or [])[:30]:
                    source_name = table_names.get(relationship.get("source_table_id"), relationship.get("source_table_id"))
                    target_name = table_names.get(relationship.get("target_table_id"), relationship.get("target_table_id"))
                    lines.append(
                        f"- `{source_name}`.{','.join(relationship.get('source_columns', []) or [])} -> `{target_name}`.{','.join(relationship.get('target_columns', []) or [])}"
                    )
            lines.append("")
    else:
        lines.extend(["- No live database source was available for this report.", ""])

    lines.extend(["## Uploaded / Workspace Files", ""])
    if file_profiles:
        for profile in file_profiles:
            lines.extend(
                [
                    f"### {profile.get('path')}",
                    "",
                    f"- Size: {profile.get('size', 0)} bytes",
                    f"- Type: {profile.get('extension') or '-'}",
                    f"- Profile status: {profile.get('status')}",
                ]
            )
            if profile.get("row_count_sampled") is not None:
                lines.append(f"- Sampled rows: {profile.get('row_count_sampled')}")
            columns = profile.get("columns", []) or []
            if columns:
                lines.append("- Columns:")
                for column in columns[:30]:
                    sample_values = ", ".join(column.get("sample_values", []) or [])
                    lines.append(
                        f"  - `{column.get('name')}` ({column.get('dtype')}), non-null sample={column.get('non_null_sampled')}, examples={sample_values or '-'}"
                    )
            if profile.get("error"):
                lines.append(f"- Profiling note: {profile.get('error')}")
            lines.append("")
    else:
        lines.extend(["- No user-uploaded workspace data files were detected.", ""])

    lines.extend(
        [
            "## Analysis Contract",
            "",
            "When analyzing this workspace:",
            "",
            "1. Prefer the database relationship graph for join planning and entity grain decisions.",
            "2. Treat table primary keys, foreign keys, row counts, and high-null columns as first-order data quality signals.",
            "3. Reconcile uploaded files with database entities before feature engineering; identify duplicate keys, missing periods, and inconsistent code mappings.",
            "4. For customs risk analysis, prioritize import/export主体、商品编码、贸易方式、原产地/目的地、申报价格、数量、税则、许可证件、物流路径 and time-window anomalies.",
            "5. Before producing final conclusions, record assumptions about table joins, missing values, and derived fields.",
            "",
            "## Suggested Next Prompts",
            "",
            "- 基于 Data Exploration SKILL，识别当前数据中最适合开展风险画像的主表和关联表。",
            "- 根据字段完整性、关系图和样本值，为该数据集设计走私违规、逃证逃税、安全准入风险特征。",
            "- 生成一份分阶段分析计划：数据校验、指标构造、异常发现、证据链整理、报告输出。",
        ]
    )
    return "\n".join(lines).strip() + "\n"


@app.post("/api/db/test")
async def test_db_connection(body: dict = Body(...)):
    """测试数据库连接"""
    try:
        db_type = body.get("db_type")
        config = body.get("config", {})
        engine = build_db_engine(db_type, config)
        try:
            with engine.connect() as conn:
                conn.execute(text("SELECT 1"))
        finally:
            engine.dispose()
        return {"success": True, "message": "Connection successful"}
    except Exception as e:
        print(f"DB Test Error: {e}")
        return {"success": False, "message": format_db_error(e, body.get("db_type"))}


@app.post("/api/db/list-databases")
async def list_databases(body: dict = Body(...)):
    """根据连接参数拉取数据库名称列表，用于前端数据库下拉选择。"""
    try:
        db_type = body.get("db_type")
        config = body.get("config", {})
        databases = list_accessible_databases(db_type, config)
        return {"success": True, "databases": databases}
    except Exception as e:
        print(f"DB List Databases Error: {e}")
        return {"success": False, "message": format_db_error(e, body.get("db_type"))}


@app.post("/api/db/context/load")
async def load_db_context_into_session(body: dict = Body(...)):
    """将当前数据库结构导入会话上下文，作为分析知识库。"""
    try:
        db_type = body.get("db_type")
        config = body.get("config", {})
        session_id = str(body.get("session_id", "default") or "default")
        username = str(body.get("username", "default") or "default")

        snapshot = build_database_context_snapshot(db_type, config)
        context_text = snapshot.get("context_text", "")
        kb_entry = upsert_database_knowledge_source(username, db_type, config, snapshot)

        workspace_dir = get_session_workspace(session_id, username)
        generated_dir = Path(workspace_dir) / "generated"
        generated_dir.mkdir(parents=True, exist_ok=True)

        safe_db_name = re.sub(r"[^A-Za-z0-9_-]+", "_", str(snapshot.get("database") or "database")).strip("_") or "database"
        base_name = f"Database_Context_{safe_db_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        md_path = _save_md(context_text, base_name, str(generated_dir))

        file_url = None
        if md_path:
            file_url = build_download_url(f"{username}/{session_id}/generated/{Path(md_path).name}")

        loaded_at = datetime.now().isoformat(timespec="seconds")

        SESSION_DB_CONTEXT[session_id] = {
            "context_text": context_text,
            "source_label": snapshot.get("source_label", ""),
            "table_count": snapshot.get("table_count", 0),
            "column_count": snapshot.get("column_count", 0),
            "loaded_at": loaded_at,
            "snapshot_file": Path(md_path).name if md_path else "",
            "snapshot_url": file_url or "",
            "knowledge_summary": kb_entry.get("summary", ""),
        }

        return {
            "success": True,
            "message": "数据库上下文已导入当前会话",
            "source_label": snapshot.get("source_label", ""),
            "table_count": snapshot.get("table_count", 0),
            "column_count": snapshot.get("column_count", 0),
            "context_length": len(context_text),
            "snapshot_file": Path(md_path).name if md_path else None,
            "snapshot_url": file_url,
            "knowledge_summary": kb_entry.get("summary", ""),
            "loaded_at": loaded_at,
        }
    except Exception as e:
        print(f"DB Context Load Error: {e}")
        return {"success": False, "message": format_db_error(e, body.get("db_type"))}


@app.post("/api/db/schema/graph")
async def load_db_schema_graph(body: dict = Body(...)):
    """返回数据库表/字段/外键关系图，用于前端可视化。"""
    try:
        db_type = body.get("db_type")
        config = body.get("config", {})
        graph = build_database_schema_graph(db_type, config)
        return {"success": True, "graph": graph, "message": graph.get("summary", "关系图已生成")}
    except Exception as e:
        print(f"DB Schema Graph Error: {e}")
        return {"success": False, "message": format_db_error(e, body.get("db_type"))}


@app.post("/api/data/profile-report")
async def create_data_profile_report(body: dict = Body(...)):
    """生成当前数据库与 workspace 文件的数据探查 SKILL 文档。"""
    try:
        session_id = str(body.get("session_id", "default") or "default")
        username = str(body.get("username", "default") or "default")
        selected_database_sources = body.get("selected_database_sources", [])

        db_graphs: List[Dict[str, Any]] = []
        source_candidates: List[Dict[str, Any]] = []
        if isinstance(selected_database_sources, list):
            source_candidates.extend([item for item in selected_database_sources if isinstance(item, dict)])

        db_type = body.get("db_type")
        config = body.get("config")
        if db_type and isinstance(config, dict):
            source_candidates.append({"dbType": db_type, "config": config})

        seen_sources = set()
        for source in source_candidates[:5]:
            source_db_type = source.get("dbType") or source.get("db_type")
            source_config = source.get("config", {}) if isinstance(source.get("config", {}), dict) else {}
            source_key = json.dumps({"db_type": source_db_type, "config": source_config}, sort_keys=True, ensure_ascii=False)
            if source_key in seen_sources:
                continue
            seen_sources.add(source_key)
            try:
                db_graphs.append(build_database_schema_graph(source_db_type, source_config))
            except Exception as graph_error:
                db_graphs.append(
                    {
                        "source_label": source.get("label") or str(source_db_type or "database"),
                        "db_type": source_db_type,
                        "tables": [],
                        "relationships": [],
                        "error": format_db_error(graph_error, source_db_type),
                    }
                )

        file_profiles = _collect_workspace_file_profiles(session_id, username)
        markdown = _build_data_profile_skill_markdown(username, session_id, db_graphs, file_profiles)

        workspace_dir = get_session_workspace(session_id, username)
        generated_dir = Path(workspace_dir) / "generated"
        generated_dir.mkdir(parents=True, exist_ok=True)
        base_name = f"Data_Exploration_SKILL_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        md_path = _save_md(markdown, base_name, str(generated_dir))
        file_url = build_download_url(f"{username}/{session_id}/generated/{Path(md_path).name}")

        return {
            "success": True,
            "message": "数据探查 SKILL 文档已生成",
            "filename": Path(md_path).name,
            "file_url": file_url,
            "database_source_count": len(db_graphs),
            "file_count": len(file_profiles),
            "summary": f"已整理 {len(db_graphs)} 个数据库来源、{len(file_profiles)} 个工作区文件",
        }
    except Exception as e:
        print(f"Data Profile Report Error: {e}")
        return {"success": False, "message": str(e)}


@app.post("/api/db/generate-sql")
async def generate_sql(body: dict = Body(...)):
    """自然语言生成 SQL（自动获取真实数据库 Schema）"""
    try:
        db_type = body.get("db_type", "mysql")
        config = body.get("config", {})
        prompt = body.get("prompt", "")
        model_provider = body.get("model_provider")

        if not prompt:
            return {"success": False, "message": "Prompt is required"}

        # 尝试连接数据库，获取真实 schema
        schema_info = ""
        try:
            from sqlalchemy import inspect
            normalized_type = normalize_db_type(db_type)
            engine = build_db_engine(normalized_type, config)
            inspector = inspect(engine)
            tables = inspector.get_table_names()
            schema_lines = []
            for table_name in tables:
                columns = inspector.get_columns(table_name)
                col_defs = []
                for col in columns:
                    col_type = str(col.get("type", "unknown"))
                    nullable = "NULL" if col.get("nullable", True) else "NOT NULL"
                    default = f" DEFAULT {col['default']}" if col.get("default") else ""
                    col_defs.append(f"  {col['name']} {col_type} {nullable}{default}")
                pk = inspector.get_pk_constraint(table_name)
                if pk and pk.get("constrained_columns"):
                    col_defs.append(f"  PRIMARY KEY ({', '.join(pk['constrained_columns'])})")
                schema_lines.append(f"CREATE TABLE {table_name} (\n" + ",\n".join(col_defs) + "\n);")
            schema_info = "\n\n".join(schema_lines)
            engine.dispose()
        except Exception as e:
            print(f"Could not fetch schema from DB: {e}")
            schema_info = "(无法自动获取表结构，请确保连接信息正确)"

        system_msg = f"""你是一个精通 SQL 的专家。请根据以下数据库 Schema 和用户需求，生成适用于 {db_type} 数据库的 SQL 查询语句。
严格基于所提供的 Schema 生成查询，不要虚构不存在的表或列。
只返回 SQL 代码块，不要有任何其他文字解释。
代码块格式如下：
```sql
SELECT ...
```"""
        user_msg = f"数据库 Schema：\n{schema_info}\n\n用户需求：{prompt}"

        llm_client, llm_model, _ = get_runtime_llm(model_provider)
        response = llm_client.chat.completions.create(
            model=llm_model,
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg}
            ],
            temperature=0.1
        )

        sql_content = response.choices[0].message.content
        # 提取 SQL 代码块
        import re
        sql_match = re.search(r"```sql\n(.*?)\n```", sql_content, re.DOTALL)
        if sql_match:
            sql = sql_match.group(1).strip()
        else:
            sql = sql_content.strip().replace("```sql", "").replace("```", "")

        return {"success": True, "sql": sql}
    except Exception as e:
        print(f"Generate SQL Error: {e}")
        return {"success": False, "message": str(e)}


@app.post("/api/db/execute")
async def execute_db_sql(body: dict = Body(...)):
    """执行 SQL 并保存结果到工作区"""
    try:
        db_type = body.get("db_type")
        config = body.get("config", {})
        sql = body.get("sql")
        dataset_name = body.get("dataset_name", "query_result")
        mode = body.get("mode", "overwrite") # overwrite 或 append
        format = body.get("format", "csv") # csv 或 json
        session_id = body.get("session_id", "default")
        username = body.get("username", "default")

        if not sql:
            return {"success": False, "message": "SQL is required"}

        engine = build_db_engine(db_type, config)
        try:
            # 使用 pandas 执行查询
            df = pd.read_sql(sql, engine)
        finally:
            engine.dispose()

        workspace_dir = get_session_workspace(session_id, username)
        file_ext = ".csv" if format == "csv" else ".json"
        filename = f"{dataset_name}{file_ext}"
        file_path = Path(workspace_dir) / filename

        if mode == "append" and file_path.exists():
            if format == "csv":
                old_df = pd.read_csv(file_path)
                new_df = pd.concat([old_df, df], ignore_index=True)
                new_df.to_csv(file_path, index=False, encoding="utf-8-sig")
            else:
                old_df = pd.read_json(file_path)
                new_df = pd.concat([old_df, df], ignore_index=True)
                new_df.to_json(file_path, orient="records", force_ascii=False, indent=2)
        else:
            if format == "csv":
                df.to_csv(file_path, index=False, encoding="utf-8-sig")
            else:
                df.to_json(file_path, orient="records", force_ascii=False, indent=2)

        return {
            "success": True,
            "message": f"Successfully executed and saved to {filename}",
            "filename": filename,
            "row_count": len(df)
        }
    except Exception as e:
        print(f"Execute DB SQL Error: {e}")
        return {"success": False, "message": format_db_error(e, body.get("db_type"))}


@app.get("/api/kb/settings")
async def get_kb_settings():
    try:
        return {"success": True, "settings": load_kb_settings(raw=False)}
    except Exception as e:
        print(f"Get KB settings error: {e}")
        return {"success": False, "message": str(e)}


@app.post("/api/kb/settings")
async def update_kb_settings(body: dict = Body(...)):
    try:
        current = load_kb_settings(raw=True)
        merged = normalize_kb_settings_input(body, current)
        save_kb_settings(merged)
        return {"success": True, "settings": load_kb_settings(raw=False), "message": "知识库配置已保存"}
    except Exception as e:
        print(f"Save KB settings error: {e}")
        return {"success": False, "message": str(e)}


@app.post("/api/kb/test")
async def test_kb_settings(body: dict = Body(...)):
    try:
        provider = body.get("provider", "all")
        payload_settings = body.get("settings")
        current = load_kb_settings(raw=True)
        merged = normalize_kb_settings_input(payload_settings, current) if isinstance(payload_settings, dict) else current

        providers = ["onyx", "dify"] if provider == "all" else [provider]
        results = []
        for item in providers:
            if item not in {"onyx", "dify"}:
                return {"success": False, "message": f"Unsupported provider: {item}"}
            results.append(run_kb_provider_test(item, merged))

        for result in results:
            merged.setdefault("test_status", {})[result["provider"]] = {
                "status": result["status"],
                "message": result["message"],
                "tested_at": result["tested_at"],
            }
        save_kb_settings(merged)

        success = all(result["success"] for result in results)
        response: Dict[str, Any] = {
            "success": success,
            "provider": provider,
            "results": results,
            "settings": load_kb_settings(raw=False),
        }
        if len(results) == 1:
            response.update(results[0])
        return response
    except Exception as e:
        print(f"Test KB settings error: {e}")
        return {"success": False, "message": str(e)}


# ========== 用户本地配置持久化 API ==========

@app.get("/api/config/models")
async def get_user_model_configs(username: str = Query("default")):
    """获取用户保存的模型提供商配置列表"""
    data = _load_user_config(username, "model_providers.json", {"providers": []})
    providers = data.get("providers", [])
    selected_id = str(data.get("selected_id", "") or "").strip()
    if not selected_id and isinstance(providers, list) and providers:
        first = providers[0] if isinstance(providers[0], dict) else {}
        selected_id = str(first.get("id", "") or "").strip()
    return {"success": True, "providers": providers, "selected_id": selected_id}


@app.post("/api/config/models")
async def save_user_model_config(body: dict = Body(...)):
    """保存用户的模型提供商配置"""
    username = body.get("username", "default")
    providers = body.get("providers", [])
    selected_id = str(body.get("selected_id", "") or "").strip()
    existing = _load_user_config(username, "model_providers.json", {"providers": []})
    existing["providers"] = providers
    if not selected_id and isinstance(providers, list) and providers:
        first = providers[0] if isinstance(providers[0], dict) else {}
        selected_id = str(first.get("id", "") or "").strip()
    if selected_id:
        existing["selected_id"] = selected_id
    elif not providers:
        existing["selected_id"] = ""
    path = _save_user_config(username, "model_providers.json", existing)
    return {
        "success": True,
        "path": path,
        "selected_id": str(existing.get("selected_id", "") or ""),
        "message": "模型配置已保存到本地",
    }


@app.delete("/api/config/models")
async def delete_user_model_config(body: dict = Body(...)):
    """删除用户保存的某个模型配置"""
    username = body.get("username", "default")
    provider_id = body.get("id", "")
    data = _load_user_config(username, "model_providers.json", {"providers": []})
    data["providers"] = [p for p in data["providers"] if p.get("id") != provider_id]
    _save_user_config(username, "model_providers.json", data)
    return {"success": True, "message": "已删除"}


@app.get("/api/config/databases")
async def get_user_db_configs(username: str = Query("default")):
    """获取用户保存的数据库连接配置列表"""
    data = _load_user_config(username, "database_connections.json", {"connections": []})
    return {"success": True, "connections": data.get("connections", [])}


@app.post("/api/config/databases")
async def save_user_db_config(body: dict = Body(...)):
    """保存用户的数据库连接配置（添加/更新）"""
    username = body.get("username", "default")
    connection = body.get("connection", {})
    conn_id = connection.get("id") or f"db_{int(time_module.time())}"
    connection["id"] = conn_id

    data = _load_user_config(username, "database_connections.json", {"connections": []})
    idx = next((i for i, c in enumerate(data["connections"]) if c.get("id") == conn_id), None)
    if idx is not None:
        data["connections"][idx] = connection
    else:
        data["connections"].append(connection)

    path = _save_user_config(username, "database_connections.json", data)
    return {"success": True, "path": path, "connection": connection, "message": "数据库配置已保存到本地"}


@app.delete("/api/config/databases")
async def delete_user_db_config(body: dict = Body(...)):
    """删除用户保存的某个数据库连接配置"""
    username = body.get("username", "default")
    conn_id = body.get("id", "")
    data = _load_user_config(username, "database_connections.json", {"connections": []})
    data["connections"] = [c for c in data["connections"] if c.get("id") != conn_id]
    _save_user_config(username, "database_connections.json", data)
    return {"success": True, "message": "已删除"}


@app.get("/api/config/knowledge")
async def get_user_knowledge_settings(username: str = Query("default")):
    """获取用户的知识库设置（与 DB settings 合并）"""
    global_settings = load_kb_settings(raw=False)
    user_settings = _load_user_config(username, "knowledge_settings.json", {})
    merged = {**global_settings, **user_settings}
    return {"success": True, "settings": merged}


@app.post("/api/config/knowledge")
async def save_user_knowledge_settings(body: dict = Body(...)):
    """保存用户的知识库设置"""
    username = body.get("username", "default")
    settings = body.get("settings", {})
    _save_user_config(username, "knowledge_settings.json", settings)
    return {"success": True, "message": "知识库配置已保存到本地"}


@app.get("/api/config/analysis-history")
async def get_analysis_history_settings(username: str = Query("default")):
    settings = _load_analysis_history_settings(username)
    return {"success": True, "settings": settings}


@app.post("/api/config/analysis-history")
async def save_analysis_history_settings(body: dict = Body(...)):
    username = body.get("username", "default")
    settings = _sanitize_analysis_history_settings(body.get("settings", {}))
    _save_user_config(username, "analysis_history_settings.json", settings)
    return {"success": True, "settings": settings, "message": "分析历史配置已保存到本地"}


@app.get("/api/analysis/history")
async def list_analysis_history_runs(username: str = Query("default"), limit: int = Query(30, ge=1, le=200)):
    settings = _load_analysis_history_settings(username)
    data = _load_analysis_history_index(username)
    runs = [item for item in data.get("runs", []) if isinstance(item, dict)][:limit]
    stats = {
        "total": len(data.get("runs", [])),
        "completed": sum(1 for item in data.get("runs", []) if isinstance(item, dict) and item.get("status") == "completed"),
        "failed": sum(1 for item in data.get("runs", []) if isinstance(item, dict) and item.get("status") == "failed"),
        "warning": sum(1 for item in data.get("runs", []) if isinstance(item, dict) and item.get("status") == "warning"),
    }
    return {"success": True, "settings": settings, "runs": runs, "stats": stats}


@app.get("/api/analysis/history/{run_id}")
async def get_analysis_history_run(run_id: str, username: str = Query("default")):
    data = _load_analysis_history_index(username)
    run_summary = next(
        (item for item in data.get("runs", []) if isinstance(item, dict) and item.get("run_id") == run_id),
        None,
    )
    if not run_summary:
        raise HTTPException(status_code=404, detail="分析历史不存在")

    run_path = _get_analysis_history_run_path(username, run_id)
    events: List[Dict[str, Any]] = []
    if os.path.exists(run_path):
        with open(run_path, "r", encoding="utf-8") as f:
            for raw_line in f:
                line = raw_line.strip()
                if not line:
                    continue
                try:
                    payload = json.loads(line)
                except Exception:
                    continue
                if isinstance(payload, dict):
                    events.append(payload)

    return {
        "success": True,
        "run": run_summary,
        "events": events,
        "settings": _load_analysis_history_settings(username),
    }


@app.get("/api/config/export")
async def export_user_config(username: str = Query("default")):
    """导出用户所有配置为一个 JSON 文件（用于备份或迁移）"""
    config_dir = _get_user_config_dir(username)
    result = {"username": username, "configs": {}}
    for fname in os.listdir(config_dir):
        if fname.endswith(".json"):
            try:
                with open(os.path.join(config_dir, fname), "r", encoding="utf-8") as f:
                    result["configs"][fname] = json.load(f)
            except Exception:
                result["configs"][fname] = {}
    return result


if __name__ == "__main__":
    print("🚀 启动后端服务...")
    print(f"   - API服务: http://localhost:8200")
    print(f"   - 文件服务: http://localhost:8100")
    uvicorn.run(app, host="0.0.0.0", port=8200)
