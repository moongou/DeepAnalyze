import os
import pandas as pd
from datetime import timedelta
from pathlib import Path

# Try to import established utils
try:
    from .pdf_utils import register_chinese_fonts, get_chinese_font_name, generate_pdf as gen_pdf_rl
    from .docx_utils import generate_docx as gen_docx_impl
except (ImportError, ValueError):
    # Fallback for direct execution or different import contexts
    try:
        from pdf_utils import register_chinese_fonts, get_chinese_font_name, generate_pdf as gen_pdf_rl
        from docx_utils import generate_docx as gen_docx_impl
    except ImportError:
        register_chinese_fonts = None
        gen_pdf_rl = None
        gen_docx_impl = None

def get_fonts_dir():
    """Find fonts directory dynamically."""
    # Try using config if available
    try:
        from config import get_fonts_dir as gfd
        return gfd()
    except ImportError:
        pass

    current_dir = Path(__file__).parent.resolve()
    project_root = current_dir
    # If we are in demo/chat/
    if (project_root / "assets").exists():
        return project_root / "assets" / "fonts"

    # Try parent
    project_root = current_dir.parent
    if (project_root / "assets").exists():
        return project_root / "assets" / "fonts"

    # Fallback to absolute known path for this specific environment
    fallback = Path("/Users/m3max/IdeaProjects/DeepAnalyze/assets/fonts")
    if fallback.exists():
        return fallback

    return Path("assets/fonts")

def init_fpdf_chinese():
    """
    Returns an FPDF instance with pre-registered Chinese fonts.
    Solves the 'Undefined font: simhei' issue by ensuring all instances
    have the same registration.
    """
    try:
        from fpdf import FPDF
    except ImportError:
        print("Error: fpdf2 library not installed. Please use 'pip install fpdf2'")
        return None

    pdf = FPDF()
    fonts_dir = get_fonts_dir()

    # Fonts to register
    # Format: (name, filename)
    fonts = [
        ('SimHei', 'simhei.ttf'),
        ('SimKai', 'simkai.ttf'),
        ('STFangSong', 'STFangSong.ttf'),
        ('STHeiti', 'STHeiti.ttf')
    ]

    count = 0
    for font_name, filename in fonts:
        font_path = fonts_dir / filename
        if font_path.exists():
            # Register regular
            pdf.add_font(font_name, '', str(font_path))
            # Register bold (using same file as fallback if dedicated bold doesn't exist)
            pdf.add_font(font_name, 'B', str(font_path))
            count += 1

    if count == 0:
        print(f"Warning: No Chinese fonts found in {fonts_dir}")

    return pdf

def generate_report_pdf(md_text, output_path, title="Analysis Report"):
    """
    Generate a PDF report from Markdown text using the best available method.
    Priority: 1. reportlab (via pdf_utils), 2. fpdf2
    """
    if gen_pdf_rl:
        return gen_pdf_rl(md_text, output_path, title=title)

    # Fallback to FPDF2
    pdf = init_fpdf_chinese()
    if not pdf:
        return False

    pdf.add_page()
    pdf.set_font('SimHei', 'B', 16)
    pdf.cell(0, 10, title, ln=True, align='C')
    pdf.ln(5)

    pdf.set_font('STFangSong', '', 12)
    # Simple markdown-ish to PDF converter for fpdf2
    lines = md_text.split('\n')
    for line in lines:
        if line.startswith('# '):
            pdf.set_font('SimHei', 'B', 16)
            pdf.multi_cell(0, 10, line[2:].strip())
        elif line.startswith('## '):
            pdf.set_font('SimHei', 'B', 14)
            pdf.multi_cell(0, 10, line[3:].strip())
        elif line.startswith('### '):
            pdf.set_font('SimHei', 'B', 12)
            pdf.multi_cell(0, 10, line[4:].strip())
        else:
            pdf.set_font('STFangSong', '', 11)
            pdf.multi_cell(0, 8, line)

    pdf.output(output_path)
    return True

def generate_report_docx(md_text, output_path, title="Analysis Report"):
    """Generate a DOCX report from Markdown text."""
    if gen_docx_impl:
        return gen_docx_impl(md_text, output_path, title=title)
    return False

def generate_report_pptx(md_text, output_path, title="Analysis Report"):
    """
    Generate a PPTX report from Markdown text.
    Uses python-pptx with Chinese font support, multi-paragraph formatting, and image embedding.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import re as _re
        from datetime import datetime as _dt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Chinese font configuration
        _font_title = "SimHei"
        _font_body = "STFangSong"

        def _set_run_font(run, font_name, size_pt, bold=False, italic=False, color_rgb=None):
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = font_name
            if color_rgb:
                run.font.color.rgb = color_rgb
            try:
                from lxml import etree
                rPr = run._r.get_or_add_rPr()
                ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
                ea.set('typeface', font_name)
            except Exception:
                pass

        def _add_formatted_text(text_frame, text, font_name, size_pt, color_rgb=None):
            """Add multi-paragraph formatted text to a text frame"""
            lines = text.split('\n')
            first = True
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                if first:
                    p = text_frame.paragraphs[0]
                    first = False
                else:
                    p = text_frame.add_paragraph()
                p.space_after = Pt(4)

                is_bullet = line.startswith('- ') or line.startswith('* ') or line.startswith('• ')
                if is_bullet:
                    line = line.lstrip('-*• ').strip()
                    r = p.add_run()
                    r.text = "• "
                    _set_run_font(r, font_name, size_pt, bold=True, color_rgb=RGBColor(0, 102, 204))
                    r2 = p.add_run()
                    r2.text = line
                    _set_run_font(r2, font_name, size_pt, color_rgb=color_rgb or RGBColor(51, 51, 51))
                elif line.startswith('**') and line.endswith('**'):
                    r = p.add_run()
                    r.text = line.strip('*')
                    _set_run_font(r, font_name, size_pt + 1, bold=True, color_rgb=color_rgb or RGBColor(0, 51, 102))
                else:
                    clean = _re.sub(r'\*\*(.*?)\*\*', r'\1', line)
                    clean = _re.sub(r'\*(.*?)\*', r'\1', clean)
                    clean = _re.sub(r'`(.*?)`', r'\1', clean)
                    clean = _re.sub(r'!\[.*?\]\(.*?\)', '', clean)
                    clean = _re.sub(r'\[([^\]]*)\]\(.*?\)', r'\1', clean)
                    if clean.strip():
                        r = p.add_run()
                        r.text = clean
                        _set_run_font(r, font_name, size_pt, color_rgb=color_rgb or RGBColor(51, 51, 51))

        # === Title slide with styled background ===
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(3.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0, 51, 102)
        bg.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11.333), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        _set_run_font(run, _font_title, 36, bold=True, color_rgb=RGBColor(255, 255, 255))

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = f"DeepAnalyze 智能分析报告 | {_dt.now().strftime('%Y年%m月%d日')}"
        _set_run_font(run2, _font_body, 16, color_rgb=RGBColor(200, 220, 255))

        # Content slides
        raw_sections = _re.split(r'\n(?=#{1,3}\s+)', md_text)
        output_dir = Path(output_path).parent if output_path else None

        for section in raw_sections:
            if not section.strip():
                continue
            lines = section.strip().split('\n')
            sec_title = lines[0].strip().lstrip('#').strip() if lines else "分析内容"
            body = '\n'.join(lines[1:]).strip() if len(lines) > 1 else ""
            if not sec_title and not body:
                continue

            # Extract and find image references
            img_refs = _re.findall(r'!\[.*?\]\((.*?)\)', body)
            embedded_imgs = []
            for ref in img_refs:
                img_name = os.path.basename(ref)
                candidates = [Path(ref)]
                if output_dir:
                    candidates = [
                        output_dir / img_name,
                        output_dir / "generated" / img_name,
                        Path(ref),
                    ]
                for candidate in candidates:
                    if candidate.exists() and candidate.stat().st_size > 1024:
                        embedded_imgs.append(str(candidate))
                        break

            # Clean body text
            clean_body = _re.sub(r'!\[.*?\]\(.*?\)', '', body)
            clean_body = _re.sub(r'\[([^\]]*)\]\(.*?\)', r'\1', clean_body)
            clean_body = clean_body.strip()

            if not clean_body and not embedded_imgs:
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Title bar
            title_bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(0, 51, 102)
            title_bg.line.fill.background()

            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = sec_title[:80]
            _set_run_font(run, _font_title, 24, bold=True, color_rgb=RGBColor(255, 255, 255))

            # Body with optional image
            if embedded_imgs:
                if clean_body:
                    if len(clean_body) > 2000:
                        clean_body = clean_body[:2000] + "\n...(内容已截断)"
                    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.8))
                    tf2 = txBox2.text_frame
                    tf2.word_wrap = True
                    _add_formatted_text(tf2, clean_body, _font_body, 12)

                for idx, img in enumerate(embedded_imgs[:2]):
                    try:
                        slide.shapes.add_picture(img, Inches(6.8), Inches(1.3 + idx * 3.0), Inches(5.8), Inches(2.8))
                    except Exception:
                        pass
            elif clean_body:
                if len(clean_body) > 2500:
                    clean_body = clean_body[:2500] + "\n...(内容已截断)"
                txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.8))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                _add_formatted_text(tf2, clean_body, _font_body, 13)

        prs.save(str(output_path))
        return True
    except ImportError:
        print("python-pptx 未安装。请运行: pip install python-pptx")
        return False
    except Exception as e:
        print(f"PPTX 生成失败: {e}")
        return False

def pandas_date_sub(series_or_timestamp, days):
    """
    Safely subtract days from a pandas series or Timestamp.
    Fixes: TypeError: Addition/subtraction of integers and integer-arrays with Timestamp is no longer supported
    """
    try:
        if isinstance(series_or_timestamp, (pd.Series, pd.Index)):
            return series_or_timestamp - pd.Timedelta(days=days)
        elif isinstance(series_or_timestamp, (pd.Timestamp, pd.DatetimeIndex)):
            return series_or_timestamp - pd.Timedelta(days=days)
        else:
            # Fallback for standard datetime
            return series_or_timestamp - timedelta(days=days)
    except Exception as e:
        print(f"Date math error: {e}")
        return series_or_timestamp

def pandas_date_add(series_or_timestamp, days):
    """Safely add days to a pandas series or Timestamp."""
    try:
        if isinstance(series_or_timestamp, (pd.Series, pd.Index)):
            return series_or_timestamp + pd.Timedelta(days=days)
        elif isinstance(series_or_timestamp, (pd.Timestamp, pd.DatetimeIndex)):
            return series_or_timestamp + pd.Timedelta(days=days)
        else:
            return series_or_timestamp + timedelta(days=days)
    except Exception as e:
        print(f"Date math error: {e}")
        return series_or_timestamp
